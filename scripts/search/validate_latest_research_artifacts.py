#!/usr/bin/env python3
"""Validate latest derived research artifacts and deck/report references."""

from __future__ import annotations

# Keep categorized scripts import-compatible when executed by file path.
import sys as _sys
from pathlib import Path as _Path
_SCRIPT_ROOT = _Path(__file__).resolve().parents[1]
_REPO_ROOT = _SCRIPT_ROOT.parent
for _candidate in [_REPO_ROOT, _SCRIPT_ROOT, *_SCRIPT_ROOT.iterdir()]:
    if _candidate.is_dir():
        _value = str(_candidate)
        if _value not in _sys.path:
            _sys.path.insert(0, _value)

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
DECK = ROOT.parent / "研究进展汇报.pptx"
REPORT = ROOT / "docs/reports/2026-06-03-realtime-dual-agent-roadmap.md"


def read_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def deck_text(path: Path) -> str:
    if not path.exists():
        return ""
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def add_check(checks: list[dict[str, Any]], name: str, passed: bool, evidence: str) -> None:
    checks.append(
        {
            "name": name,
            "status": "pass" if passed else "fail",
            "evidence": evidence,
        }
    )


def missing_fragments(text: str, fragments: list[str]) -> list[str]:
    return [fragment for fragment in fragments if fragment not in text]


def first_run_wall_seconds(data: dict[str, Any]) -> object:
    runs = data.get("runs") or []
    if runs and isinstance(runs[0], dict):
        return runs[0].get("wall_seconds")
    return data.get("wall_seconds")


def find_summary_variant(data: dict[str, Any], variant: str) -> dict[str, Any]:
    for row in data.get("summary", data.get("summaries", [])):
        if row.get("variant") == variant:
            return row
    return {}


def validate() -> dict[str, Any]:
    checks: list[dict[str, Any]] = []
    snapshot = read_json(ROOT / "outputs/research_progress_snapshot/snapshot.json")
    runtime_audit = read_json(ROOT / "outputs/runtime_evidence_audit/runtime_evidence_audit.json").get("summary", {})
    timeline = read_json(ROOT / "outputs/system_timeline/summary.json")
    review = read_json(ROOT / "outputs/timeline_review_audit/llm_review_signal_timeline_audit_summary.json")
    clean_llm = read_json(ROOT / "outputs/llm_window_batch/qwen36_flash_clean_high_rule_auto_audit_full_agreement.json")
    selector_holdout = read_json(ROOT / "outputs/recover_selector_split_120/recording_holdout_summary.json")
    selector_bootstrap = read_json(ROOT / "outputs/realtime_contract_bootstrap_120/realtime_contract_bootstrap.json")
    selector_bootstrap_best = find_summary_variant(selector_bootstrap, "rule_recover_policy_sweep_best")
    omni_fusion = read_json(ROOT / "outputs/omni_guard/omni_acoustic_fusion_summary.json")
    omni_expansion = read_json(ROOT / "outputs/research_progress_snapshot/omni_expansion_manifest.json")
    omni_call_manifest = read_json(ROOT / "outputs/research_progress_snapshot/omni48_live_call_manifest.json")
    split_top3 = read_json(ROOT / "outputs/runtime_safe_llm_window_batch/deepseek_proxy_high_risk_split20_top3_parallel_comparison.json")
    split_attempt = read_json(ROOT / "outputs/runtime_safe_llm_window_batch/deepseek_proxy_high_risk_split20_top4_5_parallel_attempt_summary.json")
    qwen_split = read_json(ROOT / "outputs/runtime_safe_llm_window_batch/qwen36_flash_split20_top4_5_parallel_comparison.json")
    guard_tuning = read_json(ROOT / "outputs/runtime_safe_llm_window_batch/llm_guard_tuning_104w_summary.json")
    tuned_guard = read_json(ROOT / "outputs/runtime_safe_llm_window_batch/deepseek_proxy_high_risk_104w_tuned_passthrough_safety_summary.json")
    tuned_timeline = read_json(ROOT / "outputs/runtime_safe_llm_window_batch/tuned_v2_writeback_timeline/rule_writeback_timeline_summary.json")
    latency_ledger = read_json(ROOT / "outputs/research_progress_snapshot/runtime_latency_budget_ledger.json")
    latency_slo_audit = read_json(ROOT / "outputs/research_progress_snapshot/stage_latency_slo_audit.json")
    latency_risk_margin = read_json(ROOT / "outputs/research_progress_snapshot/latency_risk_margin_audit.json")
    latency_risk_mitigation = read_json(ROOT / "outputs/research_progress_snapshot/latency_risk_mitigation_plan.json")
    claims_manifest = read_json(ROOT / "outputs/research_progress_snapshot/claims_manifest.json")
    memory_replay = read_json(ROOT / "outputs/research_progress_snapshot/memory_update_replay.json")
    replay_manifest = read_json(ROOT / "outputs/research_progress_snapshot/runtime_replay_manifest.json")
    live_readiness = read_json(ROOT / "outputs/research_progress_snapshot/live_run_readiness.json")
    live_agent_plan = read_json(ROOT / "outputs/research_progress_snapshot/live_agent_execution_plan.json")
    live_postrun_closure = read_json(ROOT / "outputs/research_progress_snapshot/live_postrun_metrics_closure.json")
    live_output_audit = read_json(ROOT / "outputs/research_progress_snapshot/live_output_audit.json")
    live_scoring_readiness = read_json(ROOT / "outputs/research_progress_snapshot/live_scoring_readiness.json")
    post_live_scoring_execution_plan = read_json(ROOT / "outputs/research_progress_snapshot/post_live_scoring_execution_plan.json")
    post_live_scoring_launcher = read_json(ROOT / "outputs/research_progress_snapshot/post_live_scoring_launcher.json")
    post_live_scoring_receipt = read_json(ROOT / "outputs/research_progress_snapshot/post_live_scoring_receipt_audit.json")
    post_live_scoring_output_audit = read_json(ROOT / "outputs/research_progress_snapshot/post_live_scoring_output_audit.json")
    post_live_evidence_dependency_dag = read_json(ROOT / "outputs/research_progress_snapshot/post_live_evidence_dependency_dag.json")
    live_input_integrity = read_json(ROOT / "outputs/research_progress_snapshot/live_input_integrity_audit.json")
    live_execution_runbook = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_runbook.json")
    live_execution_bundle = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_bundle.json")
    live_execution_handoff = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_handoff_packet.json")
    live_command_surface_audit = read_json(ROOT / "outputs/research_progress_snapshot/live_command_surface_audit.json")
    live_execution_eligibility = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_eligibility_gate.json")
    live_execution_receipt = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_receipt_audit.json")
    live_execution_launcher = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_launcher.json")
    live_output_repair_plan = read_json(ROOT / "outputs/research_progress_snapshot/live_output_repair_plan.json")
    live_resume_state_audit = read_json(ROOT / "outputs/research_progress_snapshot/live_resume_state_audit.json")
    live_runtime_environment_audit = read_json(ROOT / "outputs/research_progress_snapshot/live_runtime_environment_audit.json")
    live_execution_timing = read_json(ROOT / "outputs/research_progress_snapshot/live_execution_timing_plan.json")
    live_retry_budget = read_json(ROOT / "outputs/research_progress_snapshot/live_retry_budget_audit.json")
    live_token_quota_budget = read_json(ROOT / "outputs/research_progress_snapshot/live_token_quota_budget_audit.json")
    live_failure_recovery = read_json(ROOT / "outputs/research_progress_snapshot/live_failure_recovery_playbook.json")
    live_metric_extraction = read_json(ROOT / "outputs/research_progress_snapshot/live_metric_extraction_contract.json")
    live_output_schema = read_json(ROOT / "outputs/research_progress_snapshot/live_output_schema_contract.json")
    post_live_acceptance_scorecard = read_json(ROOT / "outputs/research_progress_snapshot/post_live_acceptance_scorecard.json")
    post_live_latency_claim_matrix = read_json(ROOT / "outputs/research_progress_snapshot/post_live_latency_claim_matrix.json")
    post_live_time_metric_statistics = read_json(ROOT / "outputs/research_progress_snapshot/post_live_time_metric_statistics_plan.json")
    post_live_time_metric_extractor = read_json(ROOT / "outputs/research_progress_snapshot/post_live_time_metric_extractor.json")
    post_live_time_metric_receipt = read_json(ROOT / "outputs/research_progress_snapshot/post_live_time_metric_receipt_audit.json")
    post_live_promotion_preflight = read_json(ROOT / "outputs/research_progress_snapshot/post_live_promotion_preflight_audit.json")
    live_parallelism = read_json(ROOT / "outputs/research_progress_snapshot/live_parallelism_sensitivity.json")
    post_live_claim_promotion = read_json(ROOT / "outputs/research_progress_snapshot/post_live_claim_promotion_gate.json")
    post_live_claim_promotion_receipt = read_json(ROOT / "outputs/research_progress_snapshot/post_live_claim_promotion_receipt_audit.json")
    phase_result_scorecard = read_json(ROOT / "outputs/research_progress_snapshot/phase_result_scorecard.json")
    live_provider_routing = read_json(ROOT / "outputs/research_progress_snapshot/live_provider_routing_decision.json")
    report_ppt_traceability = read_json(ROOT / "outputs/research_progress_snapshot/report_ppt_traceability.json")
    selector_candidate_scan = read_json(ROOT / "outputs/research_progress_snapshot/selector_true_heldout_candidate_scan.json")
    selector_split_validation = read_json(ROOT / "outputs/research_progress_snapshot/selector_true_heldout_split_validation.json")
    selector_protocol = read_json(ROOT / "outputs/research_progress_snapshot/selector_true_heldout_protocol.json")
    split20_manifest = read_json(ROOT / "outputs/research_progress_snapshot/split20_full_live_manifest.json")
    split20_export_audit = read_json(ROOT / "outputs/research_progress_snapshot/split20_resume_export_audit.json")
    split15_export_audit = read_json(ROOT / "outputs/research_progress_snapshot/split15_stretch_reexport_audit.json")
    split_policy_optimization = read_json(ROOT / "outputs/research_progress_snapshot/split_policy_optimization.json")
    next_queue = read_json(ROOT / "outputs/research_progress_snapshot/next_experiment_queue.json")
    tuned_boundary = find_summary_variant(tuned_timeline, "rule_boundary_recover")
    tuned_recover = find_summary_variant(tuned_timeline, "rule_recover_policy_sweep_best")
    pareto_md = (ROOT / "outputs/latency_tradeoff/der_latency_pareto.md").read_text(encoding="utf-8") if (ROOT / "outputs/latency_tradeoff/der_latency_pareto.md").exists() else ""
    report_text = REPORT.read_text(encoding="utf-8") if REPORT.exists() else ""
    ppt_text = deck_text(DECK)
    runtime_audit_phrase = (
        f"pass; {int(runtime_audit.get('artifact_count', 0))} artifacts; "
        f"blocking {int(runtime_audit.get('runtime_blocking_count', -1))}"
    )
    report_required_fragments = [
        "当前进展快照",
        "scripts/misc/refresh_latest_research_artifacts.py",
        runtime_audit_phrase,
        "Selector generalization",
        "Omni fusion",
        "Guard tuning contract",
        "Split20 LLM guard optimization",
        "combined_review0.5_keepfast0.5",
        "Dual Agent LLM review signal",
        "Omni fusion no-writeback contract",
        "outputs/research_progress_snapshot/omni_expansion_manifest.md",
        "outputs/research_progress_snapshot/omni48_live_call_manifest.md",
        "selector recording-holdout validation",
        "outputs/research_progress_snapshot/claims_manifest.md",
        "outputs/research_progress_snapshot/runtime_latency_budget_ledger.md",
        "outputs/research_progress_snapshot/stage_latency_slo_audit.md",
        "outputs/research_progress_snapshot/latency_risk_margin_audit.md",
        "outputs/research_progress_snapshot/latency_risk_mitigation_plan.md",
        "outputs/research_progress_snapshot/memory_update_replay.md",
        "outputs/research_progress_snapshot/runtime_replay_manifest.md",
        "outputs/research_progress_snapshot/live_run_readiness.md",
        "outputs/research_progress_snapshot/live_agent_execution_plan.md",
        "outputs/research_progress_snapshot/live_postrun_metrics_closure.md",
        "outputs/research_progress_snapshot/live_output_audit.md",
        "outputs/research_progress_snapshot/live_scoring_readiness.md",
        "outputs/research_progress_snapshot/post_live_scoring_execution_plan.md",
        "outputs/research_progress_snapshot/post_live_scoring_launcher.md",
        "outputs/research_progress_snapshot/post_live_scoring_receipt_audit.md",
        "outputs/research_progress_snapshot/post_live_scoring_output_audit.md",
        "outputs/research_progress_snapshot/post_live_evidence_dependency_dag.md",
        "outputs/research_progress_snapshot/live_input_integrity_audit.md",
        "outputs/research_progress_snapshot/live_execution_runbook.md",
        "outputs/research_progress_snapshot/live_execution_bundle.md",
        "outputs/research_progress_snapshot/live_execution_handoff_packet.md",
        "outputs/research_progress_snapshot/live_command_surface_audit.md",
        "outputs/research_progress_snapshot/live_execution_eligibility_gate.md",
        "outputs/research_progress_snapshot/live_execution_receipt_audit.md",
        "outputs/research_progress_snapshot/live_execution_launcher.md",
        "outputs/research_progress_snapshot/live_output_repair_plan.md",
        "outputs/research_progress_snapshot/live_resume_state_audit.md",
        "outputs/research_progress_snapshot/live_runtime_environment_audit.md",
        "outputs/research_progress_snapshot/live_execution_timing_plan.md",
        "outputs/research_progress_snapshot/live_retry_budget_audit.md",
        "outputs/research_progress_snapshot/live_token_quota_budget_audit.md",
        "outputs/research_progress_snapshot/live_failure_recovery_playbook.md",
        "outputs/research_progress_snapshot/live_metric_extraction_contract.md",
        "outputs/research_progress_snapshot/live_output_schema_contract.md",
        "outputs/research_progress_snapshot/post_live_acceptance_scorecard.md",
        "outputs/research_progress_snapshot/post_live_latency_claim_matrix.md",
        "outputs/research_progress_snapshot/post_live_time_metric_statistics_plan.md",
        "outputs/research_progress_snapshot/post_live_time_metric_extractor.md",
        "outputs/research_progress_snapshot/post_live_time_metric_receipt_audit.md",
        "outputs/research_progress_snapshot/post_live_promotion_preflight_audit.md",
        "outputs/research_progress_snapshot/live_parallelism_sensitivity.md",
        "outputs/research_progress_snapshot/post_live_claim_promotion_gate.md",
        "outputs/research_progress_snapshot/post_live_claim_promotion_receipt_audit.md",
        "outputs/research_progress_snapshot/phase_result_scorecard.md",
        "outputs/research_progress_snapshot/live_provider_routing_decision.md",
        "outputs/research_progress_snapshot/report_ppt_traceability.md",
        "outputs/research_progress_snapshot/selector_true_heldout_candidate_scan.md",
        "outputs/research_progress_snapshot/selector_true_heldout_split_validation.md",
        "outputs/research_progress_snapshot/selector_true_heldout_protocol.md",
        "outputs/research_progress_snapshot/split20_full_live_manifest.md",
        "outputs/research_progress_snapshot/split20_resume_export_audit.md",
        "outputs/research_progress_snapshot/split15_stretch_reexport_audit.md",
        "outputs/research_progress_snapshot/split_policy_optimization.md",
        "outputs/research_progress_snapshot/next_experiment_queue.md",
        "selector_generalization_positive",
        "omni_fusion_label_only",
        "boundary_auto_writeback_negative_control",
        "do_not_deploy",
        "Phase result scorecard",
        "phase_result_scorecard_from_existing_artifacts_no_live_calls",
        "Result rows: `8`",
        "DeepSeek API no-go",
        "Live provider routing decision",
        "live_provider_routing_decision_no_live_calls_no_secret_values",
        "Route rows: `4`",
        "Default execute scope: `none`",
        "Runtime latency budget ledger",
        "runtime_latency_budget_ledger_from_existing_artifacts",
        "Claim-now rows: `4`",
        "Stage latency SLO audit",
        "stage_latency_slo_audit_from_latency_ledger_no_live_calls",
        "Claim-now SLO pass: `4/4`",
        "Guard P95 margin: `1.151`",
        "Latency risk margin audit",
        "latency_risk_margin_audit_from_slo_no_live_calls",
        "Tight-margin rows: `1`",
        "Guard risk level",
        "0.0177",
        "Latency risk mitigation plan",
        "latency_risk_mitigation_plan_no_live_calls",
        "Primary mitigation: `max20_resume`",
        "max15_reexport",
        "full_split20_live_104w",
        "true_heldout_selector_recordings",
        "omni_fusion_expand_48_or_120",
        "Omni48 manifest",
        "Omni48 call manifest",
        "96 calls",
        "label_only_no_timeline_writeback",
        "end_to_end_runtime_replay_manifest",
        "memory_update_audit_replay",
        "memory replay",
        "runtime replay",
        "live readiness",
        "Live Agent execution plan",
        "live_agent_execution_plan_no_live_calls",
        "Planned live calls: `382`",
        "DeepSeek resume calls: `139`",
        "Live postrun metrics closure",
        "live_postrun_metrics_closure_no_live_calls",
        "DeepSeek success calls: `8`",
        "Omni48 successful calls: `0`",
        "Live output audit",
        "live_output_audit_no_live_calls",
        "Expected live calls: `382`",
        "Missing output surfaces: `3`",
        "Live scoring readiness",
        "live_scoring_readiness_no_live_calls",
        "Scoring steps: `5`",
        "P0 scoring steps: `2`",
        "Ready to score: `0`",
        "Post-live scoring execution plan",
        "post_live_scoring_execution_plan_no_live_calls",
        "Scoring execution steps: `6`",
        "Blocked execution steps: `6`",
        "Post-live scoring launcher",
        "post_live_scoring_launcher_dry_run_no_scoring_calls",
        "Ready scoring rows: `0`",
        "Executed scoring rows: `0`",
        "Scoring execute record exists: `False`",
        "Post-live scoring receipt audit",
        "post_live_scoring_receipt_audit_no_live_or_scoring_calls_no_secret_values",
        "Scoring receipt rows: `6`",
        "Post-live scoring output audit",
        "post_live_scoring_output_audit_no_scoring_calls",
        "Scoring output rows: `6`",
        "Promotion-ready rows: `0`",
        "Post-live evidence dependency DAG",
        "post_live_evidence_dependency_dag_no_live_calls",
        "DAG nodes: `10`",
        "Blocked nodes: `10`",
        "Live input integrity audit",
        "live_input_integrity_audit_no_live_calls",
        "Input-ready surfaces: `3`",
        "DeepSeek resume 139 prompts",
        "Live execution runbook",
        "live_execution_runbook_no_live_calls_no_secret_values",
        "Steps: `7`",
        "P0 planned live calls: `139`",
        "DeepSeek primary policy: `max20`",
        "Live execution bundle",
        "live_execution_bundle_no_live_calls_no_secret_values",
        "Bundle steps: `8`",
        "Blocked/waiting steps: `8`",
        "Live execution handoff packet",
        "live_execution_handoff_packet_no_live_calls",
        "Packet rows: `7`",
        "Handoff blocked/waiting rows: `6`",
        "Credential ready: `False`",
        "Known provider quota blockers: `1`",
        "Live command surface audit",
        "live_command_surface_audit_no_live_calls",
        "Command-ready: `3`",
        "skip-existing commands 3",
        "bounded-retry commands 3",
        "Live execution eligibility gate",
        "live_execution_eligibility_gate_no_live_calls_no_secret_values",
        "Eligibility rows: `7`",
        "Ready to execute live: `False`",
        "Live execution receipt audit",
        "live_execution_receipt_audit_no_live_calls_no_secret_values",
        "Receipt rows: `6`",
        "Execute record exists: `False`",
        "duplicate outputs 0",
        "secret literal commands 0",
        "Live execution launcher",
        "live_execution_launcher_dry_run_no_live_calls",
        "Launcher selected calls: `139`",
        "Started live command calls: `0`",
        "Postrun refresh executed: `False`",
        "Execute record exists: `False`",
        "Live output repair plan",
        "live_output_repair_plan_no_live_calls_no_scoring",
        "Repair rows: `3`",
        "Missing calls: `382`",
        "Live resume state audit",
        "live_resume_state_audit_no_live_calls",
        "Clean-run surfaces: `3`",
        "skip-existing supported 3",
        "bounded retry supported 3",
        "append resume supported 0",
        "quarantine required 0",
        "Live runtime environment audit",
        "live_runtime_environment_audit_no_live_calls",
        "Checks passed: `14`",
        "modules 6/6",
        "credential ready false",
        "Live execution timing plan",
        "live_execution_timing_plan_no_live_calls",
        "DeepSeek estimated wall: `384.444`",
        "139 calls / 8 workers / 18 waves",
        "Live retry budget audit",
        "live_retry_budget_audit_no_live_calls",
        "Max attempted requests: `764`",
        "P0 max attempted requests 278",
        "DeepSeek retry ceiling wall 804.888s",
        "Live token quota budget audit",
        "live_token_quota_budget_audit_no_live_calls",
        "LLM retry token proxy ceiling: `1658856`",
        "P0 retry token proxy ceiling 801724",
        "Omni48 retry clip-model seconds ceiling 1536.0",
        "Live failure recovery playbook",
        "live_failure_recovery_playbook_no_live_calls",
        "Scenarios: `8`",
        "Current blocker scenarios 5",
        "Ready recovery actions: `8`",
        "Live metric extraction contract",
        "live_metric_extraction_contract_no_live_calls",
        "Metric contracts: `8`",
        "Time metric contracts: `3`",
        "Safety metric contracts: `2`",
        "Omni metric contracts: `2`",
        "Expected input calls 为 `676`",
        "Live output schema contract",
        "live_output_schema_contract_no_live_calls",
        "Schema contracts: `8`",
        "Required fields: `62`",
        "Live output schema contracts: `3`",
        "Scoring output schema contracts: `3`",
        "Post-live acceptance scorecard",
        "post_live_acceptance_scorecard_no_live_calls",
        "Scorecard rows: `9`",
        "Blocked rows: `6`",
        "Claim-now SLO pass: `4/4`",
        "Post-live latency claim matrix",
        "post_live_latency_claim_matrix_no_live_calls",
        "Latency claim rows: `8`",
        "Claim-now preserve rows: `4`",
        "Post-live time metric statistics plan",
        "post_live_time_metric_statistics_plan_no_live_calls",
        "Time statistic rows: `9`",
        "Formula count 9",
        "Post-live time metric extractor",
        "post_live_time_metric_extractor_no_live_calls",
        "Extractor rows: `3`",
        "Ready time metric rows 0",
        "Post-live time metric receipt audit",
        "post_live_time_metric_receipt_audit_no_live_or_scoring_calls_no_secret_values",
        "Time metric receipt rows: `6`",
        "Ready for time claim promotion: `False`",
        "Post-live promotion preflight audit",
        "post_live_promotion_preflight_audit_no_live_or_scoring_calls",
        "Preflight rows: `6`",
        "Ready for promotion review: `False`",
        "Live parallelism sensitivity",
        "live_parallelism_sensitivity_no_live_calls",
        "Recommended policy/workers: `max20` / `8`",
        "worker12 wall 256.296s",
        "Post-live claim promotion gate",
        "post_live_claim_promotion_gate_no_live_calls",
        "Gates: `8`",
        "Ready to promote: `0`",
        "Blocked: `5`",
        "promote_only_after_output_audit_scoring_slo_and_traceability_pass",
        "Post-live claim promotion receipt audit",
        "post_live_claim_promotion_receipt_audit_no_live_or_scoring_or_claim_writes_no_secret_values",
        "Claim promotion receipt rows: `6`",
        "Ready for claim write: `False`",
        "Report/PPT traceability",
        "report_ppt_traceability_from_existing_artifacts_no_live_calls",
        "Rows: `54`",
        "Fully covered rows: `54`",
        "Missing report rows: `0`",
        "Missing PPT rows: `0`",
        "selector true-heldout candidate scan",
        "selector_true_heldout_candidate_scan_no_metric_claim",
        "Eligible true-heldout recordings: `0`",
        "Missing new recordings to minimum: `8`",
        "selector true-heldout split validation",
        "selector_true_heldout_split_validation_no_metric_claim",
        "True-heldout recordings: `0`",
        "blocked_waiting_for_valid_sealed_split",
        "readiness 0/3 ready",
        "selector true-heldout protocol",
        "split20 full-live manifest",
        "resume calls 139",
        "Split policy optimization",
        "split_policy_optimization_from_existing_artifacts_no_live_calls",
        "Primary policy: `max20`",
        "Stretch policy: `max15`",
        "split20_deepseek_resume_after_top3_window_ids.txt",
        "101 parents",
        "resume export audit",
        "139 prompts",
        "split15 stretch re-export audit",
        "split15_stretch_reexport_audit_no_live_calls",
        "Export prompts: `178`",
        "split parent windows 为 58",
        "needs_new_recording_split",
        "review clean FP 可到 4/4",
        "LLM/Omni 不产生时间戳",
    ]
    ppt_required_fragments = [
        "Snapshot",
        "phasecard 8",
        "route none",
        "136/138 accept",
        "46.10s",
        "recording 8/8",
        "clean FP 4/4",
        "top3 wall 29.01s",
        "guard 104w harmful 0",
        runtime_audit_phrase,
        "Omni fusion",
        "split<=20",
        "latency ledger",
        "SLO 4/4",
        "latency risk tight1",
        "risk plan max20->max15",
        "time stats 9",
        "extractor 3",
        "timereceipt 1/6",
        "timing P0 384s",
        "retry budget 764",
        "token budget 1.66M",
        "recovery 8",
        "metrics 8",
        "schema 8",
        "scorecard 9",
        "handoff 7",
        "latclaim 8",
        "parallel 8w",
        "input ready 3/3",
        "cmd audit 3/3",
        "eligible 0",
        "receipt 0",
        "repair 3",
        "launcher 139",
        "skip 3/3",
        "retry 2x",
        "resume state 3/3",
        "env 14/14",
        "Omni96 calls",
        "readiness 0/3",
        "Agent plan",
        "postrun closure",
        "output audit",
        "scoring ready 0/5",
        "scoreplan 6",
        "scorelaunch 0",
        "scorereceipt 0",
        "scoreout 0",
        "preflight 0",
        "dag 10",
        "bundle 8",
        "promotion gate 0/8",
        "claimreceipt 3/6",
        "runbook 7",
        "heldout scan",
        "heldout scan 0/8",
        "split validation",
        "split20 resume audit",
        "stretch178",
        "split policy",
        "traceability 54/54",
        "P0 split20+heldout",
    ]
    claims = claims_manifest.get("claims", [])
    claim_ids = {claim.get("claim_id") for claim in claims}
    required_claim_ids = {
        "four_stage_realtime_route",
        "runtime_evidence_contract_clean",
        "rule_writeback_primary_correction",
        "selector_generalization_positive",
        "runtime_safe_llm_guard_zero_harm",
        "split20_latency_path_limited",
        "guard_tuning_passthrough_safe",
        "boundary_auto_writeback_negative_control",
        "clean_high_llm_audit_agrees_with_rule",
        "llm_review_memory_not_timeline",
        "omni_fusion_label_only",
        "voiceprint_rule_handles_clean_high",
    }
    claim_scope_by_id = {claim.get("claim_id"): claim.get("contract_scope") for claim in claims}
    claim_strength_by_id = {claim.get("claim_id"): claim.get("claim_strength") for claim in claims}
    claims_missing_links = [
        claim.get("claim_id", "unknown")
        for claim in claims
        if not claim.get("source_artifacts") or not claim.get("validation_checks")
    ]
    claims_missing_sources = [
        {"claim_id": claim.get("claim_id", "unknown"), "source": source}
        for claim in claims
        for source in claim.get("source_artifacts", [])
        if not (ROOT / str(source)).exists()
    ]
    memory_rows = memory_replay.get("rows", [])
    memory_failed_scans = [
        row.get("patch_id", "unknown")
        for row in memory_rows
        if row.get("forbidden_runtime_token_scan") != "pass"
    ]
    replay_rows = replay_manifest.get("rows", [])
    replay_stage_ids = {row.get("stage") for row in replay_rows}
    required_replay_stages = {
        "fast_provisional",
        "rule_writeback",
        "llm_guard",
        "llm_review_signal",
        "omni_label",
        "memory_gate",
    }
    replay_missing_sources = [
        {"stage": row.get("stage", "unknown"), "source": source}
        for row in replay_rows
        for source in row.get("source_artifacts", [])
        if not (ROOT / str(source)).exists()
    ]
    replay_failed_scans = [
        row.get("stage", "unknown")
        for row in replay_rows
        if row.get("forbidden_runtime_token_scan") != "pass"
    ]
    replay_bad_audit_rows = [
        row.get("stage", "unknown")
        for row in replay_rows
        if row.get("runtime_audit_status") not in {"pass", "claim_manifest_only"}
    ]
    replay_writeback_rows = set(replay_manifest.get("summary", {}).get("writeback_rows", []))
    next_experiments = next_queue.get("next_experiments", [])
    next_experiment_ids = {item.get("experiment_id") for item in next_experiments}
    required_next_experiment_ids = {
        "full_split20_live_104w",
        "true_heldout_selector_recordings",
        "omni_fusion_expand_48_or_120",
        "memory_update_audit_replay",
        "end_to_end_runtime_replay_manifest",
    }
    next_p0_ids = {item.get("experiment_id") for item in next_experiments if item.get("priority") == "P0"}
    next_ready_count = sum(1 for item in next_experiments if str(item.get("status", "")).startswith("ready"))
    next_prepared_count = sum(1 for item in next_experiments if str(item.get("status", "")).startswith("prepared"))
    next_blocked_count = sum(1 for item in next_experiments if str(item.get("status", "")).startswith("blocked"))
    next_completed_count = sum(1 for item in next_experiments if str(item.get("status", "")).startswith("completed"))
    next_missing_target_claims = [
        item.get("experiment_id", "unknown")
        for item in next_experiments
        if item.get("target_claim") not in claim_ids
    ]
    next_missing_required_inputs = [
        {"experiment_id": item.get("experiment_id", "unknown"), "input": required_input}
        for item in next_experiments
        for required_input in item.get("required_inputs", [])
        if str(required_input).startswith("outputs/") and not (ROOT / str(required_input)).exists()
    ]
    do_not_deploy_ids = {item.get("item_id") for item in next_queue.get("do_not_deploy", [])}
    readiness_rows = live_readiness.get("runs", [])
    readiness_by_run = {row.get("run_id"): row for row in readiness_rows}
    readiness_env = live_readiness.get("environment", {})
    readiness_missing_sources = [
        {"run_id": row.get("run_id", "unknown"), "source": source}
        for row in readiness_rows
        for source in row.get("source_artifacts", [])
        if str(source).startswith("outputs/") and not (ROOT / str(source)).exists()
    ]
    selector_protocol_gates = {gate.get("gate_id"): gate for gate in selector_protocol.get("success_gates", [])}
    split20_manifest_summary = split20_manifest.get("summary", {})
    split20_resume_surface = split20_manifest.get("resume_surface", {})
    split20_manifest_source_missing = [
        source
        for source in split20_manifest.get("source_artifacts", [])
        if str(source).startswith("outputs/") and not (ROOT / str(source)).exists()
    ]
    split20_pending_file = ROOT / str(split20_resume_surface.get("pending_window_id_file", ""))
    split20_completed_file = ROOT / str(split20_resume_surface.get("completed_window_id_file", ""))
    split20_failed_file = ROOT / str(split20_resume_surface.get("failed_window_id_file", ""))
    split20_pending_ids = split20_pending_file.read_text(encoding="utf-8").splitlines() if split20_pending_file.exists() else []
    split20_completed_ids = split20_completed_file.read_text(encoding="utf-8").splitlines() if split20_completed_file.exists() else []
    split20_failed_ids = split20_failed_file.read_text(encoding="utf-8").splitlines() if split20_failed_file.exists() else []
    split20_export_summary = split20_export_audit.get("summary", {})
    split20_export_jsonl = ROOT / str(split20_export_audit.get("export_prompt_jsonl", ""))
    split15_export_summary = split15_export_audit.get("summary", {})
    split15_export_jsonl = ROOT / str(split15_export_audit.get("export_prompt_jsonl", ""))

    add_check(
        checks,
        "snapshot_contract",
        snapshot.get("runtime_contract") == "research_progress_snapshot_from_existing_artifacts",
        str(snapshot.get("runtime_contract")),
    )
    add_check(
        checks,
        "runtime_audit_pass",
        runtime_audit.get("overall_status") == "pass" and int(runtime_audit.get("runtime_blocking_count", -1)) == 0,
        json.dumps(runtime_audit, ensure_ascii=False, sort_keys=True),
    )
    add_check(
        checks,
        "four_stage_timeline",
        all(
            float(timeline.get(key, 0.0)) > 0.0
            for key in [
                "fast_avg_delay_sec",
                "rule_writeback_avg_delay_sec",
                "llm_guard_avg_delay_sec",
                "llm_review_avg_delay_sec",
            ]
        ),
        json.dumps(
            {
                "fast": timeline.get("fast_avg_delay_sec"),
                "rule": timeline.get("rule_writeback_avg_delay_sec"),
                "guard": timeline.get("llm_guard_avg_delay_sec"),
                "review": timeline.get("llm_review_avg_delay_sec"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "review_signal_no_timeline_block",
        int(review.get("review_cases", 0)) == 4
        and int(review.get("blocks_timeline_writeback", -1)) == 0
        and int(review.get("blocks_memory_update", 0)) == 4,
        json.dumps(review, ensure_ascii=False, sort_keys=True),
    )
    add_check(
        checks,
        "clean_high_llm_full_surface",
        int(clean_llm.get("patches", 0)) == 138
        and int(clean_llm.get("llm_accepts", 0)) == 136
        and int(clean_llm.get("llm_non_accepts", 0)) == 2,
        json.dumps(
            {
                "patches": clean_llm.get("patches"),
                "accepts": clean_llm.get("llm_accepts"),
                "non_accepts": clean_llm.get("llm_non_accepts"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "selector_holdout_bootstrap_generalization",
        int(selector_holdout.get("positive_splits", 0)) == 8
        and int(selector_holdout.get("splits", 0)) == 8
        and float(selector_holdout.get("weighted_heldout_der", 1.0)) < float(selector_holdout.get("weighted_fast_der", 0.0))
        and selector_holdout.get("fixed_policy") == "ratio_le_0.65_else_uncovered"
        and float(selector_bootstrap_best.get("prob_beats_fast", 0.0)) == 1.0
        and float(selector_bootstrap_best.get("delta_ci_low", 0.0)) > 0.0,
        json.dumps(
            {
                "positive": selector_holdout.get("positive_splits"),
                "splits": selector_holdout.get("splits"),
                "heldout_der": selector_holdout.get("weighted_heldout_der"),
                "fast_der": selector_holdout.get("weighted_fast_der"),
                "delta": selector_holdout.get("weighted_delta_vs_fast"),
                "fixed_policy": selector_holdout.get("fixed_policy"),
                "bootstrap_delta": selector_bootstrap_best.get("delta_vs_fast"),
                "bootstrap_ci_low": selector_bootstrap_best.get("delta_ci_low"),
                "bootstrap_prob": selector_bootstrap_best.get("prob_beats_fast"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "omni_fusion_label_only_contract",
        int(omni_fusion.get("windows", 0)) == 12
        and omni_fusion.get("high_sentinel_recall") == "1/4 (25.0%)"
        and omni_fusion.get("review_priority_high_recall") == "4/4 (100.0%)"
        and omni_fusion.get("clean_high_sentinel_fp") == "0/4 (0.0%)"
        and omni_fusion.get("clean_review_priority_fp") == "4/4 (100.0%)"
        and omni_fusion.get("runtime_contract") == "omni_audio_jsonl_joined_with_deployable_acoustic_proxy; no_timeline_writeback",
        json.dumps(
            {
                "windows": omni_fusion.get("windows"),
                "high_sentinel_recall": omni_fusion.get("high_sentinel_recall"),
                "review_high_recall": omni_fusion.get("review_priority_high_recall"),
                "clean_sentinel_fp": omni_fusion.get("clean_high_sentinel_fp"),
                "clean_review_fp": omni_fusion.get("clean_review_priority_fp"),
                "runtime_contract": omni_fusion.get("runtime_contract"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "omni_expansion_manifest_contract",
        omni_expansion.get("runtime_contract") == "omni_expansion_manifest_ready_no_live_calls_no_timeline_writeback"
        and int(omni_expansion.get("summary", {}).get("selected_windows", 0)) == 48
        and int(omni_expansion.get("summary", {}).get("anchor_smoke_windows", 0)) == 12
        and int(omni_expansion.get("summary", {}).get("new_runtime_proxy_windows", 0)) == 36
        and int(omni_expansion.get("summary", {}).get("planned_model_calls", 0)) == 96
        and int(omni_expansion.get("summary", {}).get("audio_missing_count", -1)) == 0
        and omni_expansion.get("summary", {}).get("live_call_status") == "not_run_manifest_only"
        and omni_expansion.get("summary", {}).get("no_timeline_writeback") is True,
        json.dumps(
            {
                "contract": omni_expansion.get("runtime_contract"),
                "summary": omni_expansion.get("summary", {}),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    omni_call_summary = omni_call_manifest.get("summary", {})
    add_check(
        checks,
        "omni48_live_call_manifest_contract",
        omni_call_manifest.get("runtime_contract") == "omni48_live_call_manifest_no_live_calls_label_only"
        and omni_call_manifest.get("status") == "pass"
        and int(omni_call_summary.get("window_count", 0)) == 48
        and int(omni_call_summary.get("call_count", 0)) == 96
        and int(omni_call_summary.get("expected_call_count", 0)) == 96
        and int(omni_call_summary.get("model_count", 0)) == 2
        and omni_call_summary.get("target_models") == ["qwen3.5-omni-flash", "qwen3.5-omni-plus-2026-03-15"]
        and int(omni_call_summary.get("anchor_smoke_calls", 0)) == 24
        and int(omni_call_summary.get("new_runtime_proxy_calls", 0)) == 72
        and int(omni_call_summary.get("audio_missing_count", -1)) == 0
        and int(omni_call_summary.get("live_calls_performed", -1)) == 0
        and omni_call_summary.get("writeback_right") == "label_only_no_timeline_writeback"
        and omni_call_summary.get("no_timeline_writeback") is True,
        json.dumps(
            {
                "contract": omni_call_manifest.get("runtime_contract"),
                "status": omni_call_manifest.get("status"),
                "summary": omni_call_summary,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "split20_live_top3_guard",
        int(split_top3.get("parent_windows", 0)) == 3
        and int(split_top3.get("harmful_accepts", -1)) == 0
        and float(split_top3.get("measured_wall_seconds", 999.0)) < float(split_top3.get("original_max_call_seconds", 0.0))
        and split_attempt.get("failure_type") == "AllocationQuota.FreeTierOnly"
        and int(qwen_split.get("harmful_accepts", -1)) == 0,
        json.dumps(
            {
                "top3_parents": split_top3.get("parent_windows"),
                "top3_wall": split_top3.get("measured_wall_seconds"),
                "original_max": split_top3.get("original_max_call_seconds"),
                "split_max": split_top3.get("split_max_call_seconds"),
                "harmful": split_top3.get("harmful_accepts"),
                "deepseek_top45_failure": split_attempt.get("failure_type"),
                "qwen_top45_wall": first_run_wall_seconds(qwen_split),
                "qwen_harmful": qwen_split.get("harmful_accepts"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "guard_tuning_review_passthrough_contract",
        guard_tuning.get("best_zero_harm_policy") == "combined_review0.5_keepfast0.5"
        and int(guard_tuning.get("best_zero_harm_conservative_recovered", 0)) == 260
        and int(tuned_guard.get("safe_accepts", 0)) == 323
        and int(tuned_guard.get("harmful_accepts", -1)) == 0
        and "keep_fast_supported_passthrough_exception" in str(tuned_guard.get("runtime_contract", ""))
        and float(tuned_boundary.get("avg_der", 0.0)) > float(tuned_recover.get("avg_der", 1.0)),
        json.dumps(
            {
                "policy": guard_tuning.get("best_zero_harm_policy"),
                "recovered": guard_tuning.get("best_zero_harm_conservative_recovered"),
                "safe_accepts": tuned_guard.get("safe_accepts"),
                "harmful": tuned_guard.get("harmful_accepts"),
                "contract": tuned_guard.get("runtime_contract"),
                "boundary_der": tuned_boundary.get("avg_der"),
                "recover_best_der": tuned_recover.get("avg_der"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "pareto_has_review_signal",
        "Dual Agent LLM review signal" in pareto_md and "blocks writeback 0 / memory 4" in pareto_md,
        "Dual Agent LLM review signal" if "Dual Agent LLM review signal" in pareto_md else "missing",
    )
    add_check(
        checks,
        "claims_manifest_contract",
        claims_manifest.get("runtime_contract") == "research_claims_manifest_from_existing_artifacts"
        and len(claims) >= 12
        and not (required_claim_ids - claim_ids)
        and not claims_missing_links
        and not claims_missing_sources
        and claim_scope_by_id.get("omni_fusion_label_only") == "runtime_pass_no_timeline_writeback"
        and claim_scope_by_id.get("selector_generalization_positive") == "dev_only_validation"
        and claim_strength_by_id.get("boundary_auto_writeback_negative_control") == "do_not_deploy",
        json.dumps(
            {
                "contract": claims_manifest.get("runtime_contract"),
                "claims": len(claims),
                "missing_ids": sorted(required_claim_ids - claim_ids),
                "missing_links": claims_missing_links,
                "missing_sources": claims_missing_sources,
                "omni_scope": claim_scope_by_id.get("omni_fusion_label_only"),
                "selector_scope": claim_scope_by_id.get("selector_generalization_positive"),
                "boundary_strength": claim_strength_by_id.get("boundary_auto_writeback_negative_control"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    latency_summary = latency_ledger.get("summary", {})
    latency_rows = {row.get("stage_id"): row for row in latency_ledger.get("rows", [])}
    add_check(
        checks,
        "runtime_latency_budget_ledger_contract",
        latency_ledger.get("runtime_contract") == "runtime_latency_budget_ledger_from_existing_artifacts"
        and latency_ledger.get("status") == "pass"
        and int(latency_summary.get("row_count", 0)) == 10
        and int(latency_summary.get("claim_now_rows", 0)) == 4
        and int(latency_summary.get("smoke_only_rows", 0)) == 2
        and int(latency_summary.get("pending_or_blocked_rows", 0)) == 2
        and int(latency_summary.get("offline_budget_rows", 0)) == 1
        and int(latency_summary.get("live_calls_performed_by_builder", -1)) == 0
        and not latency_summary.get("failed_claim_rows")
        and latency_rows.get("fast_first_output", {}).get("target_status") == "pass"
        and latency_rows.get("rule_writeback", {}).get("target_status") == "pass"
        and latency_rows.get("runtime_safe_llm_guard", {}).get("target_status") == "pass"
        and latency_rows.get("split20_deepseek_full_resume", {}).get("claim_status") == "blocked_by_quota_or_missing_resume"
        and latency_rows.get("omni48_label_only_live", {}).get("claim_status") == "pending_omni48_live_outputs",
        json.dumps(
            {
                "contract": latency_ledger.get("runtime_contract"),
                "status": latency_ledger.get("status"),
                "summary": latency_summary,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    latency_slo_summary = latency_slo_audit.get("summary", {})
    latency_slo_rows = {row.get("stage_id"): row for row in latency_slo_audit.get("rows", [])}
    add_check(
        checks,
        "stage_latency_slo_audit_contract",
        latency_slo_audit.get("runtime_contract") == "stage_latency_slo_audit_from_latency_ledger_no_live_calls"
        and latency_slo_audit.get("source_ledger_contract") == "runtime_latency_budget_ledger_from_existing_artifacts"
        and latency_slo_audit.get("status") == "pass"
        and int(latency_slo_summary.get("row_count", 0)) == 10
        and int(latency_slo_summary.get("claim_now_slo_rows", 0)) == 4
        and int(latency_slo_summary.get("claim_now_slo_pass", 0)) == 4
        and int(latency_slo_summary.get("claim_now_slo_fail", -1)) == 0
        and int(latency_slo_summary.get("smoke_rows", 0)) == 2
        and int(latency_slo_summary.get("pending_or_blocked_rows", 0)) == 2
        and float(latency_slo_summary.get("min_claim_p95_margin_seconds", 0.0)) == 0.555
        and float(latency_slo_summary.get("guard_p95_margin_seconds", 0.0)) == 1.151
        and int(latency_slo_summary.get("live_calls_performed_by_builder", -1)) == 0
        and latency_slo_summary.get("no_new_metric_claim") is True
        and latency_slo_rows.get("fast_first_output", {}).get("slo_class") == "claim_now_slo_pass"
        and latency_slo_rows.get("rule_writeback", {}).get("slo_class") == "claim_now_slo_pass"
        and latency_slo_rows.get("runtime_safe_llm_guard", {}).get("slo_class") == "claim_now_slo_pass"
        and latency_slo_rows.get("split20_deepseek_full_resume", {}).get("slo_class") == "pending_or_blocked"
        and latency_slo_rows.get("omni48_label_only_live", {}).get("slo_class") == "pending_or_blocked",
        json.dumps(
            {
                "contract": latency_slo_audit.get("runtime_contract"),
                "status": latency_slo_audit.get("status"),
                "summary": latency_slo_summary,
                "classes": {key: value.get("slo_class") for key, value in latency_slo_rows.items()},
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    latency_risk_summary = latency_risk_margin.get("summary", {})
    latency_risk_rows = {row.get("stage_id"): row for row in latency_risk_margin.get("rows", [])}
    add_check(
        checks,
        "latency_risk_margin_audit_contract",
        latency_risk_margin.get("runtime_contract") == "latency_risk_margin_audit_from_slo_no_live_calls"
        and latency_risk_margin.get("status") == "pass"
        and latency_risk_margin.get("source_contracts", {}).get("slo_audit") == "stage_latency_slo_audit_from_latency_ledger_no_live_calls"
        and latency_risk_margin.get("source_contracts", {}).get("latency_ledger") == "runtime_latency_budget_ledger_from_existing_artifacts"
        and latency_risk_margin.get("source_contracts", {}).get("promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and int(latency_risk_summary.get("row_count", 0)) == 10
        and int(latency_risk_summary.get("claim_now_rows", 0)) == 4
        and int(latency_risk_summary.get("tight_margin_rows", 0)) == 1
        and int(latency_risk_summary.get("watch_rows", -1)) == 0
        and int(latency_risk_summary.get("blocked_rows", 0)) == 2
        and int(latency_risk_summary.get("non_claimable_rows", 0)) == 4
        and not latency_risk_summary.get("failed_claim_rows")
        and latency_risk_summary.get("tight_margin_stage_ids") == ["runtime_safe_llm_guard"]
        and sorted(latency_risk_summary.get("blocked_stage_ids", [])) == ["omni48_label_only_live", "split20_deepseek_full_resume"]
        and latency_risk_summary.get("guard_risk_level") == "tight_margin"
        and float(latency_risk_summary.get("guard_p95_margin_seconds", 0.0)) == 1.151
        and float(latency_risk_summary.get("guard_p95_margin_ratio", 0.0)) == 0.0177
        and int(latency_risk_summary.get("post_live_ready_to_promote", -1)) == 0
        and int(latency_risk_summary.get("live_ready_runs", -1)) == 0
        and int(latency_risk_summary.get("live_calls_performed_by_builder", -1)) == 0
        and latency_risk_summary.get("no_new_metric_claim") is True
        and latency_risk_rows.get("fast_first_output", {}).get("risk_level") == "comfortable"
        and latency_risk_rows.get("runtime_safe_llm_guard", {}).get("risk_level") == "tight_margin"
        and latency_risk_rows.get("split20_simulated_policy", {}).get("risk_level") == "planning_only"
        and latency_risk_rows.get("split20_qwen_backup_top45", {}).get("risk_level") == "fallback_only"
        and latency_risk_rows.get("omni48_label_only_live", {}).get("risk_level") == "blocked",
        json.dumps(
            {
                "contract": latency_risk_margin.get("runtime_contract"),
                "status": latency_risk_margin.get("status"),
                "summary": latency_risk_summary,
                "risk_levels": {key: value.get("risk_level") for key, value in sorted(latency_risk_rows.items())},
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    latency_mitigation_summary = latency_risk_mitigation.get("summary", {})
    latency_mitigation_actions = {row.get("action_id"): row for row in latency_risk_mitigation.get("actions", [])}
    required_latency_mitigation_actions = {
        "preserve_current_slo_with_guard_risk_tag",
        "deepseek_max20_resume_primary",
        "max15_stretch_reexport",
        "qwen_backup_not_primary_latency",
        "omni48_label_only_not_guard_latency",
        "selector_true_heldout_not_latency_mitigation",
    }
    add_check(
        checks,
        "latency_risk_mitigation_plan_contract",
        latency_risk_mitigation.get("runtime_contract") == "latency_risk_mitigation_plan_no_live_calls"
        and latency_risk_mitigation.get("status") == "blocked_waiting_for_primary_live_resume"
        and latency_risk_mitigation.get("source_contracts", {}).get("latency_risk_margin") == "latency_risk_margin_audit_from_slo_no_live_calls"
        and latency_risk_mitigation.get("source_contracts", {}).get("split_policy_optimization") == "split_policy_optimization_from_existing_artifacts_no_live_calls"
        and latency_risk_mitigation.get("source_contracts", {}).get("split15_stretch_reexport_audit") == "split15_stretch_reexport_audit_no_live_calls"
        and int(latency_mitigation_summary.get("action_count", 0)) == 6
        and int(latency_mitigation_summary.get("p0_action_count", 0)) == 2
        and int(latency_mitigation_summary.get("active_current_claim_count", 0)) == 1
        and int(latency_mitigation_summary.get("mitigation_ready_count", -1)) == 0
        and int(latency_mitigation_summary.get("blocked_action_count", 0)) == 3
        and int(latency_mitigation_summary.get("fallback_only_count", 0)) == 1
        and int(latency_mitigation_summary.get("stretch_candidate_count", 0)) == 1
        and latency_mitigation_summary.get("guard_risk_level") == "tight_margin"
        and float(latency_mitigation_summary.get("guard_p95_margin_seconds", 0.0)) == 1.151
        and float(latency_mitigation_summary.get("guard_p95_margin_ratio", 0.0)) == 0.0177
        and latency_mitigation_summary.get("primary_policy") == "max20"
        and int(latency_mitigation_summary.get("primary_resume_calls", 0)) == 139
        and float(latency_mitigation_summary.get("primary_simulated_p95_call_seconds", 0.0)) == 21.358
        and float(latency_mitigation_summary.get("primary_token_multiplier", 0.0)) == 1.118
        and latency_mitigation_summary.get("stretch_policy") == "max15"
        and int(latency_mitigation_summary.get("stretch_calls", 0)) == 178
        and float(latency_mitigation_summary.get("stretch_simulated_p95_call_seconds", 0.0)) == 18.047
        and float(latency_mitigation_summary.get("stretch_token_multiplier", 0.0)) == 1.182
        and latency_mitigation_summary.get("stretch_requires_reexport") is True
        and int(latency_mitigation_summary.get("stretch_call_delta", 0)) == 31
        and float(latency_mitigation_summary.get("stretch_p95_call_gain_seconds", 0.0)) == 3.311
        and latency_mitigation_summary.get("stretch_export_status") == "pass"
        and int(latency_mitigation_summary.get("stretch_export_prompts", 0)) == 178
        and int(latency_mitigation_summary.get("stretch_export_parent_windows", 0)) == 104
        and int(latency_mitigation_summary.get("stretch_export_split_parent_windows", 0)) == 58
        and latency_mitigation_summary.get("stretch_export_prompt_jsonl") == "outputs/research_progress_snapshot/split15_stretch_reexport_prompts.jsonl"
        and int(latency_mitigation_summary.get("stretch_export_live_calls_performed", -1)) == 0
        and int(latency_mitigation_summary.get("post_live_ready_to_promote", -1)) == 0
        and int(latency_mitigation_summary.get("live_ready_runs", -1)) == 0
        and int(latency_mitigation_summary.get("runbook_p0_planned_live_calls", 0)) == 139
        and int(latency_mitigation_summary.get("missing_output_surfaces", 0)) == 3
        and int(latency_mitigation_summary.get("live_calls_performed_by_builder", -1)) == 0
        and latency_mitigation_summary.get("no_secret_values_written") is True
        and latency_mitigation_summary.get("no_new_metric_claim") is True
        and required_latency_mitigation_actions.issubset(latency_mitigation_actions)
        and latency_mitigation_actions.get("deepseek_max20_resume_primary", {}).get("claim_boundary") == "blocked_until_full_surface_live_output_and_scoring"
        and latency_mitigation_actions.get("max15_stretch_reexport", {}).get("claim_boundary") == "stretch_plan_no_new_metric_claim"
        and latency_mitigation_actions.get("qwen_backup_not_primary_latency", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and latency_mitigation_actions.get("omni48_label_only_not_guard_latency", {}).get("claim_boundary") == "label_only_no_timeline_writeback",
        json.dumps(
            {
                "contract": latency_risk_mitigation.get("runtime_contract"),
                "status": latency_risk_mitigation.get("status"),
                "summary": latency_mitigation_summary,
                "missing_actions": sorted(required_latency_mitigation_actions - set(latency_mitigation_actions)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "runtime_replay_manifest_contract",
        replay_manifest.get("runtime_contract") == "end_to_end_runtime_replay_manifest_no_eval_context"
        and len(replay_rows) == 6
        and not (required_replay_stages - replay_stage_ids)
        and replay_writeback_rows == {"fast_provisional", "rule_writeback"}
        and not replay_manifest.get("summary", {}).get("failed_rows")
        and not replay_missing_sources
        and not replay_failed_scans
        and not replay_bad_audit_rows,
        json.dumps(
            {
                "contract": replay_manifest.get("runtime_contract"),
                "rows": len(replay_rows),
                "missing_stages": sorted(required_replay_stages - replay_stage_ids),
                "writeback_rows": sorted(replay_writeback_rows),
                "failed_rows": replay_manifest.get("summary", {}).get("failed_rows"),
                "missing_sources": replay_missing_sources,
                "failed_scans": replay_failed_scans,
                "bad_audit_rows": replay_bad_audit_rows,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "memory_update_replay_contract",
        memory_replay.get("runtime_contract") == "memory_update_replay_review_signal_blocks_memory_only"
        and int(memory_replay.get("summary", {}).get("review_cases", 0)) == 4
        and int(memory_replay.get("summary", {}).get("memory_candidates_before", 0)) == 4
        and int(memory_replay.get("summary", {}).get("memory_updates_allowed_after", -1)) == 0
        and int(memory_replay.get("summary", {}).get("memory_updates_blocked", 0)) == 4
        and int(memory_replay.get("summary", {}).get("timeline_writebacks_preserved", 0)) == 4
        and not memory_replay.get("summary", {}).get("failed_rows")
        and not memory_failed_scans,
        json.dumps(
            {
                "contract": memory_replay.get("runtime_contract"),
                "summary": memory_replay.get("summary", {}),
                "failed_scans": memory_failed_scans,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "next_experiment_queue_contract",
        next_queue.get("runtime_contract") == "research_next_experiment_queue_from_validated_claims"
        and len(next_experiments) >= 5
        and not (required_next_experiment_ids - next_experiment_ids)
        and {"full_split20_live_104w", "true_heldout_selector_recordings"}.issubset(next_p0_ids)
        and next_blocked_count >= 1
        and next_prepared_count >= 1
        and next_completed_count >= 2
        and "boundary_auto_writeback" in do_not_deploy_ids
        and not next_missing_target_claims
        and not next_missing_required_inputs,
        json.dumps(
            {
                "contract": next_queue.get("runtime_contract"),
                "experiments": len(next_experiments),
                "missing_ids": sorted(required_next_experiment_ids - next_experiment_ids),
                "p0_ids": sorted(next_p0_ids),
                "ready_count": next_ready_count,
                "prepared_count": next_prepared_count,
                "blocked_count": next_blocked_count,
                "completed_count": next_completed_count,
                "do_not_deploy": sorted(do_not_deploy_ids),
                "missing_target_claims": next_missing_target_claims,
                "missing_required_inputs": next_missing_required_inputs,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "live_run_readiness_contract",
        live_readiness.get("runtime_contract") == "live_run_readiness_non_secret_no_live_calls"
        and live_readiness.get("secret_policy") == "env_presence_only_no_secret_values_written"
        and live_readiness.get("summary", {}).get("non_secret") is True
        and int(live_readiness.get("summary", {}).get("live_calls_performed", -1)) == 0
        and int(live_readiness.get("summary", {}).get("run_count", 0)) == 3
        and int(live_readiness.get("summary", {}).get("p0_blocked_count", 0)) >= 1
        and readiness_env.get("config_defaults_not_counted_as_credentials") is True
        and int(readiness_by_run.get("omni48_live", {}).get("planned_calls", 0)) == 96
        and int(readiness_by_run.get("omni48_live", {}).get("planned_windows", 0)) == 48
        and int(readiness_by_run.get("split20_deepseek_full", {}).get("planned_calls", 0)) == 139
        and int(readiness_by_run.get("split20_deepseek_full", {}).get("planned_windows", 0)) == 101
        and "--window-id-file outputs/research_progress_snapshot/split20_deepseek_resume_after_top3_window_ids.txt"
        in readiness_by_run.get("split20_deepseek_full", {}).get("run_command", "")
        and "--skip-existing-output" in readiness_by_run.get("split20_deepseek_full", {}).get("run_command", "")
        and "--max-call-attempts 2" in readiness_by_run.get("split20_deepseek_full", {}).get("run_command", "")
        and "--retry-backoff-seconds 2.0" in readiness_by_run.get("split20_deepseek_full", {}).get("run_command", "")
        and any(
            "AllocationQuota.FreeTierOnly" in str(blocker)
            for blocker in readiness_by_run.get("split20_deepseek_full", {}).get("blockers", [])
        )
        and not readiness_missing_sources,
        json.dumps(
            {
                "contract": live_readiness.get("runtime_contract"),
                "summary": live_readiness.get("summary", {}),
                "dashscope_env": readiness_env.get("dashscope_env_present", {}),
                "omni": {
                    "status": readiness_by_run.get("omni48_live", {}).get("status"),
                    "calls": readiness_by_run.get("omni48_live", {}).get("planned_calls"),
                    "windows": readiness_by_run.get("omni48_live", {}).get("planned_windows"),
                },
                "deepseek": {
                    "status": readiness_by_run.get("split20_deepseek_full", {}).get("status"),
                    "calls": readiness_by_run.get("split20_deepseek_full", {}).get("planned_calls"),
                    "windows": readiness_by_run.get("split20_deepseek_full", {}).get("planned_windows"),
                    "blockers": readiness_by_run.get("split20_deepseek_full", {}).get("blockers", []),
                },
                "missing_sources": readiness_missing_sources,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_agent_summary = live_agent_plan.get("summary", {})
    live_agent_steps = live_agent_plan.get("steps", [])
    add_check(
        checks,
        "live_agent_execution_plan_contract",
        live_agent_plan.get("runtime_contract") == "live_agent_execution_plan_no_live_calls"
        and live_agent_plan.get("status") == "pass"
        and int(live_agent_summary.get("step_count", 0)) == 5
        and int(live_agent_summary.get("live_step_count", 0)) == 3
        and int(live_agent_summary.get("planned_live_calls", 0)) == 382
        and int(live_agent_summary.get("p0_planned_live_calls", 0)) == 139
        and int(live_agent_summary.get("deepseek_resume_calls", 0)) == 139
        and int(live_agent_summary.get("omni_label_only_calls", 0)) == 96
        and int(live_agent_summary.get("qwen_backup_calls", 0)) == 147
        and int(live_agent_summary.get("live_calls_performed", -1)) == 0
        and live_agent_summary.get("no_secret_values_written") is True
        and not live_agent_plan.get("missing_sources")
        and {step.get("step_id") for step in live_agent_steps}
        == {
            "credential_preflight",
            "split20_deepseek_resume_after_top3",
            "split20_qwen_backup_full_surface",
            "omni48_label_only_live",
            "postrun_refresh_and_validation",
        },
        json.dumps(
            {
                "contract": live_agent_plan.get("runtime_contract"),
                "status": live_agent_plan.get("status"),
                "missing_sources": live_agent_plan.get("missing_sources"),
                "summary": live_agent_summary,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_postrun_summary = live_postrun_closure.get("summary", {})
    live_postrun_surfaces = {row.get("surface_id"): row for row in live_postrun_closure.get("surfaces", [])}
    add_check(
        checks,
        "live_postrun_metrics_closure_contract",
        live_postrun_closure.get("runtime_contract") == "live_postrun_metrics_closure_no_live_calls"
        and live_postrun_closure.get("status") == "pending_live_outputs"
        and int(live_postrun_summary.get("split20_expected_calls", 0)) == 147
        and int(live_postrun_summary.get("split20_expected_parent_windows", 0)) == 104
        and int(live_postrun_summary.get("deepseek_success_calls", 0)) == 8
        and int(live_postrun_summary.get("deepseek_resume_expected_calls", 0)) == 139
        and int(live_postrun_summary.get("deepseek_resume_successful_calls", -1)) == 0
        and int(live_postrun_summary.get("deepseek_quota_failed_calls", 0)) == 4
        and int(live_postrun_summary.get("omni48_expected_calls", 0)) == 96
        and int(live_postrun_summary.get("omni48_successful_calls", -1)) == 0
        and live_postrun_summary.get("split20_latency_claim_status") == "blocked_by_quota_or_missing_resume"
        and live_postrun_summary.get("omni48_latency_claim_status") == "pending_omni48_live_outputs"
        and int(live_postrun_summary.get("live_calls_performed_by_builder", -1)) == 0
        and {"deepseek_resume_after_top3", "qwen_full_backup", "omni48_label_only"}.issubset(live_postrun_surfaces)
        and live_postrun_surfaces.get("deepseek_resume_after_top3", {}).get("status") == "pending_live_output"
        and live_postrun_surfaces.get("omni48_label_only", {}).get("status") == "pending_live_output",
        json.dumps(
            {
                "contract": live_postrun_closure.get("runtime_contract"),
                "status": live_postrun_closure.get("status"),
                "summary": live_postrun_summary,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_output_summary = live_output_audit.get("summary", {})
    live_output_surfaces = {row.get("surface_id"): row for row in live_output_audit.get("surfaces", [])}
    add_check(
        checks,
        "live_output_audit_contract",
        live_output_audit.get("runtime_contract") == "live_output_audit_no_live_calls"
        and live_output_audit.get("status") == "pending_live_outputs"
        and int(live_output_summary.get("surface_count", 0)) == 3
        and int(live_output_summary.get("expected_live_calls", 0)) == 382
        and int(live_output_summary.get("observed_live_output_rows", -1)) == 0
        and int(live_output_summary.get("successful_live_output_rows", -1)) == 0
        and int(live_output_summary.get("missing_output_surfaces", 0)) == 3
        and int(live_output_summary.get("partial_or_invalid_surfaces", -1)) == 0
        and int(live_output_summary.get("claim_ready_surfaces", -1)) == 0
        and int(live_output_summary.get("live_calls_performed_by_auditor", -1)) == 0
        and {"deepseek_resume_after_top3", "qwen_full_backup", "omni48_label_only"}.issubset(live_output_surfaces)
        and int(live_output_surfaces.get("deepseek_resume_after_top3", {}).get("expected_calls", 0)) == 139
        and int(live_output_surfaces.get("qwen_full_backup", {}).get("expected_calls", 0)) == 147
        and int(live_output_surfaces.get("omni48_label_only", {}).get("expected_calls", 0)) == 96
        and all(row.get("claim_gate") == "blocked_missing_output" for row in live_output_surfaces.values()),
        json.dumps(
            {
                "contract": live_output_audit.get("runtime_contract"),
                "status": live_output_audit.get("status"),
                "summary": live_output_summary,
                "surfaces": {
                    key: {
                        "expected_calls": value.get("expected_calls"),
                        "status": value.get("status"),
                        "claim_gate": value.get("claim_gate"),
                    }
                    for key, value in live_output_surfaces.items()
                },
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_scoring_summary = live_scoring_readiness.get("summary", {})
    live_scoring_steps = {row.get("scoring_id"): row for row in live_scoring_readiness.get("scoring_steps", [])}
    required_scoring_ids = {
        "deepseek_resume_safety",
        "deepseek_full_split20_comparison",
        "qwen_full_backup_safety",
        "qwen_full_backup_comparison",
        "omni48_label_summary",
    }
    add_check(
        checks,
        "live_scoring_readiness_contract",
        live_scoring_readiness.get("runtime_contract") == "live_scoring_readiness_no_live_calls"
        and live_scoring_readiness.get("status") == "blocked_waiting_live_outputs"
        and int(live_scoring_summary.get("scoring_step_count", 0)) == 5
        and int(live_scoring_summary.get("ready_to_score_steps", -1)) == 0
        and int(live_scoring_summary.get("blocked_steps", 0)) == 5
        and int(live_scoring_summary.get("p0_scoring_steps", 0)) == 2
        and int(live_scoring_summary.get("unique_live_output_calls", 0)) == 382
        and int(live_scoring_summary.get("deepseek_resume_expected_calls", 0)) == 139
        and int(live_scoring_summary.get("qwen_full_expected_calls", 0)) == 147
        and int(live_scoring_summary.get("omni48_expected_calls", 0)) == 96
        and int(live_scoring_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_scoring_summary.get("no_scoring_commands_executed") is True
        and required_scoring_ids.issubset(live_scoring_steps)
        and live_scoring_steps.get("deepseek_resume_safety", {}).get("priority") == "P0"
        and live_scoring_steps.get("deepseek_full_split20_comparison", {}).get("priority") == "P0"
        and "--batch-jsonl outputs/runtime_safe_llm_window_batch/deepseek_proxy_high_risk_104w_split20_resume_after_top3.jsonl"
        in live_scoring_steps.get("deepseek_resume_safety", {}).get("scoring_command", "")
        and "scripts/analysis/summarize_split_llm_runs.py"
        in live_scoring_steps.get("deepseek_full_split20_comparison", {}).get("scoring_command", "")
        and "scripts/analysis/summarize_omni_window_batch.py"
        in live_scoring_steps.get("omni48_label_summary", {}).get("scoring_command", ""),
        json.dumps(
            {
                "contract": live_scoring_readiness.get("runtime_contract"),
                "status": live_scoring_readiness.get("status"),
                "summary": live_scoring_summary,
                "scoring_ids": sorted(live_scoring_steps),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    scoring_execution_summary = post_live_scoring_execution_plan.get("summary", {})
    scoring_execution_rows = {
        row.get("scoring_execution_id"): row
        for row in post_live_scoring_execution_plan.get("rows", [])
    }
    required_scoring_execution_ids = {
        "deepseek_resume_safety_score",
        "deepseek_full_split20_comparison_score",
        "omni48_label_summary_score",
        "qwen_full_backup_safety_score",
        "qwen_full_backup_comparison_score",
        "promotion_refresh_validation",
    }
    add_check(
        checks,
        "post_live_scoring_execution_plan_contract",
        post_live_scoring_execution_plan.get("runtime_contract") == "post_live_scoring_execution_plan_no_live_calls"
        and post_live_scoring_execution_plan.get("status") == "blocked_waiting_live_outputs"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("live_metric_extraction_contract") == "live_metric_extraction_contract_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("live_output_schema_contract") == "live_output_schema_contract_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("post_live_acceptance_scorecard") == "post_live_acceptance_scorecard_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("post_live_latency_claim_matrix") == "post_live_latency_claim_matrix_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_scoring_execution_plan.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(scoring_execution_summary.get("scoring_execution_steps", 0)) == 6
        and int(scoring_execution_summary.get("p0_execution_steps", 0)) == 3
        and int(scoring_execution_summary.get("p1_execution_steps", 0)) == 3
        and int(scoring_execution_summary.get("blocked_execution_steps", 0)) == 6
        and int(scoring_execution_summary.get("ready_execution_steps", -1)) == 0
        and int(scoring_execution_summary.get("scoring_commands", 0)) == 6
        and int(scoring_execution_summary.get("p0_scoring_steps", 0)) == 2
        and int(scoring_execution_summary.get("readiness_ready_to_score_steps", -1)) == 0
        and int(scoring_execution_summary.get("missing_output_surfaces", 0)) == 3
        and int(scoring_execution_summary.get("expected_live_calls", 0)) == 382
        and int(scoring_execution_summary.get("metric_contract_count", 0)) == 8
        and int(scoring_execution_summary.get("schema_contract_count", 0)) == 8
        and int(scoring_execution_summary.get("scorecard_rows", 0)) == 9
        and int(scoring_execution_summary.get("latency_claim_rows", 0)) == 8
        and int(scoring_execution_summary.get("ready_to_promote_count", -1)) == 0
        and int(scoring_execution_summary.get("traceability_rows", 0)) == 54
        and int(scoring_execution_summary.get("live_calls_performed_by_builder", -1)) == 0
        and scoring_execution_summary.get("no_scoring_commands_executed") is True
        and scoring_execution_summary.get("no_secret_values_written") is True
        and scoring_execution_summary.get("no_new_metric_claim") is True
        and required_scoring_execution_ids.issubset(scoring_execution_rows)
        and scoring_execution_rows.get("deepseek_resume_safety_score", {}).get("step_order") == 1
        and scoring_execution_rows.get("deepseek_full_split20_comparison_score", {}).get("step_order") == 2
        and scoring_execution_rows.get("deepseek_full_split20_comparison_score", {}).get("claim_boundary") == "required_before_full_surface_latency_claim"
        and scoring_execution_rows.get("omni48_label_summary_score", {}).get("claim_boundary") == "label_only_no_timeline_writeback"
        and scoring_execution_rows.get("qwen_full_backup_comparison_score", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and scoring_execution_rows.get("promotion_refresh_validation", {}).get("claim_boundary") == "required_before_report_ppt_claim_promotion"
        and "scripts/analysis/analyze_runtime_safe_llm_guard.py" in scoring_execution_rows.get("deepseek_resume_safety_score", {}).get("command", "")
        and "scripts/analysis/summarize_split_llm_runs.py" in scoring_execution_rows.get("deepseek_full_split20_comparison_score", {}).get("command", "")
        and "scripts/misc/refresh_latest_research_artifacts.py" in scoring_execution_rows.get("promotion_refresh_validation", {}).get("command", ""),
        json.dumps(
            {
                "contract": post_live_scoring_execution_plan.get("runtime_contract"),
                "status": post_live_scoring_execution_plan.get("status"),
                "summary": scoring_execution_summary,
                "missing_ids": sorted(required_scoring_execution_ids - set(scoring_execution_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    scoring_launcher_summary = post_live_scoring_launcher.get("summary", {})
    scoring_launcher_rows = {
        row.get("launcher_step_id"): row
        for row in post_live_scoring_launcher.get("rows", [])
    }
    add_check(
        checks,
        "post_live_scoring_launcher_contract",
        post_live_scoring_launcher.get("runtime_contract") == "post_live_scoring_launcher_dry_run_no_scoring_calls"
        and post_live_scoring_launcher.get("secret_policy") == "commands_scanned_no_secret_values_written"
        and post_live_scoring_launcher.get("status") == "dry_run_blocked_waiting_live_outputs_or_execute_flag"
        and post_live_scoring_launcher.get("source_contracts", {}).get("post_live_scoring_execution_plan") == "post_live_scoring_execution_plan_no_live_calls"
        and post_live_scoring_launcher.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and post_live_scoring_launcher.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and post_live_scoring_launcher.get("source_contracts", {}).get("live_output_repair_plan") == "live_output_repair_plan_no_live_calls_no_scoring"
        and post_live_scoring_launcher.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and post_live_scoring_launcher.get("execution_blockers") == [
            "dry_run_only_requires_execute_scoring_flag",
            "no_ready_scoring_commands",
        ]
        and int(scoring_launcher_summary.get("launcher_steps", 0)) == 6
        and int(scoring_launcher_summary.get("available_scoring_rows", 0)) == 6
        and int(scoring_launcher_summary.get("ready_scoring_rows", -1)) == 0
        and int(scoring_launcher_summary.get("selected_scoring_rows", -1)) == 0
        and int(scoring_launcher_summary.get("p0_ready_scoring_rows", -1)) == 0
        and int(scoring_launcher_summary.get("p1_ready_scoring_rows", -1)) == 0
        and scoring_launcher_summary.get("scoring_scope") == "p0"
        and scoring_launcher_summary.get("execute_scoring") is False
        and scoring_launcher_summary.get("execution_allowed") is False
        and int(scoring_launcher_summary.get("execution_blocker_count", 0)) == 2
        and int(scoring_launcher_summary.get("executed_scoring_rows", -1)) == 0
        and int(scoring_launcher_summary.get("passed_scoring_rows", -1)) == 0
        and int(scoring_launcher_summary.get("failed_scoring_rows", -1)) == 0
        and int(scoring_launcher_summary.get("readiness_ready_to_score_steps", -1)) == 0
        and int(scoring_launcher_summary.get("output_missing_surfaces", 0)) == 3
        and int(scoring_launcher_summary.get("repair_scoring_ready_rows", -1)) == 0
        and int(scoring_launcher_summary.get("traceability_rows", 0)) == 54
        and scoring_launcher_summary.get("scoring_execute_record_exists") is False
        and scoring_launcher_summary.get("scoring_execute_record_path") == "outputs/research_progress_snapshot/post_live_scoring_launcher_execute_latest.json"
        and scoring_launcher_summary.get("latest_scoring_execute_status") == ""
        and scoring_launcher_summary.get("latest_scoring_execute_runtime_contract") == ""
        and scoring_launcher_summary.get("latest_scoring_execute_scope") == ""
        and int(scoring_launcher_summary.get("latest_scoring_executed_rows", -1)) == 0
        and int(scoring_launcher_summary.get("latest_scoring_passed_rows", -1)) == 0
        and int(scoring_launcher_summary.get("latest_scoring_failed_rows", -1)) == 0
        and int(scoring_launcher_summary.get("live_calls_performed_by_launcher", -1)) == 0
        and scoring_launcher_summary.get("no_live_calls_performed") is True
        and scoring_launcher_summary.get("no_scoring_commands_executed") is True
        and scoring_launcher_summary.get("no_secret_values_written") is True
        and scoring_launcher_summary.get("no_new_metric_claim") is True
        and required_scoring_execution_ids.issubset(scoring_launcher_rows)
        and all(row.get("eligible") is False for row in scoring_launcher_rows.values())
        and all(row.get("selected") is False for row in scoring_launcher_rows.values())
        and all(row.get("status") == "blocked_waiting_prerequisite" for row in scoring_launcher_rows.values()),
        json.dumps(
            {
                "contract": post_live_scoring_launcher.get("runtime_contract"),
                "status": post_live_scoring_launcher.get("status"),
                "summary": scoring_launcher_summary,
                "missing_ids": sorted(required_scoring_execution_ids - set(scoring_launcher_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    scoring_receipt_summary = post_live_scoring_receipt.get("summary", {})
    scoring_receipt_rows = {
        row.get("receipt_id"): row
        for row in post_live_scoring_receipt.get("rows", [])
    }
    required_scoring_receipt_ids = {
        "scoring_execute_record_presence",
        "scoring_scope_alignment",
        "scoring_command_result_receipt",
        "scoring_output_audit_after_execute",
        "time_metric_extractor_after_scoring",
        "promotion_preflight_after_scoring",
    }
    add_check(
        checks,
        "post_live_scoring_receipt_audit_contract",
        post_live_scoring_receipt.get("runtime_contract") == "post_live_scoring_receipt_audit_no_live_or_scoring_calls_no_secret_values"
        and post_live_scoring_receipt.get("status") == "blocked_no_scoring_receipt_or_outputs"
        and post_live_scoring_receipt.get("source_contracts", {}).get("post_live_scoring_launcher") == "post_live_scoring_launcher_dry_run_no_scoring_calls"
        and post_live_scoring_receipt.get("source_contracts", {}).get("post_live_scoring_execution_plan") == "post_live_scoring_execution_plan_no_live_calls"
        and post_live_scoring_receipt.get("source_contracts", {}).get("post_live_scoring_output_audit") == "post_live_scoring_output_audit_no_scoring_calls"
        and post_live_scoring_receipt.get("source_contracts", {}).get("post_live_time_metric_extractor") == "post_live_time_metric_extractor_no_live_calls"
        and post_live_scoring_receipt.get("source_contracts", {}).get("post_live_promotion_preflight_audit") == "post_live_promotion_preflight_audit_no_live_or_scoring_calls"
        and post_live_scoring_receipt.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(scoring_receipt_summary.get("receipt_rows", 0)) == 6
        and int(scoring_receipt_summary.get("pass_rows", -1)) == 0
        and int(scoring_receipt_summary.get("blocked_rows", 0)) == 6
        and int(scoring_receipt_summary.get("missing_source_rows", -1)) == 0
        and scoring_receipt_summary.get("scoring_execute_record_exists") is False
        and scoring_receipt_summary.get("scoring_execute_record_path") == "outputs/research_progress_snapshot/post_live_scoring_launcher_execute_latest.json"
        and scoring_receipt_summary.get("latest_scoring_execute_status") == ""
        and scoring_receipt_summary.get("latest_scoring_execute_runtime_contract") == ""
        and scoring_receipt_summary.get("latest_scoring_execute_scope") == ""
        and int(scoring_receipt_summary.get("executed_scoring_rows", -1)) == 0
        and int(scoring_receipt_summary.get("passed_scoring_rows", -1)) == 0
        and int(scoring_receipt_summary.get("failed_scoring_rows", -1)) == 0
        and int(scoring_receipt_summary.get("expected_ready_scoring_rows", -1)) == 0
        and int(scoring_receipt_summary.get("selected_scoring_rows", -1)) == 0
        and int(scoring_receipt_summary.get("available_scoring_rows", 0)) == 6
        and int(scoring_receipt_summary.get("plan_scoring_execution_steps", 0)) == 6
        and int(scoring_receipt_summary.get("readiness_ready_to_score_steps", -1)) == 0
        and int(scoring_receipt_summary.get("scoring_output_promotion_ready_rows", -1)) == 0
        and int(scoring_receipt_summary.get("scoring_output_missing_artifacts", 0)) == 12
        and int(scoring_receipt_summary.get("computed_time_metric_rows", -1)) == 0
        and int(scoring_receipt_summary.get("ready_time_metric_rows", -1)) == 0
        and scoring_receipt_summary.get("ready_for_promotion_review") is False
        and scoring_receipt_summary.get("promotion_preflight_ready") is False
        and int(scoring_receipt_summary.get("traceability_rows", 0)) == 54
        and int(scoring_receipt_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(scoring_receipt_summary.get("live_calls_performed_by_builder", -1)) == 0
        and scoring_receipt_summary.get("no_live_calls_performed_by_auditor") is True
        and scoring_receipt_summary.get("no_scoring_commands_executed_by_auditor") is True
        and scoring_receipt_summary.get("no_secret_values_written") is True
        and scoring_receipt_summary.get("no_new_metric_claim") is True
        and required_scoring_receipt_ids.issubset(scoring_receipt_rows)
        and all(row.get("status") == "blocked" for row in scoring_receipt_rows.values())
        and all(row.get("source_artifacts_exist") is True for row in scoring_receipt_rows.values()),
        json.dumps(
            {
                "contract": post_live_scoring_receipt.get("runtime_contract"),
                "status": post_live_scoring_receipt.get("status"),
                "summary": scoring_receipt_summary,
                "missing_ids": sorted(required_scoring_receipt_ids - set(scoring_receipt_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    scoring_output_summary = post_live_scoring_output_audit.get("summary", {})
    scoring_output_rows = {
        row.get("scoring_execution_id"): row
        for row in post_live_scoring_output_audit.get("rows", [])
    }
    add_check(
        checks,
        "post_live_scoring_output_audit_contract",
        post_live_scoring_output_audit.get("runtime_contract") == "post_live_scoring_output_audit_no_scoring_calls"
        and post_live_scoring_output_audit.get("status") == "blocked_waiting_scoring_outputs"
        and post_live_scoring_output_audit.get("source_contracts", {}).get("post_live_scoring_execution_plan") == "post_live_scoring_execution_plan_no_live_calls"
        and post_live_scoring_output_audit.get("source_contracts", {}).get("post_live_scoring_launcher") == "post_live_scoring_launcher_dry_run_no_scoring_calls"
        and post_live_scoring_output_audit.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and post_live_scoring_output_audit.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_scoring_output_audit.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(scoring_output_summary.get("scoring_output_rows", 0)) == 6
        and int(scoring_output_summary.get("p0_output_rows", 0)) == 3
        and int(scoring_output_summary.get("p1_output_rows", 0)) == 3
        and int(scoring_output_summary.get("total_output_artifacts", 0)) == 16
        and int(scoring_output_summary.get("existing_output_artifacts", -1)) == 4
        and int(scoring_output_summary.get("missing_output_artifacts", -1)) == 12
        and int(scoring_output_summary.get("all_output_artifacts_exist_rows", -1)) == 1
        and int(scoring_output_summary.get("missing_output_rows", -1)) == 5
        and int(scoring_output_summary.get("blocked_current_state_rows", 0)) == 6
        and int(scoring_output_summary.get("promotion_ready_rows", -1)) == 0
        and int(scoring_output_summary.get("ready_to_score_steps", -1)) == 0
        and int(scoring_output_summary.get("ready_to_promote_count", -1)) == 0
        and int(scoring_output_summary.get("scoring_launcher_executed_rows", -1)) == 0
        and scoring_output_summary.get("scoring_execute_record_exists") is False
        and int(scoring_output_summary.get("traceability_rows", 0)) == 54
        and int(scoring_output_summary.get("live_calls_performed_by_builder", -1)) == 0
        and scoring_output_summary.get("no_live_calls_performed") is True
        and scoring_output_summary.get("no_scoring_commands_executed") is True
        and scoring_output_summary.get("no_secret_values_written") is True
        and scoring_output_summary.get("no_new_metric_claim") is True
        and required_scoring_execution_ids.issubset(scoring_output_rows)
        and all(row.get("current_state_blocked") is True for row in scoring_output_rows.values())
        and all(row.get("ready_for_promotion_gate") is False for row in scoring_output_rows.values())
        and scoring_output_rows.get("promotion_refresh_validation", {}).get("all_output_artifacts_exist") is True
        and scoring_output_rows.get("deepseek_resume_safety_score", {}).get("missing_output_artifact_count") == 3,
        json.dumps(
            {
                "contract": post_live_scoring_output_audit.get("runtime_contract"),
                "status": post_live_scoring_output_audit.get("status"),
                "summary": scoring_output_summary,
                "missing_ids": sorted(required_scoring_execution_ids - set(scoring_output_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    evidence_dag_summary = post_live_evidence_dependency_dag.get("summary", {})
    evidence_dag_rows = {
        row.get("node_id"): row
        for row in post_live_evidence_dependency_dag.get("rows", [])
    }
    required_evidence_dag_ids = {
        "live_outputs_complete",
        "output_schema_clean",
        "deepseek_resume_safety_score",
        "deepseek_split20_latency_score",
        "omni48_label_metrics",
        "qwen_backup_metrics",
        "metric_extraction_complete",
        "latency_claim_matrix_update",
        "promotion_gate_pass",
        "report_ppt_refresh_validation",
    }
    add_check(
        checks,
        "post_live_evidence_dependency_dag_contract",
        post_live_evidence_dependency_dag.get("runtime_contract") == "post_live_evidence_dependency_dag_no_live_calls"
        and post_live_evidence_dependency_dag.get("status") == "blocked_waiting_live_outputs"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("post_live_scoring_execution_plan") == "post_live_scoring_execution_plan_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("post_live_scoring_output_audit") == "post_live_scoring_output_audit_no_scoring_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("live_metric_extraction_contract") == "live_metric_extraction_contract_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("live_output_schema_contract") == "live_output_schema_contract_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("post_live_acceptance_scorecard") == "post_live_acceptance_scorecard_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("post_live_latency_claim_matrix") == "post_live_latency_claim_matrix_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_evidence_dependency_dag.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(evidence_dag_summary.get("dag_nodes", 0)) == 10
        and int(evidence_dag_summary.get("p0_dag_nodes", 0)) == 8
        and int(evidence_dag_summary.get("p1_dag_nodes", 0)) == 2
        and int(evidence_dag_summary.get("blocked_nodes", 0)) == 10
        and int(evidence_dag_summary.get("fallback_only_nodes", 0)) == 1
        and int(evidence_dag_summary.get("label_only_nodes", 0)) == 1
        and int(evidence_dag_summary.get("ready_nodes", -1)) == 0
        and int(evidence_dag_summary.get("expected_live_calls", 0)) == 382
        and int(evidence_dag_summary.get("missing_output_surfaces", 0)) == 3
        and int(evidence_dag_summary.get("ready_to_score_steps", -1)) == 0
        and int(evidence_dag_summary.get("scoring_execution_steps", 0)) == 6
        and int(evidence_dag_summary.get("scoring_output_promotion_ready_rows", -1)) == 0
        and int(evidence_dag_summary.get("metric_contract_count", 0)) == 8
        and int(evidence_dag_summary.get("schema_contract_count", 0)) == 8
        and int(evidence_dag_summary.get("scorecard_rows", 0)) == 9
        and int(evidence_dag_summary.get("latency_claim_rows", 0)) == 8
        and int(evidence_dag_summary.get("ready_to_promote_count", -1)) == 0
        and int(evidence_dag_summary.get("traceability_rows", 0)) == 54
        and int(evidence_dag_summary.get("live_calls_performed_by_builder", -1)) == 0
        and evidence_dag_summary.get("no_scoring_commands_executed") is True
        and evidence_dag_summary.get("no_secret_values_written") is True
        and evidence_dag_summary.get("no_new_metric_claim") is True
        and required_evidence_dag_ids.issubset(evidence_dag_rows)
        and evidence_dag_rows.get("live_outputs_complete", {}).get("node_order") == 1
        and evidence_dag_rows.get("output_schema_clean", {}).get("depends_on") == ["live_outputs_complete"]
        and evidence_dag_rows.get("deepseek_resume_safety_score", {}).get("depends_on") == ["output_schema_clean"]
        and evidence_dag_rows.get("deepseek_split20_latency_score", {}).get("depends_on") == ["deepseek_resume_safety_score"]
        and evidence_dag_rows.get("omni48_label_metrics", {}).get("claim_boundary") == "label_only_no_timeline_writeback"
        and evidence_dag_rows.get("qwen_backup_metrics", {}).get("claim_boundary") == "fallback_only_not_primary_claim"
        and evidence_dag_rows.get("metric_extraction_complete", {}).get("depends_on") == [
            "deepseek_resume_safety_score",
            "deepseek_split20_latency_score",
            "omni48_label_metrics",
            "qwen_backup_metrics",
        ]
        and evidence_dag_rows.get("latency_claim_matrix_update", {}).get("depends_on") == ["metric_extraction_complete"]
        and evidence_dag_rows.get("promotion_gate_pass", {}).get("depends_on") == ["latency_claim_matrix_update"]
        and evidence_dag_rows.get("report_ppt_refresh_validation", {}).get("depends_on") == ["promotion_gate_pass"],
        json.dumps(
            {
                "contract": post_live_evidence_dependency_dag.get("runtime_contract"),
                "status": post_live_evidence_dependency_dag.get("status"),
                "summary": evidence_dag_summary,
                "missing_ids": sorted(required_evidence_dag_ids - set(evidence_dag_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_input_summary = live_input_integrity.get("summary", {})
    live_input_rows = {row.get("surface_id"): row for row in live_input_integrity.get("rows", [])}
    required_live_input_ids = {
        "deepseek_resume_inputs",
        "qwen_full_backup_inputs",
        "omni48_label_only_inputs",
    }
    add_check(
        checks,
        "live_input_integrity_audit_contract",
        live_input_integrity.get("runtime_contract") == "live_input_integrity_audit_no_live_calls"
        and live_input_integrity.get("status") == "inputs_ready_waiting_credentials_or_quota"
        and live_input_integrity.get("source_contracts", {}).get("split20_resume_export_audit") == "split20_resume_export_audit_no_live_calls"
        and live_input_integrity.get("source_contracts", {}).get("split20_full_live_manifest") == "split20_full_live_manifest_no_live_calls"
        and live_input_integrity.get("source_contracts", {}).get("omni48_live_call_manifest") == "omni48_live_call_manifest_no_live_calls_label_only"
        and live_input_integrity.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_input_integrity.get("source_contracts", {}).get("live_parallelism_sensitivity") == "live_parallelism_sensitivity_no_live_calls"
        and int(live_input_summary.get("surface_count", 0)) == 3
        and int(live_input_summary.get("input_ready_surfaces", 0)) == 3
        and int(live_input_summary.get("p0_input_ready_surfaces", 0)) == 1
        and int(live_input_summary.get("missing_input_surfaces", -1)) == 0
        and int(live_input_summary.get("output_missing_surfaces", 0)) == 3
        and int(live_input_summary.get("credential_or_quota_blocked_runs", 0)) == 3
        and int(live_input_summary.get("deepseek_resume_prompt_calls", 0)) == 139
        and int(live_input_summary.get("deepseek_resume_parent_windows", 0)) == 101
        and int(live_input_summary.get("deepseek_resume_pending_ids", 0)) == 101
        and int(live_input_summary.get("qwen_full_prompt_calls", 0)) == 147
        and int(live_input_summary.get("qwen_full_parent_windows", 0)) == 104
        and int(live_input_summary.get("omni48_manifest_windows", 0)) == 48
        and int(live_input_summary.get("omni48_planned_calls", 0)) == 96
        and live_input_summary.get("recommended_policy") == "max20"
        and int(live_input_summary.get("recommended_workers", 0)) == 8
        and int(live_input_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_input_summary.get("no_secret_values_written") is True
        and live_input_summary.get("no_new_metric_claim") is True
        and required_live_input_ids.issubset(live_input_rows)
        and live_input_integrity.get("input_ready_surface_ids") == [
            "deepseek_resume_inputs",
            "qwen_full_backup_inputs",
            "omni48_label_only_inputs",
        ]
        and live_input_integrity.get("output_missing_surface_ids") == [
            "deepseek_resume_inputs",
            "qwen_full_backup_inputs",
            "omni48_label_only_inputs",
        ]
        and live_input_rows.get("deepseek_resume_inputs", {}).get("input_ready") is True
        and live_input_rows.get("deepseek_resume_inputs", {}).get("writeback_right") == "block_or_quarantine_only"
        and live_input_rows.get("qwen_full_backup_inputs", {}).get("input_ready") is True
        and live_input_rows.get("omni48_label_only_inputs", {}).get("writeback_right") == "label_only_no_timeline_writeback",
        json.dumps(
            {
                "contract": live_input_integrity.get("runtime_contract"),
                "status": live_input_integrity.get("status"),
                "summary": live_input_summary,
                "missing_ids": sorted(required_live_input_ids - set(live_input_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_runbook_summary = live_execution_runbook.get("summary", {})
    live_runbook_steps = {row.get("step_id"): row for row in live_execution_runbook.get("steps", [])}
    required_runbook_steps = {
        "credential_preflight",
        "deepseek_resume_primary",
        "post_live_output_audit",
        "deepseek_resume_safety_and_comparison",
        "omni48_label_only_live",
        "qwen_full_backup_optional",
        "refresh_report_ppt_validation",
    }
    add_check(
        checks,
        "live_execution_runbook_contract",
        live_execution_runbook.get("runtime_contract") == "live_execution_runbook_no_live_calls_no_secret_values"
        and live_execution_runbook.get("secret_policy") == "env_presence_only_no_secret_values_written"
        and live_execution_runbook.get("status") == "blocked_waiting_for_credentials_or_live_outputs"
        and int(live_runbook_summary.get("runbook_step_count", 0)) == 7
        and int(live_runbook_summary.get("p0_steps", 0)) == 5
        and int(live_runbook_summary.get("p1_steps", 0)) == 2
        and int(live_runbook_summary.get("blocked_steps", 0)) == 5
        and int(live_runbook_summary.get("planned_live_calls_total", 0)) == 382
        and int(live_runbook_summary.get("p0_planned_live_calls", 0)) == 139
        and live_runbook_summary.get("deepseek_primary_policy") == "max20"
        and int(live_runbook_summary.get("claim_now_slo_pass", 0)) == 4
        and int(live_runbook_summary.get("claim_now_slo_rows", 0)) == 4
        and int(live_runbook_summary.get("ready_runs", -1)) == 0
        and int(live_runbook_summary.get("missing_output_surfaces", 0)) == 3
        and int(live_runbook_summary.get("ready_to_score_steps", -1)) == 0
        and int(live_runbook_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_runbook_summary.get("no_secret_values_written") is True
        and required_runbook_steps.issubset(live_runbook_steps)
        and live_runbook_steps.get("deepseek_resume_primary", {}).get("planned_live_calls") == 139
        and "--window-id-file outputs/research_progress_snapshot/split20_deepseek_resume_after_top3_window_ids.txt"
        in live_runbook_steps.get("deepseek_resume_primary", {}).get("command", "")
        and all(
            "--skip-existing-output" in live_runbook_steps.get(step_id, {}).get("command", "")
            for step_id in ["deepseek_resume_primary", "omni48_label_only_live", "qwen_full_backup_optional"]
        )
        and all(
            "--max-call-attempts 2" in live_runbook_steps.get(step_id, {}).get("command", "")
            for step_id in ["deepseek_resume_primary", "omni48_label_only_live", "qwen_full_backup_optional"]
        )
        and all(
            "--retry-backoff-seconds 2.0" in live_runbook_steps.get(step_id, {}).get("command", "")
            for step_id in ["deepseek_resume_primary", "omni48_label_only_live", "qwen_full_backup_optional"]
        )
        and "python scripts/live/build_live_output_audit.py"
        in live_runbook_steps.get("post_live_output_audit", {}).get("command", "")
        and "python scripts/misc/refresh_latest_research_artifacts.py"
        in live_runbook_steps.get("refresh_report_ppt_validation", {}).get("command", ""),
        json.dumps(
            {
                "contract": live_execution_runbook.get("runtime_contract"),
                "status": live_execution_runbook.get("status"),
                "summary": live_runbook_summary,
                "steps": sorted(live_runbook_steps),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_bundle_summary = live_execution_bundle.get("summary", {})
    live_bundle_rows = {
        row.get("bundle_step_id"): row
        for row in live_execution_bundle.get("rows", [])
    }
    required_bundle_ids = {
        "credential_preflight",
        "p0_deepseek_resume_live",
        "p0_output_audit",
        "p0_safety_then_latency_scoring",
        "p0_metrics_promotion_refresh",
        "p1_omni48_label_live",
        "p1_qwen_backup_live",
        "final_report_ppt_validation",
    }
    add_check(
        checks,
        "live_execution_bundle_contract",
        live_execution_bundle.get("runtime_contract") == "live_execution_bundle_no_live_calls_no_secret_values"
        and live_execution_bundle.get("secret_policy") == "env_presence_only_no_secret_values_written"
        and live_execution_bundle.get("status") == "blocked_waiting_credentials_quota_or_live_outputs"
        and live_execution_bundle.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("live_agent_execution_plan") == "live_agent_execution_plan_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("live_execution_runbook") == "live_execution_runbook_no_live_calls_no_secret_values"
        and live_execution_bundle.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("live_execution_handoff_packet") == "live_execution_handoff_packet_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("post_live_evidence_dependency_dag") == "post_live_evidence_dependency_dag_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("live_execution_timing_plan") == "live_execution_timing_plan_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("live_retry_budget_audit") == "live_retry_budget_audit_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("live_token_quota_budget_audit") == "live_token_quota_budget_audit_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("post_live_scoring_execution_plan") == "post_live_scoring_execution_plan_no_live_calls"
        and live_execution_bundle.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(live_bundle_summary.get("bundle_steps", 0)) == 8
        and int(live_bundle_summary.get("p0_bundle_steps", 0)) == 6
        and int(live_bundle_summary.get("p1_bundle_steps", 0)) == 2
        and int(live_bundle_summary.get("blocked_or_waiting_steps", 0)) == 8
        and int(live_bundle_summary.get("live_call_steps", 0)) == 3
        and int(live_bundle_summary.get("command_ready_count", 0)) == 3
        and int(live_bundle_summary.get("command_count", 0)) == 3
        and int(live_bundle_summary.get("planned_live_calls", 0)) == 382
        and int(live_bundle_summary.get("p0_planned_live_calls", 0)) == 139
        and int(live_bundle_summary.get("p1_planned_live_calls", 0)) == 243
        and int(live_bundle_summary.get("ready_runs", -1)) == 0
        and live_bundle_summary.get("credential_ready") is False
        and int(live_bundle_summary.get("known_provider_quota_blockers", 0)) == 1
        and int(live_bundle_summary.get("missing_output_surfaces", 0)) == 3
        and int(live_bundle_summary.get("dag_nodes", 0)) == 10
        and float(live_bundle_summary.get("deepseek_estimated_wall_seconds", 0.0)) == 384.444
        and int(live_bundle_summary.get("max_attempted_requests", 0)) == 764
        and int(live_bundle_summary.get("llm_retry_token_proxy_ceiling", 0)) == 1658856
        and int(live_bundle_summary.get("traceability_rows", 0)) == 54
        and int(live_bundle_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_bundle_summary.get("no_scoring_commands_executed") is True
        and live_bundle_summary.get("no_secret_values_written") is True
        and live_bundle_summary.get("no_new_metric_claim") is True
        and required_bundle_ids.issubset(live_bundle_rows)
        and live_bundle_rows.get("p0_deepseek_resume_live", {}).get("planned_live_calls") == 139
        and "--window-id-file outputs/research_progress_snapshot/split20_deepseek_resume_after_top3_window_ids.txt"
        in live_bundle_rows.get("p0_deepseek_resume_live", {}).get("command", "")
        and live_bundle_rows.get("p0_output_audit", {}).get("depends_on") == ["p0_deepseek_resume_live"]
        and live_bundle_rows.get("p0_safety_then_latency_scoring", {}).get("depends_on") == ["p0_output_audit"]
        and "scripts/analysis/analyze_runtime_safe_llm_guard.py"
        in live_bundle_rows.get("p0_safety_then_latency_scoring", {}).get("command", "")
        and "scripts/analysis/summarize_split_llm_runs.py"
        in live_bundle_rows.get("p0_safety_then_latency_scoring", {}).get("command", "")
        and live_bundle_rows.get("p1_omni48_label_live", {}).get("claim_boundary") == "label_only_no_timeline_writeback"
        and live_bundle_rows.get("p1_qwen_backup_live", {}).get("claim_boundary") == "fallback_only_not_primary_claim"
        and live_bundle_rows.get("final_report_ppt_validation", {}).get("depends_on") == ["p0_metrics_promotion_refresh"],
        json.dumps(
            {
                "contract": live_execution_bundle.get("runtime_contract"),
                "status": live_execution_bundle.get("status"),
                "summary": live_bundle_summary,
                "missing_ids": sorted(required_bundle_ids - set(live_bundle_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_handoff_summary = live_execution_handoff.get("summary", {})
    live_handoff_rows = {row.get("packet_id"): row for row in live_execution_handoff.get("rows", [])}
    required_handoff_ids = {
        "credential_quota_preflight",
        "deepseek_resume_primary_command",
        "deepseek_postrun_scoring_gate",
        "claim_preservation_boundary",
        "omni48_label_only_boundary",
        "qwen_full_backup_boundary",
        "final_refresh_validation_sync",
    }
    add_check(
        checks,
        "live_execution_handoff_packet_contract",
        live_execution_handoff.get("runtime_contract") == "live_execution_handoff_packet_no_live_calls"
        and live_execution_handoff.get("secret_policy") == "env_presence_only_no_secret_values_written"
        and live_execution_handoff.get("status") == "blocked_waiting_credentials_or_quota"
        and live_execution_handoff.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_runtime_environment_audit") == "live_runtime_environment_audit_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_input_integrity_audit") == "live_input_integrity_audit_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_execution_runbook") == "live_execution_runbook_no_live_calls_no_secret_values"
        and live_execution_handoff.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_execution_timing_plan") == "live_execution_timing_plan_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_retry_budget_audit") == "live_retry_budget_audit_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_token_quota_budget_audit") == "live_token_quota_budget_audit_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_failure_recovery_playbook") == "live_failure_recovery_playbook_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and live_execution_handoff.get("source_contracts", {}).get("post_live_acceptance_scorecard") == "post_live_acceptance_scorecard_no_live_calls"
        and int(live_handoff_summary.get("packet_rows", 0)) == 7
        and int(live_handoff_summary.get("p0_packet_rows", 0)) == 5
        and int(live_handoff_summary.get("p1_packet_rows", 0)) == 2
        and int(live_handoff_summary.get("handoff_blocked_rows", 0)) == 6
        and live_handoff_summary.get("credential_ready") is False
        and int(live_handoff_summary.get("known_provider_quota_blockers", 0)) == 1
        and int(live_handoff_summary.get("ready_runs", -1)) == 0
        and int(live_handoff_summary.get("blocked_runs", 0)) == 3
        and int(live_handoff_summary.get("input_ready_surfaces", 0)) == 3
        and int(live_handoff_summary.get("command_ready_count", 0)) == 3
        and int(live_handoff_summary.get("command_count", 0)) == 3
        and int(live_handoff_summary.get("planned_live_calls", 0)) == 382
        and int(live_handoff_summary.get("p0_planned_live_calls", 0)) == 139
        and int(live_handoff_summary.get("runbook_steps", 0)) == 7
        and float(live_handoff_summary.get("deepseek_estimated_wall_seconds", 0.0)) == 384.444
        and int(live_handoff_summary.get("max_attempted_requests", 0)) == 764
        and int(live_handoff_summary.get("llm_retry_token_proxy_ceiling", 0)) == 1658856
        and int(live_handoff_summary.get("ready_recovery_actions", 0)) == 8
        and int(live_handoff_summary.get("missing_output_surfaces", 0)) == 3
        and int(live_handoff_summary.get("ready_to_score_steps", -1)) == 0
        and int(live_handoff_summary.get("ready_to_promote_count", -1)) == 0
        and int(live_handoff_summary.get("scorecard_rows", 0)) == 9
        and int(live_handoff_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_handoff_summary.get("no_secret_values_written") is True
        and live_handoff_summary.get("no_new_metric_claim") is True
        and required_handoff_ids.issubset(live_handoff_rows)
        and "llm_window_batch_policy_eval.py" in live_handoff_rows.get("deepseek_resume_primary_command", {}).get("required_action", "")
        and live_handoff_rows.get("claim_preservation_boundary", {}).get("claim_boundary") == "no_new_metric_claim"
        and live_handoff_rows.get("omni48_label_only_boundary", {}).get("claim_boundary") == "label_only_no_timeline_writeback"
        and live_handoff_rows.get("qwen_full_backup_boundary", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim",
        json.dumps(
            {
                "contract": live_execution_handoff.get("runtime_contract"),
                "status": live_execution_handoff.get("status"),
                "summary": live_handoff_summary,
                "missing_ids": sorted(required_handoff_ids - set(live_handoff_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_command_summary = live_command_surface_audit.get("summary", {})
    live_command_rows = {row.get("command_id"): row for row in live_command_surface_audit.get("rows", [])}
    required_live_command_ids = {
        "deepseek_resume_primary",
        "omni48_label_only_live",
        "qwen_full_backup_optional",
    }
    add_check(
        checks,
        "live_command_surface_audit_contract",
        live_command_surface_audit.get("runtime_contract") == "live_command_surface_audit_no_live_calls"
        and live_command_surface_audit.get("secret_policy") == "commands_scanned_no_secret_values_written"
        and live_command_surface_audit.get("status") == "commands_ready_waiting_credentials_or_quota"
        and live_command_surface_audit.get("source_contracts", {}).get("live_execution_runbook") == "live_execution_runbook_no_live_calls_no_secret_values"
        and live_command_surface_audit.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_command_surface_audit.get("source_contracts", {}).get("live_input_integrity_audit") == "live_input_integrity_audit_no_live_calls"
        and live_command_surface_audit.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and int(live_command_summary.get("command_count", 0)) == 3
        and int(live_command_summary.get("command_ready_count", 0)) == 3
        and int(live_command_summary.get("p0_command_ready_count", 0)) == 1
        and int(live_command_summary.get("missing_input_commands", -1)) == 0
        and int(live_command_summary.get("duplicate_output_paths", -1)) == 0
        and int(live_command_summary.get("secret_literal_commands", -1)) == 0
        and int(live_command_summary.get("missing_required_flag_commands", -1)) == 0
        and int(live_command_summary.get("skip_existing_output_commands", 0)) == 3
        and int(live_command_summary.get("bounded_retry_commands", 0)) == 3
        and int(live_command_summary.get("planned_live_calls", 0)) == 382
        and int(live_command_summary.get("p0_planned_live_calls", 0)) == 139
        and int(live_command_summary.get("deepseek_resume_calls", 0)) == 139
        and int(live_command_summary.get("qwen_backup_calls", 0)) == 147
        and int(live_command_summary.get("omni48_calls", 0)) == 96
        and int(live_command_summary.get("readiness_ready_runs", -1)) == 0
        and int(live_command_summary.get("input_ready_surfaces", 0)) == 3
        and int(live_command_summary.get("missing_output_surfaces", 0)) == 3
        and int(live_command_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_command_summary.get("no_secret_values_written") is True
        and live_command_summary.get("no_new_metric_claim") is True
        and required_live_command_ids.issubset(live_command_rows)
        and all(row.get("command_ready") is True for row in live_command_rows.values())
        and all(row.get("skip_existing_output") is True for row in live_command_rows.values())
        and all(str(row.get("max_call_attempts", "")) == "2" for row in live_command_rows.values())
        and all(str(row.get("retry_backoff_seconds", "")) == "2.0" for row in live_command_rows.values())
        and live_command_rows.get("deepseek_resume_primary", {}).get("planned_live_calls") == 139
        and live_command_rows.get("deepseek_resume_primary", {}).get("output_jsonl")
        == "outputs/runtime_safe_llm_window_batch/deepseek_proxy_high_risk_104w_split20_resume_after_top3.jsonl"
        and live_command_rows.get("omni48_label_only_live", {}).get("writeback_right") == "label_only_no_timeline_writeback"
        and live_command_rows.get("qwen_full_backup_optional", {}).get("planned_live_calls") == 147,
        json.dumps(
            {
                "contract": live_command_surface_audit.get("runtime_contract"),
                "status": live_command_surface_audit.get("status"),
                "summary": live_command_summary,
                "missing_ids": sorted(required_live_command_ids - set(live_command_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    eligibility_summary = live_execution_eligibility.get("summary", {})
    eligibility_rows = {
        row.get("gate_id"): row
        for row in live_execution_eligibility.get("rows", [])
    }
    required_eligibility_ids = {
        "live_input_surface_gate",
        "live_command_surface_gate",
        "live_runtime_credential_gate",
        "live_readiness_gate",
        "live_launcher_execute_gate",
        "operator_handoff_gate",
        "post_live_promotion_preflight_gate",
    }
    add_check(
        checks,
        "live_execution_eligibility_gate_contract",
        live_execution_eligibility.get("runtime_contract") == "live_execution_eligibility_gate_no_live_calls_no_secret_values"
        and live_execution_eligibility.get("status") == "blocked_waiting_credentials_quota_or_execute_flag"
        and live_execution_eligibility.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("live_input_integrity_audit") == "live_input_integrity_audit_no_live_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("live_runtime_environment_audit") == "live_runtime_environment_audit_no_live_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("live_execution_handoff_packet") == "live_execution_handoff_packet_no_live_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("live_execution_launcher") == "live_execution_launcher_dry_run_no_live_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("post_live_promotion_preflight_audit") == "post_live_promotion_preflight_audit_no_live_or_scoring_calls"
        and live_execution_eligibility.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(eligibility_summary.get("eligibility_rows", 0)) == 7
        and int(eligibility_summary.get("pass_rows", 0)) == 2
        and int(eligibility_summary.get("blocked_rows", 0)) == 5
        and int(eligibility_summary.get("missing_source_rows", -1)) == 0
        and eligibility_summary.get("ready_to_execute_live") is False
        and int(eligibility_summary.get("input_ready_surfaces", 0)) == 3
        and int(eligibility_summary.get("command_ready_count", 0)) == 3
        and eligibility_summary.get("credential_ready") is False
        and int(eligibility_summary.get("known_provider_quota_blockers", 0)) == 1
        and int(eligibility_summary.get("ready_runs", -1)) == 0
        and int(eligibility_summary.get("selected_live_calls", 0)) == 139
        and int(eligibility_summary.get("p0_selected_live_calls", 0)) == 139
        and eligibility_summary.get("execute_live") is False
        and eligibility_summary.get("execution_allowed") is False
        and int(eligibility_summary.get("handoff_blocked_rows", 0)) == 6
        and eligibility_summary.get("promotion_preflight_ready") is False
        and int(eligibility_summary.get("traceability_rows", 0)) == 54
        and int(eligibility_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(eligibility_summary.get("live_calls_performed_by_builder", -1)) == 0
        and eligibility_summary.get("no_live_calls_performed") is True
        and eligibility_summary.get("no_secret_values_written") is True
        and eligibility_summary.get("no_new_metric_claim") is True
        and live_execution_eligibility.get("recommended_first_execute_command") == "python scripts/live/run_live_execution_sequence.py --execute-live --live-scope p0"
        and required_eligibility_ids.issubset(eligibility_rows)
        and eligibility_rows.get("live_input_surface_gate", {}).get("status") == "pass"
        and eligibility_rows.get("live_command_surface_gate", {}).get("status") == "pass"
        and eligibility_rows.get("live_runtime_credential_gate", {}).get("status") == "blocked"
        and eligibility_rows.get("live_readiness_gate", {}).get("status") == "blocked"
        and eligibility_rows.get("live_launcher_execute_gate", {}).get("status") == "blocked"
        and eligibility_rows.get("operator_handoff_gate", {}).get("status") == "blocked"
        and eligibility_rows.get("post_live_promotion_preflight_gate", {}).get("status") == "blocked"
        and all(row.get("source_artifacts_exist") is True for row in eligibility_rows.values()),
        json.dumps(
            {
                "contract": live_execution_eligibility.get("runtime_contract"),
                "status": live_execution_eligibility.get("status"),
                "summary": eligibility_summary,
                "missing_ids": sorted(required_eligibility_ids - set(eligibility_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    receipt_summary = live_execution_receipt.get("summary", {})
    receipt_rows = {
        row.get("receipt_id"): row
        for row in live_execution_receipt.get("rows", [])
    }
    required_receipt_ids = {
        "execute_record_presence",
        "execute_scope_alignment",
        "live_command_result_receipt",
        "postrun_refresh_receipt",
        "output_audit_after_execute",
        "promotion_preflight_after_execute",
    }
    add_check(
        checks,
        "live_execution_receipt_audit_contract",
        live_execution_receipt.get("runtime_contract") == "live_execution_receipt_audit_no_live_calls_no_secret_values"
        and live_execution_receipt.get("status") == "blocked_no_execute_receipt_or_outputs"
        and live_execution_receipt.get("source_contracts", {}).get("live_execution_launcher") == "live_execution_launcher_dry_run_no_live_calls"
        and live_execution_receipt.get("source_contracts", {}).get("live_execution_eligibility_gate") == "live_execution_eligibility_gate_no_live_calls_no_secret_values"
        and live_execution_receipt.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_execution_receipt.get("source_contracts", {}).get("post_live_promotion_preflight_audit") == "post_live_promotion_preflight_audit_no_live_or_scoring_calls"
        and live_execution_receipt.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(receipt_summary.get("receipt_rows", 0)) == 6
        and int(receipt_summary.get("pass_rows", -1)) == 0
        and int(receipt_summary.get("blocked_rows", 0)) == 6
        and int(receipt_summary.get("missing_source_rows", -1)) == 0
        and receipt_summary.get("execute_record_exists") is False
        and receipt_summary.get("execute_record_path") == "outputs/research_progress_snapshot/live_execution_launcher_execute_latest.json"
        and receipt_summary.get("latest_execute_status") == ""
        and receipt_summary.get("latest_execute_runtime_contract") == ""
        and receipt_summary.get("latest_execute_live_scope") == ""
        and int(receipt_summary.get("started_live_command_calls", -1)) == 0
        and int(receipt_summary.get("passed_live_command_calls", -1)) == 0
        and int(receipt_summary.get("failed_live_command_calls", -1)) == 0
        and int(receipt_summary.get("failed_live_command_rows", -1)) == 0
        and receipt_summary.get("postrun_refresh_executed") is False
        and int(receipt_summary.get("live_calls_performed_by_launcher", -1)) == 0
        and int(receipt_summary.get("expected_selected_live_calls", 0)) == 139
        and receipt_summary.get("eligibility_ready_to_execute_live") is False
        and int(receipt_summary.get("observed_live_output_rows", -1)) == 0
        and int(receipt_summary.get("claim_ready_surfaces", -1)) == 0
        and int(receipt_summary.get("missing_output_surfaces", 0)) == 3
        and receipt_summary.get("ready_for_postrun_scoring_review") is False
        and int(receipt_summary.get("traceability_rows", 0)) == 54
        and int(receipt_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(receipt_summary.get("live_calls_performed_by_builder", -1)) == 0
        and receipt_summary.get("no_live_calls_performed_by_auditor") is True
        and receipt_summary.get("no_secret_values_written") is True
        and receipt_summary.get("no_new_metric_claim") is True
        and required_receipt_ids.issubset(receipt_rows)
        and all(row.get("status") == "blocked" for row in receipt_rows.values())
        and all(row.get("source_artifacts_exist") is True for row in receipt_rows.values()),
        json.dumps(
            {
                "contract": live_execution_receipt.get("runtime_contract"),
                "status": live_execution_receipt.get("status"),
                "summary": receipt_summary,
                "missing_ids": sorted(required_receipt_ids - set(receipt_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    launcher_summary = live_execution_launcher.get("summary", {})
    launcher_rows = {
        row.get("launcher_step_id"): row
        for row in live_execution_launcher.get("rows", [])
    }
    required_launcher_ids = {
        "credential_preflight_refresh",
        "live_deepseek_resume_primary",
        "live_omni48_label_only_live",
        "live_qwen_full_backup_optional",
        "postrun_refresh_validation",
    }
    add_check(
        checks,
        "live_execution_launcher_contract",
        live_execution_launcher.get("runtime_contract") == "live_execution_launcher_dry_run_no_live_calls"
        and live_execution_launcher.get("secret_policy") == "env_presence_only_no_secret_values_written"
        and live_execution_launcher.get("status") == "dry_run_blocked_waiting_credentials_or_execute_flag"
        and live_execution_launcher.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_execution_launcher.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_execution_launcher.get("source_contracts", {}).get("live_execution_bundle") == "live_execution_bundle_no_live_calls_no_secret_values"
        and live_execution_launcher.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(launcher_summary.get("launcher_steps", 0)) == 5
        and int(launcher_summary.get("available_live_command_rows", 0)) == 3
        and int(launcher_summary.get("selected_live_command_rows", 0)) == 1
        and int(launcher_summary.get("available_live_calls", 0)) == 382
        and int(launcher_summary.get("selected_live_calls", 0)) == 139
        and int(launcher_summary.get("p0_selected_live_calls", 0)) == 139
        and int(launcher_summary.get("p1_selected_live_calls", -1)) == 0
        and launcher_summary.get("live_scope") == "p0"
        and launcher_summary.get("execute_live") is False
        and launcher_summary.get("credential_ready") is False
        and launcher_summary.get("artifact_credential_ready") is False
        and int(launcher_summary.get("known_provider_quota_blockers", 0)) == 1
        and launcher_summary.get("execution_allowed") is False
        and int(launcher_summary.get("execution_blocker_count", 0)) == 1
        and int(launcher_summary.get("failed_execution_rows", -1)) == 0
        and int(launcher_summary.get("executed_live_command_rows", -1)) == 0
        and int(launcher_summary.get("failed_live_command_rows", -1)) == 0
        and int(launcher_summary.get("started_live_command_calls", -1)) == 0
        and int(launcher_summary.get("passed_live_command_calls", -1)) == 0
        and int(launcher_summary.get("failed_live_command_calls", -1)) == 0
        and launcher_summary.get("postrun_refresh_executed") is False
        and launcher_summary.get("postrun_refresh_blocked_by_live_failures") is False
        and int(launcher_summary.get("traceability_rows", 0)) == 54
        and int(launcher_summary.get("live_calls_performed_by_launcher", -1)) == 0
        and launcher_summary.get("execute_record_exists") is False
        and launcher_summary.get("execute_record_path") == "outputs/research_progress_snapshot/live_execution_launcher_execute_latest.json"
        and launcher_summary.get("latest_execute_status") == ""
        and launcher_summary.get("latest_execute_runtime_contract") == ""
        and launcher_summary.get("latest_execute_live_scope") == ""
        and int(launcher_summary.get("latest_execute_started_live_command_calls", -1)) == 0
        and int(launcher_summary.get("latest_execute_passed_live_command_calls", -1)) == 0
        and int(launcher_summary.get("latest_execute_failed_live_command_calls", -1)) == 0
        and int(launcher_summary.get("latest_execute_live_calls_performed_by_launcher", -1)) == 0
        and int(launcher_summary.get("latest_execute_failed_live_command_rows", -1)) == 0
        and launcher_summary.get("latest_execute_postrun_refresh_executed") is False
        and launcher_summary.get("dry_run_no_live_calls") is True
        and launcher_summary.get("no_secret_values_written") is True
        and launcher_summary.get("no_new_metric_claim") is True
        and live_execution_launcher.get("execution_blockers") == ["dry_run_only_requires_execute_live_flag"]
        and required_launcher_ids.issubset(launcher_rows)
        and launcher_rows.get("live_deepseek_resume_primary", {}).get("selected") is True
        and launcher_rows.get("live_deepseek_resume_primary", {}).get("planned_live_calls") == 139
        and launcher_rows.get("live_deepseek_resume_primary", {}).get("status") == "dry_run_selected_not_executed"
        and launcher_rows.get("live_omni48_label_only_live", {}).get("selected") is False
        and launcher_rows.get("live_qwen_full_backup_optional", {}).get("selected") is False
        and launcher_rows.get("postrun_refresh_validation", {}).get("status") == "dry_run_postrun_not_executed",
        json.dumps(
            {
                "contract": live_execution_launcher.get("runtime_contract"),
                "status": live_execution_launcher.get("status"),
                "summary": launcher_summary,
                "missing_ids": sorted(required_launcher_ids - set(launcher_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    repair_summary = live_output_repair_plan.get("summary", {})
    repair_rows = {row.get("surface_id"): row for row in live_output_repair_plan.get("rows", [])}
    required_repair_ids = {
        "deepseek_resume_after_top3",
        "qwen_full_backup",
        "omni48_label_only",
    }
    add_check(
        checks,
        "live_output_repair_plan_contract",
        live_output_repair_plan.get("runtime_contract") == "live_output_repair_plan_no_live_calls_no_scoring"
        and live_output_repair_plan.get("status") == "waiting_live_outputs"
        and live_output_repair_plan.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_output_repair_plan.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_output_repair_plan.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and int(repair_summary.get("repair_rows", 0)) == 3
        and int(repair_summary.get("clean_run_rows", 0)) == 3
        and int(repair_summary.get("skip_existing_rerun_rows", -1)) == 0
        and int(repair_summary.get("quarantine_required_rows", -1)) == 0
        and int(repair_summary.get("scoring_ready_rows", -1)) == 0
        and int(repair_summary.get("expected_calls", 0)) == 382
        and int(repair_summary.get("observed_calls", -1)) == 0
        and int(repair_summary.get("successful_calls", -1)) == 0
        and int(repair_summary.get("failed_calls", -1)) == 0
        and int(repair_summary.get("missing_calls", 0)) == 382
        and int(repair_summary.get("skip_existing_supported_rows", 0)) == 3
        and int(repair_summary.get("bounded_retry_supported_rows", 0)) == 3
        and int(repair_summary.get("live_calls_performed_by_builder", -1)) == 0
        and repair_summary.get("no_scoring_commands_executed") is True
        and repair_summary.get("no_secret_values_written") is True
        and repair_summary.get("no_new_metric_claim") is True
        and required_repair_ids.issubset(repair_rows)
        and repair_rows.get("deepseek_resume_after_top3", {}).get("repair_action") == "clean_run_waiting_credentials_or_quota"
        and repair_rows.get("deepseek_resume_after_top3", {}).get("priority") == "P0"
        and int(repair_rows.get("deepseek_resume_after_top3", {}).get("expected_calls", 0)) == 139
        and repair_rows.get("deepseek_resume_after_top3", {}).get("skip_existing_supported") is True
        and repair_rows.get("deepseek_resume_after_top3", {}).get("bounded_retry_supported") is True
        and "llm_window_batch_policy_eval.py" in repair_rows.get("deepseek_resume_after_top3", {}).get("next_command", "")
        and repair_rows.get("omni48_label_only", {}).get("repair_action") == "clean_run_waiting_credentials_or_quota"
        and int(repair_rows.get("omni48_label_only", {}).get("expected_calls", 0)) == 96
        and repair_rows.get("qwen_full_backup", {}).get("repair_action") == "clean_run_waiting_credentials_or_quota"
        and int(repair_rows.get("qwen_full_backup", {}).get("expected_calls", 0)) == 147,
        json.dumps(
            {
                "contract": live_output_repair_plan.get("runtime_contract"),
                "status": live_output_repair_plan.get("status"),
                "summary": repair_summary,
                "missing_ids": sorted(required_repair_ids - set(repair_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_resume_summary = live_resume_state_audit.get("summary", {})
    live_resume_rows = {row.get("surface_id"): row for row in live_resume_state_audit.get("rows", [])}
    required_live_resume_ids = {
        "deepseek_resume_after_top3",
        "qwen_full_backup",
        "omni48_label_only",
    }
    add_check(
        checks,
        "live_resume_state_audit_contract",
        live_resume_state_audit.get("runtime_contract") == "live_resume_state_audit_no_live_calls"
        and live_resume_state_audit.get("status") == "clean_run_ready_waiting_credentials_or_quota"
        and live_resume_state_audit.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_resume_state_audit.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_resume_state_audit.get("source_contracts", {}).get("live_execution_runbook") == "live_execution_runbook_no_live_calls_no_secret_values"
        and live_resume_state_audit.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and live_resume_state_audit.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and int(live_resume_summary.get("surface_count", 0)) == 3
        and int(live_resume_summary.get("clean_run_surfaces", 0)) == 3
        and int(live_resume_summary.get("partial_or_invalid_surfaces", -1)) == 0
        and int(live_resume_summary.get("completed_output_surfaces", -1)) == 0
        and int(live_resume_summary.get("current_commands_safe_to_run", 0)) == 3
        and int(live_resume_summary.get("p0_current_commands_safe_to_run", 0)) == 1
        and int(live_resume_summary.get("append_resume_supported_surfaces", -1)) == 0
        and int(live_resume_summary.get("skip_existing_supported_surfaces", -1)) == 3
        and int(live_resume_summary.get("bounded_retry_supported_surfaces", -1)) == 3
        and int(live_resume_summary.get("quarantine_required_surfaces", -1)) == 0
        and int(live_resume_summary.get("planned_live_calls", 0)) == 382
        and int(live_resume_summary.get("p0_planned_live_calls", 0)) == 139
        and int(live_resume_summary.get("observed_live_output_rows", -1)) == 0
        and int(live_resume_summary.get("successful_live_output_rows", -1)) == 0
        and int(live_resume_summary.get("missing_live_calls", 0)) == 382
        and int(live_resume_summary.get("ready_to_score_steps", -1)) == 0
        and int(live_resume_summary.get("ready_to_promote_gates", -1)) == 0
        and int(live_resume_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_resume_summary.get("no_secret_values_written") is True
        and live_resume_summary.get("no_new_metric_claim") is True
        and required_live_resume_ids.issubset(live_resume_rows)
        and all(row.get("output_state") == "missing_output_clean_run" for row in live_resume_rows.values())
        and all(row.get("current_command_safe_to_run") is True for row in live_resume_rows.values())
        and all(row.get("append_resume_supported") is False for row in live_resume_rows.values())
        and all(row.get("skip_existing_supported") is True for row in live_resume_rows.values())
        and all(row.get("bounded_retry_supported") is True for row in live_resume_rows.values())
        and all(str(row.get("max_call_attempts", "")) == "2" for row in live_resume_rows.values())
        and all(str(row.get("retry_backoff_seconds", "")) == "2.0" for row in live_resume_rows.values())
        and all(
            row.get("runner_write_mode") == "merge_successful_existing_then_overwrite_output_jsonl_and_csv"
            for row in live_resume_rows.values()
        )
        and live_resume_rows.get("deepseek_resume_after_top3", {}).get("planned_calls") == 139
        and live_resume_rows.get("omni48_label_only", {}).get("recommended_action") == "run_current_command_when_credentials_and_quota_are_ready",
        json.dumps(
            {
                "contract": live_resume_state_audit.get("runtime_contract"),
                "status": live_resume_state_audit.get("status"),
                "summary": live_resume_summary,
                "missing_ids": sorted(required_live_resume_ids - set(live_resume_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    live_env_summary = live_runtime_environment_audit.get("summary", {})
    live_env_checks = {row.get("check_id"): row for row in live_runtime_environment_audit.get("checks", [])}
    required_live_env_checks = {
        "python_version",
        "module:openai",
        "module:numpy",
        "module:soundfile",
        "module:alimeeting_diarization_bench.config",
        "module:scripts.llm_policy_agent_eval",
        "module:scripts.omni_audio_guard_smoke",
        "script:scripts/llm/llm_window_batch_policy_eval.py",
        "script:scripts/llm/omni_guard_window_batch.py",
        "script:scripts/misc/refresh_latest_research_artifacts.py",
        "output_dir:outputs/omni_guard",
        "output_dir:outputs/research_progress_snapshot",
        "output_dir:outputs/runtime_safe_llm_window_batch",
        "omni48_audio_manifest",
    }
    add_check(
        checks,
        "live_runtime_environment_audit_contract",
        live_runtime_environment_audit.get("runtime_contract") == "live_runtime_environment_audit_no_live_calls"
        and live_runtime_environment_audit.get("secret_policy") == "env_presence_only_no_secret_values_written"
        and live_runtime_environment_audit.get("status") == "runtime_ready_waiting_credentials_or_quota"
        and live_runtime_environment_audit.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_runtime_environment_audit.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_runtime_environment_audit.get("source_contracts", {}).get("live_resume_state_audit") == "live_resume_state_audit_no_live_calls"
        and live_runtime_environment_audit.get("source_contracts", {}).get("live_input_integrity_audit") == "live_input_integrity_audit_no_live_calls"
        and live_runtime_environment_audit.get("environment", {}).get("config_defaults_not_counted_as_credentials") is True
        and live_runtime_environment_audit.get("environment", {}).get("dashscope_like_api_key_present") is False
        and int(live_env_summary.get("check_count", 0)) == 14
        and int(live_env_summary.get("passed_checks", 0)) == 14
        and int(live_env_summary.get("failed_checks", -1)) == 0
        and int(live_env_summary.get("module_checks", 0)) == 6
        and int(live_env_summary.get("module_passed", 0)) == 6
        and int(live_env_summary.get("script_checks", 0)) == 3
        and int(live_env_summary.get("script_passed", 0)) == 3
        and int(live_env_summary.get("output_dir_checks", 0)) == 3
        and int(live_env_summary.get("output_dir_passed", 0)) == 3
        and live_env_summary.get("audio_manifest_passed") is True
        and live_env_summary.get("credential_ready") is False
        and int(live_env_summary.get("known_provider_quota_blockers", 0)) >= 1
        and int(live_env_summary.get("command_ready_count", 0)) == 3
        and int(live_env_summary.get("resume_clean_run_surfaces", 0)) == 3
        and int(live_env_summary.get("input_ready_surfaces", 0)) == 3
        and int(live_env_summary.get("live_calls_performed_by_builder", -1)) == 0
        and live_env_summary.get("no_secret_values_written") is True
        and live_env_summary.get("no_new_metric_claim") is True
        and required_live_env_checks.issubset(live_env_checks)
        and all(row.get("status") == "pass" for row in live_env_checks.values()),
        json.dumps(
            {
                "contract": live_runtime_environment_audit.get("runtime_contract"),
                "status": live_runtime_environment_audit.get("status"),
                "summary": live_env_summary,
                "missing_ids": sorted(required_live_env_checks - set(live_env_checks)),
                "blockers": live_runtime_environment_audit.get("blockers", []),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    timing_summary = live_execution_timing.get("summary", {})
    timing_rows = {row.get("timing_id"): row for row in live_execution_timing.get("rows", [])}
    required_timing_ids = {
        "credential_preflight",
        "deepseek_max20_resume_live_call",
        "deepseek_postrun_audit_and_scoring",
        "omni48_label_only_live",
        "qwen_full_backup_optional",
        "refresh_report_ppt_validation",
    }
    add_check(
        checks,
        "live_execution_timing_plan_contract",
        live_execution_timing.get("runtime_contract") == "live_execution_timing_plan_no_live_calls"
        and live_execution_timing.get("status") == "blocked_waiting_for_credentials_or_live_outputs"
        and live_execution_timing.get("source_contracts", {}).get("live_agent_execution_plan") == "live_agent_execution_plan_no_live_calls"
        and live_execution_timing.get("source_contracts", {}).get("live_execution_runbook") == "live_execution_runbook_no_live_calls_no_secret_values"
        and live_execution_timing.get("source_contracts", {}).get("latency_risk_mitigation_plan") == "latency_risk_mitigation_plan_no_live_calls"
        and int(timing_summary.get("timing_rows", 0)) == 6
        and int(timing_summary.get("p0_rows", 0)) == 4
        and int(timing_summary.get("p1_rows", 0)) == 2
        and int(timing_summary.get("known_wall_rows", 0)) == 3
        and int(timing_summary.get("unknown_wall_rows", 0)) == 3
        and int(timing_summary.get("blocked_rows", 0)) == 5
        and int(timing_summary.get("deepseek_resume_calls", 0)) == 139
        and int(timing_summary.get("deepseek_parallel_workers", 0)) == 8
        and int(timing_summary.get("deepseek_parallel_waves", 0)) == 18
        and float(timing_summary.get("deepseek_p95_call_seconds", 0.0)) == 21.358
        and float(timing_summary.get("deepseek_estimated_wall_seconds", 0.0)) == 384.444
        and int(timing_summary.get("qwen_backup_calls", 0)) == 147
        and int(timing_summary.get("qwen_parallel_workers", 0)) == 8
        and int(timing_summary.get("qwen_parallel_waves", 0)) == 19
        and float(timing_summary.get("qwen_estimated_wall_seconds", 0.0)) == 836.456
        and int(timing_summary.get("omni48_label_only_calls", 0)) == 96
        and float(timing_summary.get("omni48_clip_model_seconds_proxy", 0.0)) == 768.0
        and timing_summary.get("primary_policy") == "max20"
        and timing_summary.get("stretch_policy") == "max15"
        and int(timing_summary.get("live_calls_performed_by_builder", -1)) == 0
        and timing_summary.get("no_secret_values_written") is True
        and timing_summary.get("no_new_metric_claim") is True
        and required_timing_ids.issubset(timing_rows)
        and set(live_execution_timing.get("blocked_timing_ids", []))
        == {
            "credential_preflight",
            "deepseek_max20_resume_live_call",
            "deepseek_postrun_audit_and_scoring",
            "omni48_label_only_live",
            "qwen_full_backup_optional",
        }
        and timing_rows.get("deepseek_max20_resume_live_call", {}).get("claim_status") == "not_claimable_until_resume_output_and_scoring"
        and timing_rows.get("qwen_full_backup_optional", {}).get("claim_status") == "fallback_only_not_primary_latency_claim"
        and timing_rows.get("omni48_label_only_live", {}).get("claim_status") == "label_only_not_timeline_writeback_not_guard_latency_claim",
        json.dumps(
            {
                "contract": live_execution_timing.get("runtime_contract"),
                "status": live_execution_timing.get("status"),
                "summary": timing_summary,
                "missing_ids": sorted(required_timing_ids - set(timing_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    retry_budget_summary = live_retry_budget.get("summary", {})
    retry_budget_rows = {row.get("command_id"): row for row in live_retry_budget.get("rows", [])}
    required_retry_budget_ids = {
        "deepseek_resume_primary",
        "omni48_label_only_live",
        "qwen_full_backup_optional",
    }
    add_check(
        checks,
        "live_retry_budget_audit_contract",
        live_retry_budget.get("runtime_contract") == "live_retry_budget_audit_no_live_calls"
        and live_retry_budget.get("status") == "retry_budget_ready_waiting_credentials_or_quota"
        and live_retry_budget.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_retry_budget.get("source_contracts", {}).get("live_execution_timing_plan") == "live_execution_timing_plan_no_live_calls"
        and live_retry_budget.get("source_contracts", {}).get("live_resume_state_audit") == "live_resume_state_audit_no_live_calls"
        and int(retry_budget_summary.get("surface_count", 0)) == 3
        and int(retry_budget_summary.get("bounded_retry_surfaces", 0)) == 3
        and int(retry_budget_summary.get("planned_live_calls", 0)) == 382
        and int(retry_budget_summary.get("max_attempted_requests", 0)) == 764
        and int(retry_budget_summary.get("additional_retry_attempt_budget", 0)) == 382
        and float(retry_budget_summary.get("backoff_ceiling_seconds", 0.0)) == 764.0
        and int(retry_budget_summary.get("p0_planned_live_calls", 0)) == 139
        and int(retry_budget_summary.get("p0_max_attempted_requests", 0)) == 278
        and float(retry_budget_summary.get("p0_backoff_ceiling_seconds", 0.0)) == 278.0
        and float(retry_budget_summary.get("deepseek_retry_ceiling_wall_seconds", 0.0)) == 804.888
        and float(retry_budget_summary.get("deepseek_retry_wall_overhead_seconds", 0.0)) == 420.444
        and float(retry_budget_summary.get("qwen_retry_ceiling_wall_seconds", 0.0)) == 1710.912
        and int(retry_budget_summary.get("omni_retry_attempt_budget", 0)) == 192
        and retry_budget_summary.get("credential_ready") is False
        and int(retry_budget_summary.get("live_calls_performed_by_builder", -1)) == 0
        and retry_budget_summary.get("no_secret_values_written") is True
        and retry_budget_summary.get("no_new_metric_claim") is True
        and required_retry_budget_ids.issubset(retry_budget_rows)
        and all(row.get("max_call_attempts") == 2 for row in retry_budget_rows.values())
        and all(float(row.get("retry_backoff_seconds", 0.0)) == 2.0 for row in retry_budget_rows.values())
        and retry_budget_rows.get("deepseek_resume_primary", {}).get("max_attempted_requests") == 278
        and retry_budget_rows.get("deepseek_resume_primary", {}).get("claim_status") == "retry_budget_planning_only_no_live_metric_claim"
        and retry_budget_rows.get("omni48_label_only_live", {}).get("retry_ceiling_wall_seconds") == "",
        json.dumps(
            {
                "contract": live_retry_budget.get("runtime_contract"),
                "status": live_retry_budget.get("status"),
                "summary": retry_budget_summary,
                "missing_ids": sorted(required_retry_budget_ids - set(retry_budget_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    token_quota_summary = live_token_quota_budget.get("summary", {})
    token_quota_rows = {row.get("surface_id"): row for row in live_token_quota_budget.get("rows", [])}
    required_token_quota_ids = {
        "deepseek_resume_after_top3",
        "qwen_full_backup",
        "omni48_label_only",
    }
    add_check(
        checks,
        "live_token_quota_budget_audit_contract",
        live_token_quota_budget.get("runtime_contract") == "live_token_quota_budget_audit_no_live_calls"
        and live_token_quota_budget.get("status") == "quota_proxy_ready_waiting_credentials_or_quota"
        and live_token_quota_budget.get("token_proxy_policy") == "prompt_chars_div4_proxy_not_provider_billing_tokens"
        and live_token_quota_budget.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_token_quota_budget.get("source_contracts", {}).get("live_retry_budget_audit") == "live_retry_budget_audit_no_live_calls"
        and live_token_quota_budget.get("source_contracts", {}).get("live_input_integrity_audit") == "live_input_integrity_audit_no_live_calls"
        and int(token_quota_summary.get("surface_count", 0)) == 3
        and int(token_quota_summary.get("llm_prompt_surfaces", 0)) == 2
        and int(token_quota_summary.get("omni_surfaces", 0)) == 1
        and int(token_quota_summary.get("llm_prompt_calls", 0)) == 286
        and int(token_quota_summary.get("llm_prompt_chars", 0)) == 3317714
        and int(token_quota_summary.get("llm_token_proxy_chars_div4", 0)) == 829428
        and int(token_quota_summary.get("llm_retry_token_proxy_ceiling", 0)) == 1658856
        and int(token_quota_summary.get("p0_token_proxy_chars_div4", 0)) == 400862
        and int(token_quota_summary.get("p0_retry_token_proxy_ceiling", 0)) == 801724
        and int(token_quota_summary.get("deepseek_resume_token_proxy_chars_div4", 0)) == 400862
        and int(token_quota_summary.get("deepseek_resume_retry_token_proxy_ceiling", 0)) == 801724
        and int(token_quota_summary.get("qwen_full_token_proxy_chars_div4", 0)) == 428566
        and int(token_quota_summary.get("qwen_full_retry_token_proxy_ceiling", 0)) == 857132
        and float(token_quota_summary.get("omni48_clip_model_seconds_proxy", 0.0)) == 768.0
        and float(token_quota_summary.get("omni48_retry_clip_model_seconds_ceiling", 0.0)) == 1536.0
        and int(token_quota_summary.get("max_attempted_requests", 0)) == 764
        and int(token_quota_summary.get("live_calls_performed_by_builder", -1)) == 0
        and token_quota_summary.get("no_secret_values_written") is True
        and token_quota_summary.get("no_new_metric_claim") is True
        and required_token_quota_ids.issubset(token_quota_rows)
        and token_quota_rows.get("deepseek_resume_after_top3", {}).get("priority") == "P0"
        and token_quota_rows.get("deepseek_resume_after_top3", {}).get("max_attempted_requests") == 278
        and token_quota_rows.get("qwen_full_backup", {}).get("max_attempted_requests") == 294
        and token_quota_rows.get("omni48_label_only", {}).get("retry_clip_model_seconds_ceiling") == 1536.0
        and all(row.get("claim_status", "").endswith("no_live_metric_claim") for row in token_quota_rows.values()),
        json.dumps(
            {
                "contract": live_token_quota_budget.get("runtime_contract"),
                "status": live_token_quota_budget.get("status"),
                "summary": token_quota_summary,
                "missing_ids": sorted(required_token_quota_ids - set(token_quota_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    failure_recovery_summary = live_failure_recovery.get("summary", {})
    failure_recovery_rows = {row.get("scenario_id"): row for row in live_failure_recovery.get("rows", [])}
    required_failure_recovery_ids = {
        "missing_credentials",
        "known_provider_quota_blocker",
        "missing_output_clean_run",
        "partial_or_invalid_output",
        "retry_exhausted_errors",
        "scoring_blocked_missing_output",
        "promotion_blocked",
        "token_or_attempt_budget_ceiling",
    }
    add_check(
        checks,
        "live_failure_recovery_playbook_contract",
        live_failure_recovery.get("runtime_contract") == "live_failure_recovery_playbook_no_live_calls"
        and live_failure_recovery.get("status") == "ready_waiting_credentials_or_live_outputs"
        and live_failure_recovery.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("live_retry_budget_audit") == "live_retry_budget_audit_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("live_token_quota_budget_audit") == "live_token_quota_budget_audit_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("live_resume_state_audit") == "live_resume_state_audit_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("live_runtime_environment_audit") == "live_runtime_environment_audit_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and live_failure_recovery.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and int(failure_recovery_summary.get("scenario_count", 0)) == 8
        and int(failure_recovery_summary.get("p0_scenario_count", 0)) == 7
        and int(failure_recovery_summary.get("current_blocker_scenarios", 0)) == 5
        and int(failure_recovery_summary.get("future_recovery_scenarios", 0)) == 2
        and int(failure_recovery_summary.get("planning_guardrail_scenarios", 0)) == 1
        and int(failure_recovery_summary.get("ready_recovery_actions", 0)) == 8
        and int(failure_recovery_summary.get("planned_live_calls", 0)) == 382
        and int(failure_recovery_summary.get("max_attempted_requests", 0)) == 764
        and int(failure_recovery_summary.get("llm_retry_token_proxy_ceiling", 0)) == 1658856
        and int(failure_recovery_summary.get("missing_output_surfaces", 0)) == 3
        and int(failure_recovery_summary.get("partial_or_invalid_surfaces", -1)) == 0
        and int(failure_recovery_summary.get("ready_to_score_steps", -1)) == 0
        and int(failure_recovery_summary.get("ready_to_promote_count", -1)) == 0
        and int(failure_recovery_summary.get("known_provider_quota_blockers", 0)) == 1
        and failure_recovery_summary.get("credential_ready") is False
        and int(failure_recovery_summary.get("skip_existing_supported_surfaces", 0)) == 3
        and int(failure_recovery_summary.get("bounded_retry_supported_surfaces", 0)) == 3
        and int(failure_recovery_summary.get("live_calls_performed_by_builder", -1)) == 0
        and failure_recovery_summary.get("no_secret_values_written") is True
        and failure_recovery_summary.get("no_new_metric_claim") is True
        and required_failure_recovery_ids.issubset(failure_recovery_rows)
        and failure_recovery_rows.get("missing_credentials", {}).get("status") == "current_blocker"
        and failure_recovery_rows.get("known_provider_quota_blocker", {}).get("claim_status") == "blocked_by_provider_quota_no_metric_claim"
        and failure_recovery_rows.get("missing_output_clean_run", {}).get("status") == "current_blocker"
        and failure_recovery_rows.get("partial_or_invalid_output", {}).get("status") == "future_recovery_path"
        and failure_recovery_rows.get("retry_exhausted_errors", {}).get("status") == "future_recovery_path"
        and failure_recovery_rows.get("scoring_blocked_missing_output", {}).get("current_observation") == "ready_to_score_steps=0/5"
        and failure_recovery_rows.get("promotion_blocked", {}).get("claim_status") == "promote_only_after_output_audit_scoring_slo_and_traceability_pass"
        and failure_recovery_rows.get("token_or_attempt_budget_ceiling", {}).get("priority") == "P1",
        json.dumps(
            {
                "contract": live_failure_recovery.get("runtime_contract"),
                "status": live_failure_recovery.get("status"),
                "summary": failure_recovery_summary,
                "missing_ids": sorted(required_failure_recovery_ids - set(failure_recovery_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    metric_extraction_summary = live_metric_extraction.get("summary", {})
    metric_extraction_rows = {row.get("metric_id"): row for row in live_metric_extraction.get("rows", [])}
    required_metric_extraction_ids = {
        "deepseek_resume_safety_zero_harm",
        "deepseek_resume_call_latency",
        "deepseek_resume_token_multiplier",
        "qwen_backup_safety_zero_harm",
        "qwen_backup_call_latency",
        "omni48_label_quality",
        "omni48_call_latency",
        "post_live_promotion_sync",
    }
    add_check(
        checks,
        "live_metric_extraction_contract",
        live_metric_extraction.get("runtime_contract") == "live_metric_extraction_contract_no_live_calls"
        and live_metric_extraction.get("status") == "blocked_waiting_live_outputs"
        and live_metric_extraction.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and live_metric_extraction.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_metric_extraction.get("source_contracts", {}).get("live_postrun_metrics_closure") == "live_postrun_metrics_closure_no_live_calls"
        and live_metric_extraction.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and live_metric_extraction.get("source_contracts", {}).get("live_execution_timing_plan") == "live_execution_timing_plan_no_live_calls"
        and live_metric_extraction.get("source_contracts", {}).get("live_failure_recovery_playbook") == "live_failure_recovery_playbook_no_live_calls"
        and int(metric_extraction_summary.get("metric_contract_count", 0)) == 8
        and int(metric_extraction_summary.get("p0_metric_contracts", 0)) == 4
        and int(metric_extraction_summary.get("time_metric_contracts", 0)) == 3
        and int(metric_extraction_summary.get("safety_metric_contracts", 0)) == 2
        and int(metric_extraction_summary.get("omni_metric_contracts", 0)) == 2
        and int(metric_extraction_summary.get("promotion_metric_contracts", 0)) == 1
        and int(metric_extraction_summary.get("blocked_metric_contracts", 0)) == 8
        and int(metric_extraction_summary.get("required_scoring_steps", 0)) == 5
        and int(metric_extraction_summary.get("ready_to_score_steps", -1)) == 0
        and int(metric_extraction_summary.get("expected_live_calls", 0)) == 382
        and int(metric_extraction_summary.get("expected_input_calls", 0)) == 676
        and int(metric_extraction_summary.get("deepseek_resume_expected_calls", 0)) == 139
        and int(metric_extraction_summary.get("qwen_full_expected_calls", 0)) == 147
        and int(metric_extraction_summary.get("omni48_expected_calls", 0)) == 96
        and int(metric_extraction_summary.get("missing_output_surfaces", 0)) == 3
        and int(metric_extraction_summary.get("ready_to_promote_count", -1)) == 0
        and float(metric_extraction_summary.get("deepseek_estimated_wall_seconds", 0.0)) == 384.444
        and int(metric_extraction_summary.get("live_calls_performed_by_builder", -1)) == 0
        and metric_extraction_summary.get("no_scoring_commands_executed") is True
        and metric_extraction_summary.get("no_secret_values_written") is True
        and metric_extraction_summary.get("no_new_metric_claim") is True
        and required_metric_extraction_ids.issubset(metric_extraction_rows)
        and metric_extraction_rows.get("deepseek_resume_safety_zero_harm", {}).get("promotion_gate") == "deepseek_split20_resume_safety"
        and "harmful_accepts" in metric_extraction_rows.get("deepseek_resume_safety_zero_harm", {}).get("statistic_fields", "")
        and metric_extraction_rows.get("deepseek_resume_call_latency", {}).get("metric_family") == "llm_latency"
        and "wall_seconds" in metric_extraction_rows.get("deepseek_resume_call_latency", {}).get("statistic_fields", "")
        and metric_extraction_rows.get("deepseek_resume_token_multiplier", {}).get("claim_status") == "planning_support_until_full_split20_comparison_exists"
        and metric_extraction_rows.get("qwen_backup_call_latency", {}).get("claim_status") == "fallback_only_not_primary_latency_claim"
        and metric_extraction_rows.get("omni48_label_quality", {}).get("claim_status") == "label_only_no_timeline_writeback"
        and "p95_call_seconds" in metric_extraction_rows.get("omni48_call_latency", {}).get("statistic_fields", "")
        and metric_extraction_rows.get("post_live_promotion_sync", {}).get("claim_status") == "promote_only_after_output_audit_scoring_slo_and_traceability_pass",
        json.dumps(
            {
                "contract": live_metric_extraction.get("runtime_contract"),
                "status": live_metric_extraction.get("status"),
                "summary": metric_extraction_summary,
                "missing_ids": sorted(required_metric_extraction_ids - set(metric_extraction_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    output_schema_summary = live_output_schema.get("summary", {})
    output_schema_rows = {row.get("schema_id"): row for row in live_output_schema.get("rows", [])}
    required_output_schema_ids = {
        "deepseek_resume_llm_success_jsonl",
        "qwen_full_llm_success_jsonl",
        "omni48_success_jsonl",
        "bounded_retry_error_row",
        "llm_safety_summary",
        "split20_comparison_summary",
        "omni48_metric_summary",
        "promotion_traceability_summary",
    }
    add_check(
        checks,
        "live_output_schema_contract",
        live_output_schema.get("runtime_contract") == "live_output_schema_contract_no_live_calls"
        and live_output_schema.get("status") == "blocked_waiting_live_outputs"
        and live_output_schema.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and live_output_schema.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and live_output_schema.get("source_contracts", {}).get("live_metric_extraction_contract") == "live_metric_extraction_contract_no_live_calls"
        and live_output_schema.get("source_contracts", {}).get("live_failure_recovery_playbook") == "live_failure_recovery_playbook_no_live_calls"
        and live_output_schema.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_output_schema.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and int(output_schema_summary.get("schema_contract_count", 0)) == 8
        and int(output_schema_summary.get("p0_schema_contracts", 0)) == 5
        and int(output_schema_summary.get("live_output_schema_contracts", 0)) == 3
        and int(output_schema_summary.get("scoring_output_schema_contracts", 0)) == 3
        and int(output_schema_summary.get("promotion_schema_contracts", 0)) == 1
        and int(output_schema_summary.get("error_row_schema_contracts", 0)) == 1
        and int(output_schema_summary.get("required_field_count", 0)) == 62
        and int(output_schema_summary.get("expected_live_output_rows", 0)) == 382
        and int(output_schema_summary.get("missing_output_surfaces", 0)) == 3
        and int(output_schema_summary.get("ready_to_score_steps", -1)) == 0
        and int(output_schema_summary.get("metric_contract_count", 0)) == 8
        and int(output_schema_summary.get("live_calls_performed_by_builder", -1)) == 0
        and output_schema_summary.get("no_schema_validation_executed_on_missing_outputs") is True
        and output_schema_summary.get("no_secret_values_written") is True
        and output_schema_summary.get("no_new_metric_claim") is True
        and required_output_schema_ids.issubset(output_schema_rows)
        and output_schema_rows.get("deepseek_resume_llm_success_jsonl", {}).get("expected_rows") == 139
        and "patch_decisions" in output_schema_rows.get("deepseek_resume_llm_success_jsonl", {}).get("required_fields", "")
        and output_schema_rows.get("qwen_full_llm_success_jsonl", {}).get("expected_rows") == 147
        and output_schema_rows.get("omni48_success_jsonl", {}).get("expected_rows") == 96
        and "schema_ok" in output_schema_rows.get("omni48_success_jsonl", {}).get("required_fields", "")
        and "retry_backoff_seconds" in output_schema_rows.get("bounded_retry_error_row", {}).get("required_fields", "")
        and "harmful_accepts" in output_schema_rows.get("llm_safety_summary", {}).get("required_fields", "")
        and "token_multiplier" in output_schema_rows.get("split20_comparison_summary", {}).get("required_fields", "")
        and "p95_call_seconds" in output_schema_rows.get("omni48_metric_summary", {}).get("required_fields", "")
        and output_schema_rows.get("promotion_traceability_summary", {}).get("claim_status") == "required_before_report_ppt_claim_promotion",
        json.dumps(
            {
                "contract": live_output_schema.get("runtime_contract"),
                "status": live_output_schema.get("status"),
                "summary": output_schema_summary,
                "missing_ids": sorted(required_output_schema_ids - set(output_schema_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    acceptance_summary = post_live_acceptance_scorecard.get("summary", {})
    acceptance_rows = {
        row.get("scorecard_id"): row
        for row in post_live_acceptance_scorecard.get("rows", [])
    }
    required_acceptance_ids = {
        "current_claim_now_slo_preserve",
        "deepseek_resume_output_coverage",
        "deepseek_resume_output_schema",
        "deepseek_resume_safety_zero_harm",
        "deepseek_split20_latency_evidence",
        "omni48_output_schema",
        "omni48_label_metrics",
        "qwen_backup_fallback_boundary",
        "report_ppt_promotion_sync",
    }
    add_check(
        checks,
        "post_live_acceptance_scorecard_contract",
        post_live_acceptance_scorecard.get("runtime_contract") == "post_live_acceptance_scorecard_no_live_calls"
        and post_live_acceptance_scorecard.get("status") == "blocked_waiting_live_outputs"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("stage_latency_slo_audit") == "stage_latency_slo_audit_from_latency_ledger_no_live_calls"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("live_scoring_readiness") == "live_scoring_readiness_no_live_calls"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("live_metric_extraction_contract") == "live_metric_extraction_contract_no_live_calls"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("live_output_schema_contract") == "live_output_schema_contract_no_live_calls"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_acceptance_scorecard.get("source_contracts", {}).get("latency_risk_margin_audit") == "latency_risk_margin_audit_from_slo_no_live_calls"
        and int(acceptance_summary.get("scorecard_rows", 0)) == 9
        and int(acceptance_summary.get("p0_scorecard_rows", 0)) == 6
        and int(acceptance_summary.get("p1_scorecard_rows", 0)) == 3
        and int(acceptance_summary.get("preserve_pass_rows", 0)) == 2
        and int(acceptance_summary.get("blocked_rows", 0)) == 6
        and int(acceptance_summary.get("fallback_only_rows", 0)) == 1
        and int(acceptance_summary.get("claim_now_slo_pass", 0)) == 4
        and int(acceptance_summary.get("claim_now_slo_rows", 0)) == 4
        and float(acceptance_summary.get("guard_p95_margin_seconds", 0.0)) == 1.151
        and int(acceptance_summary.get("expected_live_calls", 0)) == 382
        and int(acceptance_summary.get("missing_output_surfaces", 0)) == 3
        and int(acceptance_summary.get("ready_to_score_steps", -1)) == 0
        and int(acceptance_summary.get("metric_contract_count", 0)) == 8
        and int(acceptance_summary.get("schema_contract_count", 0)) == 8
        and int(acceptance_summary.get("ready_to_promote_count", -1)) == 0
        and int(acceptance_summary.get("traceability_rows", 0)) == 54
        and int(acceptance_summary.get("live_calls_performed_by_builder", -1)) == 0
        and acceptance_summary.get("no_secret_values_written") is True
        and acceptance_summary.get("no_new_metric_claim") is True
        and required_acceptance_ids.issubset(acceptance_rows)
        and acceptance_rows.get("deepseek_resume_output_coverage", {}).get("current_status") == "blocked_missing_output"
        and "harmful_accepts == 0" in acceptance_rows.get("deepseek_resume_safety_zero_harm", {}).get("pass_condition", "")
        and "no timeline writeback" in acceptance_rows.get("omni48_label_metrics", {}).get("claim_effect", "")
        and acceptance_rows.get("qwen_backup_fallback_boundary", {}).get("current_status") == "fallback_only"
        and acceptance_rows.get("report_ppt_promotion_sync", {}).get("current_status") == "preserve_pass",
        json.dumps(
            {
                "contract": post_live_acceptance_scorecard.get("runtime_contract"),
                "status": post_live_acceptance_scorecard.get("status"),
                "summary": acceptance_summary,
                "missing_ids": sorted(required_acceptance_ids - set(acceptance_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    latency_claim_summary = post_live_latency_claim_matrix.get("summary", {})
    latency_claim_rows = {
        row.get("latency_claim_id"): row
        for row in post_live_latency_claim_matrix.get("rows", [])
    }
    required_latency_claim_ids = {
        "fast_first_output_current",
        "rule_writeback_current",
        "runtime_safe_guard_current",
        "llm_review_signal_current",
        "deepseek_split20_full_surface",
        "omni48_label_latency",
        "qwen_full_backup_latency",
        "report_ppt_latency_sync",
    }
    add_check(
        checks,
        "post_live_latency_claim_matrix_contract",
        post_live_latency_claim_matrix.get("runtime_contract") == "post_live_latency_claim_matrix_no_live_calls"
        and post_live_latency_claim_matrix.get("status") == "blocked_waiting_live_outputs"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("runtime_latency_budget_ledger") == "runtime_latency_budget_ledger_from_existing_artifacts"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("stage_latency_slo_audit") == "stage_latency_slo_audit_from_latency_ledger_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("latency_risk_margin_audit") == "latency_risk_margin_audit_from_slo_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("live_execution_timing_plan") == "live_execution_timing_plan_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("live_metric_extraction_contract") == "live_metric_extraction_contract_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("post_live_acceptance_scorecard") == "post_live_acceptance_scorecard_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("live_execution_handoff_packet") == "live_execution_handoff_packet_no_live_calls"
        and post_live_latency_claim_matrix.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(latency_claim_summary.get("latency_claim_rows", 0)) == 8
        and int(latency_claim_summary.get("p0_latency_claim_rows", 0)) == 6
        and int(latency_claim_summary.get("p1_latency_claim_rows", 0)) == 2
        and int(latency_claim_summary.get("claim_now_preserve_rows", 0)) == 4
        and int(latency_claim_summary.get("blocked_or_waiting_rows", 0)) == 4
        and int(latency_claim_summary.get("fallback_only_rows", 0)) == 1
        and int(latency_claim_summary.get("label_only_rows", 0)) == 1
        and int(latency_claim_summary.get("tight_margin_rows", 0)) == 1
        and int(latency_claim_summary.get("claim_now_slo_pass", 0)) == 4
        and int(latency_claim_summary.get("claim_now_slo_rows", 0)) == 4
        and float(latency_claim_summary.get("guard_p95_margin_seconds", 0.0)) == 1.151
        and int(latency_claim_summary.get("expected_live_calls", 0)) == 382
        and int(latency_claim_summary.get("missing_output_surfaces", 0)) == 3
        and float(latency_claim_summary.get("deepseek_estimated_wall_seconds", 0.0)) == 384.444
        and float(latency_claim_summary.get("qwen_estimated_wall_seconds", 0.0)) == 836.456
        and int(latency_claim_summary.get("omni48_label_only_calls", 0)) == 96
        and int(latency_claim_summary.get("ready_to_promote_count", -1)) == 0
        and int(latency_claim_summary.get("traceability_rows", 0)) == 54
        and int(latency_claim_summary.get("handoff_packet_rows", 0)) == 7
        and int(latency_claim_summary.get("live_calls_performed_by_builder", -1)) == 0
        and latency_claim_summary.get("no_secret_values_written") is True
        and latency_claim_summary.get("no_new_metric_claim") is True
        and required_latency_claim_ids.issubset(latency_claim_rows)
        and latency_claim_rows.get("runtime_safe_guard_current", {}).get("claim_boundary") == "claim_now_latency_preserve_no_broader_claim"
        and latency_claim_rows.get("deepseek_split20_full_surface", {}).get("promotion_gate") == "deepseek_split20_resume_latency"
        and latency_claim_rows.get("omni48_label_latency", {}).get("claim_boundary") == "label_only_latency_not_guard_or_timeline_claim"
        and latency_claim_rows.get("qwen_full_backup_latency", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and latency_claim_rows.get("report_ppt_latency_sync", {}).get("claim_boundary") == "report_ppt_sync_required_before_latency_claim_promotion",
        json.dumps(
            {
                "contract": post_live_latency_claim_matrix.get("runtime_contract"),
                "status": post_live_latency_claim_matrix.get("status"),
                "summary": latency_claim_summary,
                "missing_ids": sorted(required_latency_claim_ids - set(latency_claim_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    time_stats_summary = post_live_time_metric_statistics.get("summary", {})
    time_stats_rows = {
        row.get("time_metric_id"): row
        for row in post_live_time_metric_statistics.get("rows", [])
    }
    required_time_stat_ids = {
        "fast_first_output_latency_current",
        "rule_writeback_latency_current",
        "runtime_safe_guard_latency_current",
        "llm_review_signal_latency_current",
        "deepseek_resume_call_latency_stats",
        "deepseek_parent_completion_latency_stats",
        "qwen_backup_latency_stats",
        "omni48_label_latency_stats",
        "report_ppt_time_claim_refresh",
    }
    add_check(
        checks,
        "post_live_time_metric_statistics_plan_contract",
        post_live_time_metric_statistics.get("runtime_contract") == "post_live_time_metric_statistics_plan_no_live_calls"
        and post_live_time_metric_statistics.get("status") == "blocked_waiting_live_outputs"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("runtime_latency_budget_ledger") == "runtime_latency_budget_ledger_from_existing_artifacts"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("stage_latency_slo_audit") == "stage_latency_slo_audit_from_latency_ledger_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("latency_risk_margin_audit") == "latency_risk_margin_audit_from_slo_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("live_metric_extraction_contract") == "live_metric_extraction_contract_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("live_output_schema_contract") == "live_output_schema_contract_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("post_live_latency_claim_matrix") == "post_live_latency_claim_matrix_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("live_execution_bundle") == "live_execution_bundle_no_live_calls_no_secret_values"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_time_metric_statistics.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(time_stats_summary.get("time_stat_rows", 0)) == 9
        and int(time_stats_summary.get("p0_time_stat_rows", 0)) == 7
        and int(time_stats_summary.get("p1_time_stat_rows", 0)) == 2
        and int(time_stats_summary.get("claim_now_preserve_rows", 0)) == 4
        and int(time_stats_summary.get("blocked_or_waiting_rows", 0)) == 5
        and int(time_stats_summary.get("fallback_only_rows", 0)) == 1
        and int(time_stats_summary.get("label_only_rows", 0)) == 1
        and int(time_stats_summary.get("stage_arrival_rows", 0)) == 4
        and int(time_stats_summary.get("post_live_stat_rows", 0)) == 4
        and int(time_stats_summary.get("formula_count", 0)) == 9
        and int(time_stats_summary.get("claim_now_slo_pass", 0)) == 4
        and int(time_stats_summary.get("claim_now_slo_rows", 0)) == 4
        and float(time_stats_summary.get("guard_p95_margin_seconds", 0.0)) == 1.151
        and int(time_stats_summary.get("metric_contract_count", 0)) == 8
        and int(time_stats_summary.get("schema_contract_count", 0)) == 8
        and int(time_stats_summary.get("latency_claim_rows", 0)) == 8
        and int(time_stats_summary.get("planned_live_calls", 0)) == 382
        and int(time_stats_summary.get("p0_planned_live_calls", 0)) == 139
        and int(time_stats_summary.get("expected_live_calls", 0)) == 382
        and int(time_stats_summary.get("missing_output_surfaces", 0)) == 3
        and int(time_stats_summary.get("ready_to_promote_count", -1)) == 0
        and int(time_stats_summary.get("traceability_rows", 0)) == 54
        and int(time_stats_summary.get("live_calls_performed_by_builder", -1)) == 0
        and time_stats_summary.get("no_scoring_commands_executed") is True
        and time_stats_summary.get("no_secret_values_written") is True
        and time_stats_summary.get("no_new_metric_claim") is True
        and required_time_stat_ids.issubset(time_stats_rows)
        and time_stats_rows.get("fast_first_output_latency_current", {}).get("claim_boundary") == "claim_now_latency_preserve"
        and time_stats_rows.get("runtime_safe_guard_latency_current", {}).get("claim_boundary") == "claim_now_latency_preserve_no_broader_claim"
        and time_stats_rows.get("deepseek_resume_call_latency_stats", {}).get("expected_rows") == 139
        and "avg(call_seconds)" in time_stats_rows.get("deepseek_resume_call_latency_stats", {}).get("statistic_formula", "")
        and time_stats_rows.get("deepseek_parent_completion_latency_stats", {}).get("claim_boundary") == "required_before_full_surface_latency_claim"
        and time_stats_rows.get("qwen_backup_latency_stats", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and time_stats_rows.get("omni48_label_latency_stats", {}).get("claim_boundary") == "label_only_latency_not_guard_or_timeline_claim"
        and time_stats_rows.get("report_ppt_time_claim_refresh", {}).get("claim_boundary") == "report_ppt_sync_required_before_time_metric_claim_promotion",
        json.dumps(
            {
                "contract": post_live_time_metric_statistics.get("runtime_contract"),
                "status": post_live_time_metric_statistics.get("status"),
                "summary": time_stats_summary,
                "missing_ids": sorted(required_time_stat_ids - set(time_stats_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    extractor_summary = post_live_time_metric_extractor.get("summary", {})
    extractor_rows = {
        row.get("time_metric_id"): row
        for row in post_live_time_metric_extractor.get("rows", [])
    }
    required_extractor_ids = {
        "deepseek_resume_time_metric_extract",
        "qwen_backup_time_metric_extract",
        "omni48_label_time_metric_extract",
    }
    add_check(
        checks,
        "post_live_time_metric_extractor_contract",
        post_live_time_metric_extractor.get("runtime_contract") == "post_live_time_metric_extractor_no_live_calls"
        and post_live_time_metric_extractor.get("status") == "blocked_waiting_live_outputs"
        and post_live_time_metric_extractor.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and post_live_time_metric_extractor.get("source_contracts", {}).get("post_live_time_metric_statistics_plan") == "post_live_time_metric_statistics_plan_no_live_calls"
        and post_live_time_metric_extractor.get("source_contracts", {}).get("live_output_schema_contract") == "live_output_schema_contract_no_live_calls"
        and post_live_time_metric_extractor.get("source_contracts", {}).get("live_execution_bundle") == "live_execution_bundle_no_live_calls_no_secret_values"
        and post_live_time_metric_extractor.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_time_metric_extractor.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(extractor_summary.get("extractor_rows", 0)) == 3
        and int(extractor_summary.get("p0_extractor_rows", 0)) == 1
        and int(extractor_summary.get("p1_extractor_rows", 0)) == 2
        and int(extractor_summary.get("missing_output_rows", 0)) == 3
        and int(extractor_summary.get("blocked_extractor_rows", 0)) == 3
        and int(extractor_summary.get("ready_time_metric_rows", -1)) == 0
        and int(extractor_summary.get("computed_time_metric_rows", -1)) == 0
        and int(extractor_summary.get("expected_rows_total", 0)) == 382
        and int(extractor_summary.get("observed_rows_total", -1)) == 0
        and int(extractor_summary.get("successful_rows_total", -1)) == 0
        and int(extractor_summary.get("parse_error_rows", -1)) == 0
        and int(extractor_summary.get("retry_rows_total", -1)) == 0
        and int(extractor_summary.get("expected_live_calls", 0)) == 382
        and int(extractor_summary.get("missing_output_surfaces", 0)) == 3
        and int(extractor_summary.get("time_stat_rows", 0)) == 9
        and int(extractor_summary.get("schema_contract_count", 0)) == 8
        and int(extractor_summary.get("planned_live_calls", 0)) == 382
        and int(extractor_summary.get("p0_planned_live_calls", 0)) == 139
        and int(extractor_summary.get("ready_to_promote_count", -1)) == 0
        and int(extractor_summary.get("traceability_rows", 0)) == 54
        and int(extractor_summary.get("live_calls_performed_by_builder", -1)) == 0
        and extractor_summary.get("no_scoring_commands_executed") is True
        and extractor_summary.get("no_secret_values_written") is True
        and extractor_summary.get("no_new_metric_claim") is True
        and required_extractor_ids.issubset(extractor_rows)
        and extractor_rows.get("deepseek_resume_time_metric_extract", {}).get("status") == "blocked_missing_output"
        and extractor_rows.get("deepseek_resume_time_metric_extract", {}).get("expected_rows") == 139
        and extractor_rows.get("deepseek_resume_time_metric_extract", {}).get("claim_boundary") == "not_claimable_until_resume_output_audit_scoring_and_traceability"
        and extractor_rows.get("qwen_backup_time_metric_extract", {}).get("status") == "blocked_missing_output"
        and extractor_rows.get("qwen_backup_time_metric_extract", {}).get("expected_rows") == 147
        and extractor_rows.get("qwen_backup_time_metric_extract", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and extractor_rows.get("omni48_label_time_metric_extract", {}).get("status") == "blocked_missing_output"
        and extractor_rows.get("omni48_label_time_metric_extract", {}).get("expected_rows") == 96
        and extractor_rows.get("omni48_label_time_metric_extract", {}).get("claim_boundary") == "label_only_latency_not_guard_or_timeline_claim",
        json.dumps(
            {
                "contract": post_live_time_metric_extractor.get("runtime_contract"),
                "status": post_live_time_metric_extractor.get("status"),
                "summary": extractor_summary,
                "missing_ids": sorted(required_extractor_ids - set(extractor_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    time_receipt_summary = post_live_time_metric_receipt.get("summary", {})
    time_receipt_rows = {
        row.get("receipt_id"): row
        for row in post_live_time_metric_receipt.get("rows", [])
    }
    required_time_receipt_ids = {
        "time_statistics_plan_alignment",
        "time_extractor_surface_coverage",
        "time_metric_computation_receipt",
        "time_metric_parse_retry_quality",
        "scoring_receipt_dependency",
        "promotion_preflight_after_time_metrics",
    }
    add_check(
        checks,
        "post_live_time_metric_receipt_audit_contract",
        post_live_time_metric_receipt.get("runtime_contract") == "post_live_time_metric_receipt_audit_no_live_or_scoring_calls_no_secret_values"
        and post_live_time_metric_receipt.get("status") == "blocked_waiting_time_metric_evidence"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("post_live_time_metric_statistics_plan") == "post_live_time_metric_statistics_plan_no_live_calls"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("post_live_time_metric_extractor") == "post_live_time_metric_extractor_no_live_calls"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("post_live_scoring_receipt_audit") == "post_live_scoring_receipt_audit_no_live_or_scoring_calls_no_secret_values"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("post_live_scoring_output_audit") == "post_live_scoring_output_audit_no_scoring_calls"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("post_live_latency_claim_matrix") == "post_live_latency_claim_matrix_no_live_calls"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("post_live_promotion_preflight_audit") == "post_live_promotion_preflight_audit_no_live_or_scoring_calls"
        and post_live_time_metric_receipt.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(time_receipt_summary.get("receipt_rows", 0)) == 6
        and int(time_receipt_summary.get("pass_rows", -1)) == 1
        and int(time_receipt_summary.get("blocked_rows", 0)) == 5
        and int(time_receipt_summary.get("missing_source_rows", -1)) == 0
        and int(time_receipt_summary.get("time_stat_rows", 0)) == 9
        and int(time_receipt_summary.get("post_live_stat_rows", 0)) == 4
        and int(time_receipt_summary.get("formula_count", 0)) == 9
        and int(time_receipt_summary.get("extractor_rows", 0)) == 3
        and int(time_receipt_summary.get("ready_time_metric_rows", -1)) == 0
        and int(time_receipt_summary.get("computed_time_metric_rows", -1)) == 0
        and int(time_receipt_summary.get("missing_output_rows", 0)) == 3
        and int(time_receipt_summary.get("expected_rows_total", 0)) == 382
        and int(time_receipt_summary.get("observed_rows_total", -1)) == 0
        and int(time_receipt_summary.get("successful_rows_total", -1)) == 0
        and int(time_receipt_summary.get("parse_error_rows", -1)) == 0
        and int(time_receipt_summary.get("retry_rows_total", -1)) == 0
        and int(time_receipt_summary.get("expected_live_calls", 0)) == 382
        and int(time_receipt_summary.get("missing_output_surfaces", 0)) == 3
        and time_receipt_summary.get("scoring_receipt_ready_for_promotion_review") is False
        and time_receipt_summary.get("scoring_execute_record_exists") is False
        and int(time_receipt_summary.get("scoring_output_promotion_ready_rows", -1)) == 0
        and time_receipt_summary.get("promotion_preflight_ready") is False
        and time_receipt_summary.get("ready_for_time_claim_promotion") is False
        and int(time_receipt_summary.get("traceability_rows", 0)) == 54
        and int(time_receipt_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(time_receipt_summary.get("live_calls_performed_by_builder", -1)) == 0
        and time_receipt_summary.get("no_live_calls_performed_by_auditor") is True
        and time_receipt_summary.get("no_scoring_commands_executed_by_auditor") is True
        and time_receipt_summary.get("no_secret_values_written") is True
        and time_receipt_summary.get("no_new_metric_claim") is True
        and required_time_receipt_ids.issubset(time_receipt_rows)
        and time_receipt_rows.get("time_statistics_plan_alignment", {}).get("status") == "pass"
        and all(
            row.get("status") == "blocked"
            for receipt_id, row in time_receipt_rows.items()
            if receipt_id != "time_statistics_plan_alignment"
        )
        and all(row.get("source_artifacts_exist") is True for row in time_receipt_rows.values()),
        json.dumps(
            {
                "contract": post_live_time_metric_receipt.get("runtime_contract"),
                "status": post_live_time_metric_receipt.get("status"),
                "summary": time_receipt_summary,
                "missing_ids": sorted(required_time_receipt_ids - set(time_receipt_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    promotion_preflight_summary = post_live_promotion_preflight.get("summary", {})
    promotion_preflight_rows = {
        row.get("preflight_id"): row
        for row in post_live_promotion_preflight.get("rows", [])
    }
    required_promotion_preflight_ids = {
        "live_execution_preflight",
        "scoring_execution_preflight",
        "time_metric_preflight",
        "claim_promotion_gate_preflight",
        "report_ppt_traceability_preflight",
        "claim_boundary_safety_preflight",
    }
    add_check(
        checks,
        "post_live_promotion_preflight_audit_contract",
        post_live_promotion_preflight.get("runtime_contract") == "post_live_promotion_preflight_audit_no_live_or_scoring_calls"
        and post_live_promotion_preflight.get("status") == "blocked_waiting_post_live_evidence"
        and post_live_promotion_preflight.get("source_contracts", {}).get("live_execution_launcher") == "live_execution_launcher_dry_run_no_live_calls"
        and post_live_promotion_preflight.get("source_contracts", {}).get("post_live_scoring_launcher") == "post_live_scoring_launcher_dry_run_no_scoring_calls"
        and post_live_promotion_preflight.get("source_contracts", {}).get("post_live_scoring_output_audit") == "post_live_scoring_output_audit_no_scoring_calls"
        and post_live_promotion_preflight.get("source_contracts", {}).get("post_live_time_metric_extractor") == "post_live_time_metric_extractor_no_live_calls"
        and post_live_promotion_preflight.get("source_contracts", {}).get("post_live_time_metric_statistics_plan") == "post_live_time_metric_statistics_plan_no_live_calls"
        and post_live_promotion_preflight.get("source_contracts", {}).get("post_live_claim_promotion_gate") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_promotion_preflight.get("source_contracts", {}).get("report_ppt_traceability") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and int(promotion_preflight_summary.get("preflight_rows", 0)) == 6
        and int(promotion_preflight_summary.get("pass_rows", 0)) == 2
        and int(promotion_preflight_summary.get("blocked_rows", 0)) == 4
        and int(promotion_preflight_summary.get("missing_source_rows", -1)) == 0
        and promotion_preflight_summary.get("ready_for_promotion_review") is False
        and int(promotion_preflight_summary.get("post_live_ready_to_promote_count", -1)) == 0
        and int(promotion_preflight_summary.get("promotion_blocked_count", 0)) == 5
        and int(promotion_preflight_summary.get("promotion_preserve_count", 0)) == 2
        and int(promotion_preflight_summary.get("promotion_fallback_only_count", 0)) == 1
        and int(promotion_preflight_summary.get("scoring_output_promotion_ready_rows", -1)) == 0
        and int(promotion_preflight_summary.get("ready_time_metric_rows", -1)) == 0
        and int(promotion_preflight_summary.get("computed_time_metric_rows", -1)) == 0
        and int(promotion_preflight_summary.get("traceability_rows", 0)) == 54
        and int(promotion_preflight_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(promotion_preflight_summary.get("live_calls_performed_by_builder", -1)) == 0
        and promotion_preflight_summary.get("no_live_calls_performed") is True
        and promotion_preflight_summary.get("no_scoring_commands_executed") is True
        and promotion_preflight_summary.get("no_secret_values_written") is True
        and promotion_preflight_summary.get("no_new_metric_claim") is True
        and required_promotion_preflight_ids.issubset(promotion_preflight_rows)
        and promotion_preflight_rows.get("live_execution_preflight", {}).get("status") == "blocked"
        and promotion_preflight_rows.get("scoring_execution_preflight", {}).get("status") == "blocked"
        and promotion_preflight_rows.get("time_metric_preflight", {}).get("status") == "blocked"
        and promotion_preflight_rows.get("claim_promotion_gate_preflight", {}).get("status") == "blocked"
        and promotion_preflight_rows.get("report_ppt_traceability_preflight", {}).get("status") == "pass"
        and promotion_preflight_rows.get("claim_boundary_safety_preflight", {}).get("status") == "pass"
        and all(row.get("source_artifacts_exist") is True for row in promotion_preflight_rows.values()),
        json.dumps(
            {
                "contract": post_live_promotion_preflight.get("runtime_contract"),
                "status": post_live_promotion_preflight.get("status"),
                "summary": promotion_preflight_summary,
                "missing_ids": sorted(required_promotion_preflight_ids - set(promotion_preflight_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    parallel_summary = live_parallelism.get("summary", {})
    parallel_rows = {
        (row.get("policy"), int(row.get("workers", 0))): row
        for row in live_parallelism.get("rows", [])
    }
    add_check(
        checks,
        "live_parallelism_sensitivity_contract",
        live_parallelism.get("runtime_contract") == "live_parallelism_sensitivity_no_live_calls"
        and live_parallelism.get("status") == "planning_only_blocked_waiting_for_credentials_or_quota"
        and live_parallelism.get("source_contracts", {}).get("split_policy_optimization") == "split_policy_optimization_from_existing_artifacts_no_live_calls"
        and live_parallelism.get("source_contracts", {}).get("live_execution_timing_plan") == "live_execution_timing_plan_no_live_calls"
        and live_parallelism.get("source_contracts", {}).get("latency_risk_mitigation_plan") == "latency_risk_mitigation_plan_no_live_calls"
        and int(parallel_summary.get("row_count", 0)) == 20
        and int(parallel_summary.get("policy_count", 0)) == 5
        and int(parallel_summary.get("worker_count", 0)) == 4
        and parallel_summary.get("recommended_policy") == "max20"
        and int(parallel_summary.get("recommended_workers", 0)) == 8
        and float(parallel_summary.get("recommended_estimated_wall_seconds", 0.0)) == 384.444
        and int(parallel_summary.get("recommended_waves", 0)) == 18
        and int(parallel_summary.get("recommended_calls", 0)) == 139
        and parallel_summary.get("recommended_worker_risk") == "current_runbook_default"
        and float(parallel_summary.get("max20_worker12_estimated_wall_seconds", 0.0)) == 256.296
        and float(parallel_summary.get("max20_worker12_wall_gain_seconds", 0.0)) == 128.148
        and float(parallel_summary.get("stretch_max15_workers8_estimated_wall_seconds", 0.0)) == 415.081
        and parallel_summary.get("stretch_max15_requires_reexport") is True
        and float(parallel_summary.get("qwen_backup_workers8_estimated_wall_seconds", 0.0)) == 836.456
        and int(parallel_summary.get("live_calls_performed_by_builder", -1)) == 0
        and parallel_summary.get("no_secret_values_written") is True
        and parallel_summary.get("no_new_metric_claim") is True
        and parallel_rows.get(("max20", 8), {}).get("recommendation") == "recommended_p0_default"
        and parallel_rows.get(("max20", 12), {}).get("recommendation") == "speedup_candidate_after_quota_stable"
        and parallel_rows.get(("max15", 8), {}).get("recommendation") == "stretch_reexport_only"
        and parallel_rows.get(("max15", 8), {}).get("requires_reexport") is True
        and parallel_rows.get(("max20", 8), {}).get("claim_status") == "planning_only_no_live_metric_claim",
        json.dumps(
            {
                "contract": live_parallelism.get("runtime_contract"),
                "status": live_parallelism.get("status"),
                "summary": parallel_summary,
                "sample_rows": {
                    "max20_8": parallel_rows.get(("max20", 8), {}),
                    "max20_12": parallel_rows.get(("max20", 12), {}),
                    "max15_8": parallel_rows.get(("max15", 8), {}),
                },
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    selector_candidate_summary = selector_candidate_scan.get("summary", {})
    add_check(
        checks,
        "selector_true_heldout_candidate_scan_contract",
        selector_candidate_scan.get("runtime_contract") == "selector_true_heldout_candidate_scan_no_metric_claim"
        and selector_candidate_scan.get("status") == "not_enough_new_local_recordings"
        and int(selector_candidate_summary.get("local_recordings", 0)) == 8
        and int(selector_candidate_summary.get("development_recordings", 0)) == 8
        and int(selector_candidate_summary.get("local_dev_overlap", 0)) == 8
        and int(selector_candidate_summary.get("eligible_true_heldout_recordings", -1)) == 0
        and int(selector_candidate_summary.get("missing_new_recordings_to_minimum", 0)) == 8
        and selector_candidate_summary.get("no_metric_claim") is True
        and selector_candidate_summary.get("sealed_split_written") is False
        and not selector_candidate_scan.get("eligible_recordings"),
        json.dumps(
            {
                "contract": selector_candidate_scan.get("runtime_contract"),
                "status": selector_candidate_scan.get("status"),
                "summary": selector_candidate_summary,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    selector_split_summary = selector_split_validation.get("summary", {})
    selector_split_gates = {gate.get("gate_id"): gate for gate in selector_split_validation.get("gates", [])}
    add_check(
        checks,
        "selector_true_heldout_split_validation_contract",
        selector_split_validation.get("runtime_contract") == "selector_true_heldout_split_validation_no_metric_claim"
        and selector_split_validation.get("status") == "blocked_waiting_for_valid_sealed_split"
        and selector_split_summary.get("split_exists") is False
        and int(selector_split_summary.get("rows", -1)) == 0
        and int(selector_split_summary.get("true_heldout_recordings", -1)) == 0
        and int(selector_split_summary.get("development_recordings", 0)) == 8
        and int(selector_split_summary.get("missing_required_columns", 0)) == 4
        and int(selector_split_summary.get("missing_new_recordings_to_minimum", 0)) == 8
        and selector_split_summary.get("no_metric_claim") is True
        and selector_split_summary.get("sealed_split_written_by_validator") is False
        and "missing_sealed_split_file" in selector_split_validation.get("blockers", [])
        and "not_enough_true_heldout_recordings" in selector_split_validation.get("blockers", [])
        and selector_split_gates.get("sealed_split_exists", {}).get("status") == "blocked"
        and selector_split_gates.get("no_metric_claim", {}).get("status") == "pass",
        json.dumps(
            {
                "contract": selector_split_validation.get("runtime_contract"),
                "status": selector_split_validation.get("status"),
                "summary": selector_split_summary,
                "blockers": selector_split_validation.get("blockers", []),
                "gate_statuses": {key: value.get("status") for key, value in selector_split_gates.items()},
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "selector_true_heldout_protocol_contract",
        selector_protocol.get("runtime_contract") == "selector_true_heldout_protocol_no_metric_claim"
        and selector_protocol.get("protocol_status") == "needs_new_recording_split"
        and selector_protocol.get("fixed_policy") == "ratio_le_0.65_else_uncovered"
        and selector_protocol.get("no_metric_claim") is True
        and int(selector_protocol.get("development_evidence", {}).get("dev_recording_count", 0)) == 8
        and int(selector_protocol.get("development_evidence", {}).get("dev_windows", 0)) == 120
        and selector_protocol.get("development_evidence", {}).get("development_scope") == "dev_only_validation_same_sampled_pool"
        and selector_protocol.get("sealed_split_state", {}).get("exists") is False
        and "missing_sealed_split_file" in selector_protocol.get("blockers", [])
        and "not_enough_true_heldout_recordings" in selector_protocol.get("blockers", [])
        and selector_protocol.get("candidate_scan", {}).get("runtime_contract") == "selector_true_heldout_candidate_scan_no_metric_claim"
        and int(selector_protocol.get("candidate_scan", {}).get("eligible_true_heldout_recordings", -1)) == 0
        and int(selector_protocol.get("candidate_scan", {}).get("missing_new_recordings_to_minimum", 0)) == 8
        and selector_protocol_gates.get("fixed_policy_before_scoring", {}).get("status") == "pass"
        and selector_protocol_gates.get("runtime_feature_surface", {}).get("status") == "pass",
        json.dumps(
            {
                "contract": selector_protocol.get("runtime_contract"),
                "status": selector_protocol.get("protocol_status"),
                "fixed_policy": selector_protocol.get("fixed_policy"),
                "dev": selector_protocol.get("development_evidence", {}),
                "sealed_split": selector_protocol.get("sealed_split_state", {}),
                "blockers": selector_protocol.get("blockers", []),
                "gate_statuses": {key: value.get("status") for key, value in selector_protocol_gates.items()},
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "split20_full_live_manifest_contract",
        split20_manifest.get("runtime_contract") == "split20_full_live_manifest_no_live_calls"
        and int(split20_manifest_summary.get("parent_windows", 0)) == 104
        and int(split20_manifest_summary.get("prompt_calls", 0)) == 147
        and int(split20_manifest_summary.get("patch_references", 0)) == 1917
        and int(split20_manifest_summary.get("split_parent_windows", 0)) == 40
        and int(split20_manifest_summary.get("deepseek_completed_parent_windows", 0)) == 3
        and int(split20_manifest_summary.get("deepseek_completed_calls", 0)) == 8
        and int(split20_manifest_summary.get("deepseek_top3_harmful_accepts", -1)) == 0
        and int(split20_manifest_summary.get("deepseek_quota_failed_parent_windows", 0)) == 2
        and int(split20_manifest_summary.get("deepseek_quota_failed_calls", 0)) == 4
        and split20_manifest_summary.get("deepseek_failure_type") == "AllocationQuota.FreeTierOnly"
        and int(split20_manifest_summary.get("deepseek_resume_parent_windows", 0)) == 101
        and int(split20_manifest_summary.get("deepseek_resume_required_calls_min", 0)) == 139
        and split20_manifest_summary.get("deepseek_full_surface_status") == "blocked_by_provider_quota_or_capacity"
        and int(split20_manifest_summary.get("qwen_backup_parent_windows", 0)) == 2
        and int(split20_manifest_summary.get("qwen_backup_harmful_accepts", -1)) == 0
        and split20_manifest_summary.get("qwen_backup_latency_verdict") == "slower_than_original_max"
        and int(split20_manifest_summary.get("live_calls_performed_by_builder", -1)) == 0
        and "deepseek_full_parallel" in split20_manifest.get("run_commands", {})
        and "deepseek_resume_after_top3" in split20_manifest.get("run_commands", {})
        and "--mode call" in split20_manifest.get("run_commands", {}).get("deepseek_full_parallel", "")
        and "--window-id-file outputs/research_progress_snapshot/split20_deepseek_resume_after_top3_window_ids.txt"
        in split20_manifest.get("run_commands", {}).get("deepseek_resume_after_top3", "")
        and len(split20_pending_ids) == 101
        and len(split20_completed_ids) == 3
        and len(split20_failed_ids) == 2
        and not (set(split20_pending_ids) & set(split20_completed_ids))
        and not split20_manifest_source_missing,
        json.dumps(
            {
                "contract": split20_manifest.get("runtime_contract"),
                "summary": split20_manifest_summary,
                "resume_surface": {
                    "pending_file": split20_resume_surface.get("pending_window_id_file"),
                    "pending_ids": len(split20_pending_ids),
                    "completed_ids": len(split20_completed_ids),
                    "failed_ids": len(split20_failed_ids),
                },
                "missing_sources": split20_manifest_source_missing,
                "commands": sorted(split20_manifest.get("run_commands", {}).keys()),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "split20_resume_export_audit_contract",
        split20_export_audit.get("runtime_contract") == "split20_resume_export_audit_no_live_calls"
        and split20_export_audit.get("status") == "pass"
        and int(split20_export_summary.get("export_prompts", 0)) == 139
        and int(split20_export_summary.get("expected_prompts", 0)) == 139
        and int(split20_export_summary.get("export_parent_windows", 0)) == 101
        and int(split20_export_summary.get("expected_parent_windows", 0)) == 101
        and int(split20_export_summary.get("pending_ids", 0)) == 101
        and int(split20_export_summary.get("completed_ids", 0)) == 3
        and int(split20_export_summary.get("failed_ids", 0)) == 2
        and not split20_export_summary.get("completed_overlap")
        and not split20_export_summary.get("missing_pending")
        and not split20_export_summary.get("extra_parent_ids")
        and not split20_export_summary.get("failed_missing")
        and int(split20_export_summary.get("live_calls_performed", -1)) == 0
        and split20_export_jsonl.exists()
        and "--mode export" in split20_export_audit.get("export_command", "")
        and "--window-id-file outputs/research_progress_snapshot/split20_deepseek_resume_after_top3_window_ids.txt"
        in split20_export_audit.get("export_command", ""),
        json.dumps(
            {
                "contract": split20_export_audit.get("runtime_contract"),
                "status": split20_export_audit.get("status"),
                "summary": split20_export_summary,
                "export_jsonl_exists": split20_export_jsonl.exists(),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "split15_stretch_reexport_audit_contract",
        split15_export_audit.get("runtime_contract") == "split15_stretch_reexport_audit_no_live_calls"
        and split15_export_audit.get("status") == "pass"
        and split15_export_audit.get("source_contracts", {}).get("split_policy_optimization") == "split_policy_optimization_from_existing_artifacts_no_live_calls"
        and split15_export_audit.get("source_contracts", {}).get("split20_full_live_manifest") == "split20_full_live_manifest_no_live_calls"
        and int(split15_export_summary.get("export_prompts", 0)) == 178
        and int(split15_export_summary.get("expected_prompts", 0)) == 178
        and int(split15_export_summary.get("export_parent_windows", 0)) == 104
        and int(split15_export_summary.get("expected_parent_windows", 0)) == 104
        and int(split15_export_summary.get("split_parent_windows", 0)) == 58
        and int(split15_export_summary.get("expected_split_parent_windows", 0)) == 58
        and int(split15_export_summary.get("max_subcalls_per_parent", 0)) == 3
        and int(split15_export_summary.get("prompt_patch_references", 0)) == 1917
        and float(split15_export_summary.get("simulated_p95_call_seconds", 0.0)) == 18.047
        and float(split15_export_summary.get("token_multiplier", 0.0)) == 1.182
        and split15_export_summary.get("requires_new_prompt_export") is True
        and split15_export_summary.get("top3_live_evidence_reusable") is False
        and int(split15_export_summary.get("missing_parent_count", -1)) == 0
        and int(split15_export_summary.get("extra_parent_count", -1)) == 0
        and int(split15_export_summary.get("live_calls_performed", -1)) == 0
        and split15_export_summary.get("no_secret_values_written") is True
        and split15_export_summary.get("no_new_metric_claim") is True
        and split15_export_jsonl.exists()
        and "--mode export" in split15_export_audit.get("export_command", "")
        and "--max-patches-per-call 15" in split15_export_audit.get("export_command", ""),
        json.dumps(
            {
                "contract": split15_export_audit.get("runtime_contract"),
                "status": split15_export_audit.get("status"),
                "summary": split15_export_summary,
                "export_jsonl_exists": split15_export_jsonl.exists(),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    split_policy_summary = split_policy_optimization.get("summary", {})
    split_policy_rows = {int(row.get("max_patches_per_call", 0)): row for row in split_policy_optimization.get("policies", [])}
    add_check(
        checks,
        "split_policy_optimization_contract",
        split_policy_optimization.get("runtime_contract") == "split_policy_optimization_from_existing_artifacts_no_live_calls"
        and split_policy_optimization.get("status") == "ready_primary_resume_blocked_by_live_outputs_or_quota"
        and int(split_policy_summary.get("policy_count", 0)) == 5
        and split_policy_summary.get("primary_policy") == "max20"
        and int(split_policy_summary.get("primary_calls", 0)) == 147
        and int(split_policy_summary.get("primary_resume_calls", 0)) == 139
        and float(split_policy_summary.get("primary_simulated_p95_call_seconds", 0.0)) == 21.358
        and float(split_policy_summary.get("primary_token_multiplier", 0.0)) == 1.118
        and split_policy_summary.get("stretch_policy") == "max15"
        and int(split_policy_summary.get("stretch_calls", 0)) == 178
        and float(split_policy_summary.get("stretch_simulated_p95_call_seconds", 0.0)) == 18.047
        and split_policy_summary.get("stretch_requires_reexport") is True
        and int(split_policy_summary.get("live_calls_performed_by_builder", -1)) == 0
        and split_policy_summary.get("no_metric_claim") is True
        and split_policy_rows.get(20, {}).get("role") == "resume_primary"
        and split_policy_rows.get(20, {}).get("top3_live_evidence_reusable") is True
        and split_policy_rows.get(15, {}).get("role") == "latency_stretch_reexport"
        and split_policy_rows.get(15, {}).get("requires_new_prompt_export") is True
        and "--decisions outputs/runtime_safe_policy_agent/sortformer_diarizen_120_decisions.jsonl"
        in split_policy_optimization.get("commands", {}).get("stretch_reexport_max15", "")
        and "--patch-id-file outputs/runtime_safe_llm_window_batch/deepseek_proxy_high_risk_104w_patch_ids.txt"
        in split_policy_optimization.get("commands", {}).get("stretch_reexport_max15", ""),
        json.dumps(
            {
                "contract": split_policy_optimization.get("runtime_contract"),
                "status": split_policy_optimization.get("status"),
                "summary": split_policy_summary,
                "roles": {key: value.get("role") for key, value in sorted(split_policy_rows.items())},
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    phase_scorecard_summary = phase_result_scorecard.get("summary", {})
    phase_scorecard_rows = {row.get("result_id"): row for row in phase_result_scorecard.get("rows", [])}
    required_phase_result_ids = {
        "recording_holdout_selector_gain",
        "bootstrap_rule_recover_gain",
        "slow_model_upper_bound",
        "runtime_latency_slo_pass",
        "llm_guard_zero_harm",
        "omni_acoustic_smoke",
        "qwen_fallback_smoke",
        "deepseek_no_go",
    }
    add_check(
        checks,
        "phase_result_scorecard_contract",
        phase_result_scorecard.get("runtime_contract") == "phase_result_scorecard_from_existing_artifacts_no_live_calls"
        and phase_result_scorecard.get("status") == "pass"
        and int(phase_scorecard_summary.get("result_rows", 0)) == 8
        and int(phase_scorecard_summary.get("reportable_rows", 0)) == 7
        and int(phase_scorecard_summary.get("deepseek_no_go_rows", 0)) == 1
        and int(phase_scorecard_summary.get("claim_now_slo_pass", 0)) == 4
        and int(phase_scorecard_summary.get("claim_now_slo_rows", 0)) == 4
        and int(phase_scorecard_summary.get("selector_positive_splits", 0)) == 8
        and int(phase_scorecard_summary.get("selector_splits", 0)) == 8
        and float(phase_scorecard_summary.get("selector_weighted_delta_pp", 0.0)) == 2.05
        and float(phase_scorecard_summary.get("rule_bootstrap_delta_pp", 0.0)) == 2.17
        and int(phase_scorecard_summary.get("guard_harmful_accepts", -1)) == 0
        and int(phase_scorecard_summary.get("guard_safe_accepts", 0)) == 323
        and int(phase_scorecard_summary.get("omni_smoke_windows", 0)) == 12
        and int(phase_scorecard_summary.get("qwen_fallback_calls", 0)) == 4
        and int(phase_scorecard_summary.get("traceability_rows", 0)) == 54
        and int(phase_scorecard_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(phase_scorecard_summary.get("live_calls_performed_by_builder", -1)) == 0
        and phase_scorecard_summary.get("no_live_calls_performed") is True
        and phase_scorecard_summary.get("no_scoring_commands_executed") is True
        and phase_scorecard_summary.get("no_deepseek_api_calls") is True
        and phase_scorecard_summary.get("no_new_metric_claim") is True
        and required_phase_result_ids.issubset(phase_scorecard_rows)
        and phase_scorecard_rows.get("recording_holdout_selector_gain", {}).get("claim_boundary") == "not_true_heldout_until_new_recordings"
        and phase_scorecard_rows.get("runtime_latency_slo_pass", {}).get("claim_boundary") == "post_live_latency_not_claimed"
        and phase_scorecard_rows.get("omni_acoustic_smoke", {}).get("claim_boundary") == "label_only_no_timeline_writeback"
        and phase_scorecard_rows.get("qwen_fallback_smoke", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and phase_scorecard_rows.get("deepseek_no_go", {}).get("claim_boundary") == "do_not_use_deepseek_api"
        and all(row.get("source_artifacts_exist") is True for row in phase_scorecard_rows.values()),
        json.dumps(
            {
                "contract": phase_result_scorecard.get("runtime_contract"),
                "status": phase_result_scorecard.get("status"),
                "summary": phase_scorecard_summary,
                "missing_ids": sorted(required_phase_result_ids - set(phase_scorecard_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    provider_route_summary = live_provider_routing.get("summary", {})
    provider_route_rows = {row.get("route_id"): row for row in live_provider_routing.get("rows", [])}
    required_provider_route_ids = {
        "deepseek_resume_primary",
        "qwen_full_backup_optional",
        "omni48_label_only",
        "default_live_execute",
    }
    add_check(
        checks,
        "live_provider_routing_decision_contract",
        live_provider_routing.get("runtime_contract") == "live_provider_routing_decision_no_live_calls_no_secret_values"
        and live_provider_routing.get("status") == "blocked_no_default_primary_provider"
        and live_provider_routing.get("source_contracts", {}).get("phase_result_scorecard") == "phase_result_scorecard_from_existing_artifacts_no_live_calls"
        and live_provider_routing.get("source_contracts", {}).get("live_run_readiness") == "live_run_readiness_non_secret_no_live_calls"
        and live_provider_routing.get("source_contracts", {}).get("live_command_surface_audit") == "live_command_surface_audit_no_live_calls"
        and live_provider_routing.get("source_contracts", {}).get("live_execution_launcher") == "live_execution_launcher_dry_run_no_live_calls"
        and live_provider_routing.get("source_contracts", {}).get("live_output_audit") == "live_output_audit_no_live_calls"
        and int(provider_route_summary.get("route_rows", 0)) == 4
        and int(provider_route_summary.get("default_selected_routes", -1)) == 0
        and provider_route_summary.get("recommended_default_execute_scope") == "none"
        and provider_route_summary.get("deepseek_no_go") is True
        and int(provider_route_summary.get("deepseek_planned_calls_not_selected", 0)) == 139
        and int(provider_route_summary.get("qwen_fallback_calls", 0)) == 147
        and int(provider_route_summary.get("qwen_smoke_calls", 0)) == 4
        and int(provider_route_summary.get("qwen_smoke_harmful_accepts", -1)) == 0
        and int(provider_route_summary.get("omni_label_calls", 0)) == 96
        and provider_route_summary.get("credential_ready") is False
        and int(provider_route_summary.get("ready_runs", -1)) == 0
        and int(provider_route_summary.get("command_ready_count", 0)) == 3
        and int(provider_route_summary.get("missing_output_surfaces", 0)) == 3
        and int(provider_route_summary.get("missing_source_rows", -1)) == 0
        and int(provider_route_summary.get("live_calls_performed_by_builder", -1)) == 0
        and provider_route_summary.get("no_live_calls_performed") is True
        and provider_route_summary.get("no_secret_values_written") is True
        and provider_route_summary.get("no_scoring_commands_executed") is True
        and provider_route_summary.get("no_new_metric_claim") is True
        and live_provider_routing.get("recommended_operator_action") == "no default live execute; explicit qwen/omni fallback only after credentials are ready"
        and required_provider_route_ids.issubset(provider_route_rows)
        and provider_route_rows.get("deepseek_resume_primary", {}).get("status") == "no_go_current"
        and provider_route_rows.get("deepseek_resume_primary", {}).get("selected_for_default_execute") is False
        and provider_route_rows.get("deepseek_resume_primary", {}).get("claim_boundary") == "do_not_use_deepseek_api_by_default"
        and provider_route_rows.get("qwen_full_backup_optional", {}).get("claim_boundary") == "fallback_only_not_primary_latency_claim"
        and provider_route_rows.get("omni48_label_only", {}).get("claim_boundary") == "label_only_no_timeline_writeback"
        and provider_route_rows.get("default_live_execute", {}).get("claim_boundary") == "default_execute_blocked_by_provider_route"
        and all(row.get("source_artifacts_exist") is True for row in provider_route_rows.values()),
        json.dumps(
            {
                "contract": live_provider_routing.get("runtime_contract"),
                "status": live_provider_routing.get("status"),
                "summary": provider_route_summary,
                "missing_ids": sorted(required_provider_route_ids - set(provider_route_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    promotion_summary = post_live_claim_promotion.get("summary", {})
    promotion_gates = {row.get("gate_id"): row for row in post_live_claim_promotion.get("gates", [])}
    required_promotion_gate_ids = {
        "current_latency_slo_claims",
        "deepseek_split20_resume_latency",
        "deepseek_split20_resume_safety",
        "omni48_label_metrics",
        "qwen_full_backup_claim",
        "selector_true_heldout_claim",
        "live_execution_handoff",
        "report_ppt_sync_after_promotion",
    }
    add_check(
        checks,
        "post_live_claim_promotion_gate_contract",
        post_live_claim_promotion.get("runtime_contract") == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_claim_promotion.get("status") == "pass"
        and post_live_claim_promotion.get("promotion_policy") == "promote_only_after_output_audit_scoring_slo_and_traceability_pass"
        and int(promotion_summary.get("gate_count", 0)) == 8
        and int(promotion_summary.get("ready_to_promote_count", -1)) == 0
        and int(promotion_summary.get("blocked_count", 0)) == 5
        and int(promotion_summary.get("preserve_count", 0)) == 2
        and int(promotion_summary.get("fallback_only_count", 0)) == 1
        and int(promotion_summary.get("missing_source_rows", -1)) == 0
        and int(promotion_summary.get("live_calls_performed_by_builder", -1)) == 0
        and promotion_summary.get("no_secret_values_written") is True
        and promotion_summary.get("no_new_metric_claim") is True
        and int(promotion_summary.get("claim_now_slo_pass", 0)) == 4
        and int(promotion_summary.get("claim_now_slo_rows", 0)) == 4
        and int(promotion_summary.get("ready_runs", -1)) == 0
        and int(promotion_summary.get("missing_output_surfaces", 0)) == 3
        and int(promotion_summary.get("ready_to_score_steps", -1)) == 0
        and int(promotion_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(promotion_summary.get("traceability_rows", 0)) == 54
        and required_promotion_gate_ids.issubset(promotion_gates)
        and not post_live_claim_promotion.get("ready_to_promote_gate_ids")
        and set(post_live_claim_promotion.get("blocked_gate_ids", []))
        == {
            "deepseek_split20_resume_latency",
            "deepseek_split20_resume_safety",
            "omni48_label_metrics",
            "selector_true_heldout_claim",
            "live_execution_handoff",
        }
        and post_live_claim_promotion.get("fallback_only_gate_ids") == ["qwen_full_backup_claim"]
        and not post_live_claim_promotion.get("missing_source_gate_ids")
        and promotion_gates.get("deepseek_split20_resume_latency", {}).get("promotion_decision") == "blocked_missing_live_output"
        and promotion_gates.get("omni48_label_metrics", {}).get("promotion_decision") == "blocked_missing_live_output"
        and promotion_gates.get("selector_true_heldout_claim", {}).get("promotion_decision") == "blocked_waiting_valid_sealed_split",
        json.dumps(
            {
                "contract": post_live_claim_promotion.get("runtime_contract"),
                "status": post_live_claim_promotion.get("status"),
                "summary": promotion_summary,
                "missing_ids": sorted(required_promotion_gate_ids - set(promotion_gates)),
                "ready": post_live_claim_promotion.get("ready_to_promote_gate_ids"),
                "blocked": post_live_claim_promotion.get("blocked_gate_ids"),
                "fallback": post_live_claim_promotion.get("fallback_only_gate_ids"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    claim_promotion_receipt_summary = post_live_claim_promotion_receipt.get("summary", {})
    claim_promotion_receipt_rows = {
        row.get("receipt_id"): row
        for row in post_live_claim_promotion_receipt.get("rows", [])
    }
    required_claim_promotion_receipt_ids = {
        "promotion_gate_presence",
        "promotion_decision_receipt",
        "promotion_preflight_receipt",
        "scoring_and_time_receipts",
        "report_ppt_traceability_receipt",
        "validation_after_promotion_receipt",
    }
    add_check(
        checks,
        "post_live_claim_promotion_receipt_audit_contract",
        post_live_claim_promotion_receipt.get("runtime_contract")
        == "post_live_claim_promotion_receipt_audit_no_live_or_scoring_or_claim_writes_no_secret_values"
        and post_live_claim_promotion_receipt.get("status") == "blocked_no_claim_promotion_receipt"
        and post_live_claim_promotion_receipt.get("source_contracts", {}).get("post_live_claim_promotion_gate")
        == "post_live_claim_promotion_gate_no_live_calls"
        and post_live_claim_promotion_receipt.get("source_contracts", {}).get("post_live_promotion_preflight_audit")
        == "post_live_promotion_preflight_audit_no_live_or_scoring_calls"
        and post_live_claim_promotion_receipt.get("source_contracts", {}).get("post_live_scoring_receipt_audit")
        == "post_live_scoring_receipt_audit_no_live_or_scoring_calls_no_secret_values"
        and post_live_claim_promotion_receipt.get("source_contracts", {}).get("post_live_time_metric_receipt_audit")
        == "post_live_time_metric_receipt_audit_no_live_or_scoring_calls_no_secret_values"
        and post_live_claim_promotion_receipt.get("source_contracts", {}).get("report_ppt_traceability")
        == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and post_live_claim_promotion_receipt.get("source_contracts", {}).get("latest_artifact_validation")
        == "latest_research_artifact_validation"
        and int(claim_promotion_receipt_summary.get("receipt_rows", 0)) == 6
        and int(claim_promotion_receipt_summary.get("pass_rows", -1)) == 3
        and int(claim_promotion_receipt_summary.get("blocked_rows", 0)) == 3
        and int(claim_promotion_receipt_summary.get("missing_source_rows", -1)) == 0
        and int(claim_promotion_receipt_summary.get("promotion_gate_count", 0)) == 8
        and int(claim_promotion_receipt_summary.get("ready_to_promote_count", -1)) == 0
        and int(claim_promotion_receipt_summary.get("blocked_promotion_count", 0)) == 5
        and int(claim_promotion_receipt_summary.get("fallback_only_count", 0)) == 1
        and claim_promotion_receipt_summary.get("preflight_ready") is False
        and int(claim_promotion_receipt_summary.get("promotion_preflight_pass_rows", 0)) == 2
        and int(claim_promotion_receipt_summary.get("promotion_preflight_blocked_rows", 0)) == 4
        and claim_promotion_receipt_summary.get("scoring_receipt_ready") is False
        and claim_promotion_receipt_summary.get("time_metric_receipt_ready") is False
        and int(claim_promotion_receipt_summary.get("computed_time_metric_rows", -1)) == 0
        and claim_promotion_receipt_summary.get("report_ppt_synced") is True
        and claim_promotion_receipt_summary.get("validation_passed") is True
        and int(claim_promotion_receipt_summary.get("validation_non_self_failed_checks", -1)) == 0
        and claim_promotion_receipt_summary.get("ready_for_claim_write") is False
        and int(claim_promotion_receipt_summary.get("traceability_rows", 0)) == 54
        and int(claim_promotion_receipt_summary.get("traceability_fully_covered_rows", 0)) == 54
        and int(claim_promotion_receipt_summary.get("live_calls_performed_by_builder", -1)) == 0
        and claim_promotion_receipt_summary.get("no_live_calls_performed_by_auditor") is True
        and claim_promotion_receipt_summary.get("no_scoring_commands_executed_by_auditor") is True
        and claim_promotion_receipt_summary.get("no_claim_writes_performed_by_auditor") is True
        and claim_promotion_receipt_summary.get("no_secret_values_written") is True
        and claim_promotion_receipt_summary.get("no_new_metric_claim") is True
        and required_claim_promotion_receipt_ids.issubset(claim_promotion_receipt_rows)
        and claim_promotion_receipt_rows.get("promotion_gate_presence", {}).get("status") == "pass"
        and claim_promotion_receipt_rows.get("report_ppt_traceability_receipt", {}).get("status") == "pass"
        and claim_promotion_receipt_rows.get("validation_after_promotion_receipt", {}).get("status") == "pass"
        and claim_promotion_receipt_rows.get("promotion_decision_receipt", {}).get("status") == "blocked"
        and claim_promotion_receipt_rows.get("promotion_preflight_receipt", {}).get("status") == "blocked"
        and claim_promotion_receipt_rows.get("scoring_and_time_receipts", {}).get("status") == "blocked"
        and all(row.get("source_artifacts_exist") is True for row in claim_promotion_receipt_rows.values()),
        json.dumps(
            {
                "contract": post_live_claim_promotion_receipt.get("runtime_contract"),
                "status": post_live_claim_promotion_receipt.get("status"),
                "summary": claim_promotion_receipt_summary,
                "missing_ids": sorted(required_claim_promotion_receipt_ids - set(claim_promotion_receipt_rows)),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    traceability_summary = report_ppt_traceability.get("summary", {})
    traceability_rows = {row.get("trace_id"): row for row in report_ppt_traceability.get("rows", [])}
    required_trace_ids = {
        "phase_result_scorecard",
        "live_provider_routing_decision",
        "latency_budget_ledger",
        "stage_latency_slo_audit",
        "latency_risk_margin_audit",
        "latency_risk_mitigation_plan",
        "selector_candidate_scan",
        "selector_split_validation",
        "selector_protocol",
        "omni48_call_manifest",
        "split20_manifest",
        "split20_resume_export_audit",
        "split15_stretch_reexport_audit",
        "split_policy_optimization",
        "live_readiness",
        "live_agent_plan",
        "live_execution_runbook",
        "live_execution_bundle",
        "live_execution_handoff_packet",
        "live_command_surface_audit",
        "live_execution_eligibility_gate",
        "live_execution_receipt_audit",
        "live_execution_launcher",
        "live_output_repair_plan",
        "live_resume_state_audit",
        "live_runtime_environment_audit",
        "live_execution_timing_plan",
        "live_retry_budget_audit",
        "live_token_quota_budget_audit",
        "live_failure_recovery_playbook",
        "live_metric_extraction_contract",
        "live_output_schema_contract",
        "post_live_acceptance_scorecard",
        "post_live_latency_claim_matrix",
        "post_live_time_metric_statistics_plan",
        "post_live_time_metric_extractor",
        "post_live_time_metric_receipt_audit",
        "live_parallelism_sensitivity",
        "live_postrun_closure",
        "live_output_audit",
        "live_scoring_readiness",
        "post_live_scoring_execution_plan",
        "post_live_scoring_launcher",
        "post_live_scoring_receipt_audit",
        "post_live_scoring_output_audit",
        "post_live_evidence_dependency_dag",
        "post_live_promotion_preflight_audit",
        "live_input_integrity_audit",
        "post_live_claim_promotion_gate",
        "post_live_claim_promotion_receipt_audit",
        "runtime_replay_manifest",
        "memory_update_replay",
        "claims_manifest",
        "next_experiment_queue",
    }
    add_check(
        checks,
        "report_ppt_traceability_contract",
        report_ppt_traceability.get("runtime_contract") == "report_ppt_traceability_from_existing_artifacts_no_live_calls"
        and report_ppt_traceability.get("status") == "pass"
        and int(traceability_summary.get("traceability_rows", 0)) == 54
        and int(traceability_summary.get("source_covered_rows", 0)) == 54
        and int(traceability_summary.get("report_covered_rows", 0)) == 54
        and int(traceability_summary.get("ppt_covered_rows", 0)) == 54
        and int(traceability_summary.get("fully_covered_rows", 0)) == 54
        and int(traceability_summary.get("missing_report_rows", -1)) == 0
        and int(traceability_summary.get("missing_ppt_rows", -1)) == 0
        and int(traceability_summary.get("missing_source_rows", -1)) == 0
        and int(traceability_summary.get("live_calls_performed_by_builder", -1)) == 0
        and traceability_summary.get("no_new_metric_claim") is True
        and required_trace_ids.issubset(traceability_rows)
        and not report_ppt_traceability.get("missing_report_trace_ids")
        and not report_ppt_traceability.get("missing_ppt_trace_ids")
        and not report_ppt_traceability.get("missing_source_trace_ids")
        and traceability_rows.get("phase_result_scorecard", {}).get("claim_boundary") == "phase_scorecard_no_live_or_scoring_calls"
        and traceability_rows.get("live_provider_routing_decision", {}).get("claim_boundary") == "provider_route_blocks_default_deepseek_execute"
        and traceability_rows.get("live_execution_runbook", {}).get("claim_boundary") == "no_secret_values_no_live_calls"
        and traceability_rows.get("live_execution_bundle", {}).get("claim_boundary") == "execution_bundle_no_live_calls_no_secret_values"
        and traceability_rows.get("live_execution_handoff_packet", {}).get("claim_boundary") == "live_execution_handoff_no_live_calls_no_secret_values"
        and traceability_rows.get("live_command_surface_audit", {}).get("claim_boundary") == "command_surface_no_live_calls_no_secret_values"
        and traceability_rows.get("live_execution_eligibility_gate", {}).get("claim_boundary") == "live_execution_eligibility_no_live_calls"
        and traceability_rows.get("live_execution_receipt_audit", {}).get("claim_boundary") == "live_execution_receipt_no_live_calls"
        and traceability_rows.get("live_execution_launcher", {}).get("claim_boundary") == "launcher_dry_run_no_live_calls_execute_live_requires_flag"
        and traceability_rows.get("live_output_repair_plan", {}).get("claim_boundary") == "repair_plan_no_live_or_scoring_calls"
        and traceability_rows.get("live_resume_state_audit", {}).get("claim_boundary") == "resume_state_no_live_calls_no_metric_claim"
        and traceability_rows.get("live_runtime_environment_audit", {}).get("claim_boundary") == "runtime_env_no_live_calls_no_secret_values"
        and traceability_rows.get("live_execution_timing_plan", {}).get("claim_boundary") == "timing_plan_no_live_metric_claim"
        and traceability_rows.get("live_retry_budget_audit", {}).get("claim_boundary") == "retry_budget_no_live_metric_claim"
        and traceability_rows.get("live_token_quota_budget_audit", {}).get("claim_boundary") == "token_quota_budget_no_live_metric_claim"
        and traceability_rows.get("live_failure_recovery_playbook", {}).get("claim_boundary") == "failure_recovery_no_live_metric_claim"
        and traceability_rows.get("live_metric_extraction_contract", {}).get("claim_boundary") == "metric_extraction_schema_no_live_metric_claim"
        and traceability_rows.get("live_output_schema_contract", {}).get("claim_boundary") == "output_schema_no_live_metric_claim"
        and traceability_rows.get("post_live_acceptance_scorecard", {}).get("claim_boundary") == "acceptance_scorecard_no_live_metric_claim"
        and traceability_rows.get("post_live_latency_claim_matrix", {}).get("claim_boundary") == "latency_claim_matrix_no_live_metric_claim"
        and traceability_rows.get("post_live_time_metric_statistics_plan", {}).get("claim_boundary") == "time_statistics_plan_no_live_or_scoring_calls"
        and traceability_rows.get("post_live_time_metric_extractor", {}).get("claim_boundary") == "time_metric_extractor_no_live_or_scoring_calls"
        and traceability_rows.get("post_live_time_metric_receipt_audit", {}).get("claim_boundary") == "time_metric_receipt_no_live_or_scoring_calls"
        and traceability_rows.get("live_parallelism_sensitivity", {}).get("claim_boundary") == "parallelism_plan_no_live_metric_claim"
        and traceability_rows.get("live_input_integrity_audit", {}).get("claim_boundary") == "input_integrity_no_live_metric_claim"
        and traceability_rows.get("post_live_scoring_execution_plan", {}).get("claim_boundary") == "scoring_execution_plan_no_live_or_scoring_calls"
        and traceability_rows.get("post_live_scoring_launcher", {}).get("claim_boundary") == "scoring_launcher_dry_run_no_scoring_calls"
        and traceability_rows.get("post_live_scoring_receipt_audit", {}).get("claim_boundary") == "scoring_receipt_no_live_or_scoring_calls"
        and traceability_rows.get("post_live_scoring_output_audit", {}).get("claim_boundary") == "scoring_output_audit_no_scoring_calls"
        and traceability_rows.get("post_live_evidence_dependency_dag", {}).get("claim_boundary") == "evidence_dependency_dag_no_live_or_scoring_calls"
        and traceability_rows.get("post_live_promotion_preflight_audit", {}).get("claim_boundary") == "promotion_preflight_no_live_or_scoring_calls"
        and traceability_rows.get("latency_risk_margin_audit", {}).get("claim_boundary") == "risk_label_no_new_metric_claim"
        and traceability_rows.get("latency_risk_mitigation_plan", {}).get("claim_boundary") == "mitigation_plan_no_new_metric_claim"
        and traceability_rows.get("post_live_claim_promotion_gate", {}).get("claim_boundary") == "promote_only_after_output_audit_scoring_slo_and_traceability_pass"
        and traceability_rows.get("post_live_claim_promotion_receipt_audit", {}).get("claim_boundary") == "claim_promotion_receipt_no_claim_writes"
        and traceability_rows.get("split15_stretch_reexport_audit", {}).get("claim_boundary") == "stretch_export_only_no_live_calls"
        and traceability_rows.get("split_policy_optimization", {}).get("claim_boundary") == "no_metric_claim"
        and traceability_rows.get("selector_split_validation", {}).get("claim_boundary") == "blocked_waiting_for_valid_sealed_split",
        json.dumps(
            {
                "contract": report_ppt_traceability.get("runtime_contract"),
                "status": report_ppt_traceability.get("status"),
                "summary": traceability_summary,
                "missing_ids": sorted(required_trace_ids - set(traceability_rows)),
                "missing_report": report_ppt_traceability.get("missing_report_trace_ids"),
                "missing_ppt": report_ppt_traceability.get("missing_ppt_trace_ids"),
                "missing_sources": report_ppt_traceability.get("missing_source_trace_ids"),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "report_references_snapshot_refresh",
        not missing_fragments(report_text, report_required_fragments),
        json.dumps(
            {
                "required": len(report_required_fragments),
                "missing": missing_fragments(report_text, report_required_fragments),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    add_check(
        checks,
        "ppt_references_latest_snapshot",
        not missing_fragments(ppt_text, ppt_required_fragments),
        json.dumps(
            {
                "required": len(ppt_required_fragments),
                "missing": missing_fragments(ppt_text, ppt_required_fragments),
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )

    failed = [check for check in checks if check["status"] != "pass"]
    return {
        "status": "pass" if not failed else "fail",
        "checks": checks,
        "failed_checks": [check["name"] for check in failed],
        "runtime_contract": "latest_research_artifact_validation",
    }


def write_markdown(result: dict[str, Any], path: Path) -> None:
    lines = [
        "# Latest Research Artifact Validation",
        "",
        f"- Status: `{result['status']}`",
        f"- Checks: `{len(result['checks'])}`",
        f"- Failed: `{len(result['failed_checks'])}`",
        "",
        "| Check | Status | Evidence |",
        "|---|---|---|",
    ]
    for check in result["checks"]:
        evidence = str(check["evidence"]).replace("|", "/")
        lines.append(f"| {check['name']} | {check['status']} | {evidence} |")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output-json", type=Path, default=Path("outputs/research_progress_snapshot/latest_artifact_validation.json"))
    parser.add_argument("--output-md", type=Path, default=Path("outputs/research_progress_snapshot/latest_artifact_validation.md"))
    args = parser.parse_args()

    result = validate()
    args.output_json.parent.mkdir(parents=True, exist_ok=True)
    args.output_json.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    write_markdown(result, args.output_md)
    print(f"Wrote {args.output_json}")
    print(f"Wrote {args.output_md}")
    print(json.dumps({"status": result["status"], "failed_checks": result["failed_checks"]}, ensure_ascii=False))
    if result["status"] != "pass":
        raise SystemExit(1)


if __name__ == "__main__":
    main()
